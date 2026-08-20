"""Build editable, style-aware PPTX decks from a brief and source assets."""

from __future__ import annotations

import base64
import json
import logging
import math
import re
import sys
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from database.models import PresentationJob, PresentationJobAsset, now_utc
from database.session import SessionLocal
from services.openai_service import OpenAIService
from services.token_usage_service import record_token_usage

logger = logging.getLogger(__name__)

OUTPUT_ROOT = Path(__file__).resolve().parents[1] / "generated" / "presentations"
MAX_REFERENCE_TEXT = 70_000
MAX_TEMPLATE_CONTEXT = 32_000
MAX_BLUEPRINT_CONTEXT = 24_000
MAX_VISION_IMAGES = 4
MAX_VISION_IMAGE_BYTES = 6 * 1024 * 1024
MAX_VISION_TOTAL_BYTES = 16 * 1024 * 1024
CANVAS_W = 13.333
CANVAS_H = 7.5

# The skill is checked into the project so native template filling also works
# on a server without a user-level Codex/Claude installation.
PPT_MASTER_SCRIPTS = Path(__file__).resolve().parents[2] / "skills" / "ppt-master" / "scripts"

STYLE_PRESETS: dict[str, dict[str, Any]] = {
    # Calm blue/white briefing palette for state-owned enterprise reporting.
    "state-briefing": {"bg": "FFFFFF", "fg": "263238", "muted": "5B6472", "accent": "00479D", "accent2": "5B9BD5", "decor": "state"},
    # Reference-derived formal families. They share restrained enterprise
    # typography, but use different page rhythm and control vocabularies.
    "aviation-blue": {"bg": "FFFFFF", "fg": "17324A", "muted": "607487", "accent": "005BAC", "accent2": "4FA3D1", "decor": "aviation"},
    "aqua-planning": {"bg": "F3FBFC", "fg": "173A3E", "muted": "5A7377", "accent": "32B8C7", "accent2": "80D6DD", "decor": "aqua"},
    "security-report": {"bg": "FFFFFF", "fg": "1D2F46", "muted": "5F6F82", "accent": "07569F", "accent2": "4874CB", "decor": "security"},
    "dark-tech": {"bg": "0B1020", "fg": "F8FAFC", "muted": "A8B4CC", "accent": "5B7CFF", "accent2": "2DD4BF", "decor": "tech"},
    "swiss-minimal": {"bg": "F8F8F5", "fg": "141414", "muted": "5F625F", "accent": "E5484D", "accent2": "1C6BFF", "decor": "swiss"},
    "glassmorphism": {"bg": "11182A", "fg": "F7FAFF", "muted": "B8C2D9", "accent": "7C8CFF", "accent2": "59E1C2", "decor": "glass"},
    "data-journalism": {"bg": "F5F3ED", "fg": "202329", "muted": "62686C", "accent": "C2412D", "accent2": "137B80", "decor": "data"},
    "editorial": {"bg": "F5F1EA", "fg": "1E1E1E", "muted": "6B6259", "accent": "B84A39", "accent2": "D89A3D", "decor": "editorial"},
    "blueprint": {"bg": "102238", "fg": "EAF5FF", "muted": "A6C0D6", "accent": "5CD6FF", "accent2": "F5B94C", "decor": "blueprint"},
    "ink-notes": {"bg": "F7F5EE", "fg": "191919", "muted": "65635D", "accent": "E2574C", "accent2": "3B6CFF", "decor": "ink"},
    "photo-editorial": {"bg": "F4F1EC", "fg": "1D1D1D", "muted": "655F58", "accent": "B84A39", "accent2": "272727", "decor": "editorial"},
    "soft-rounded": {"bg": "F4F7FB", "fg": "182334", "muted": "596A7F", "accent": "5B6EFF", "accent2": "16A887", "decor": "soft"},
    "vivid-launch": {"bg": "190E35", "fg": "FFF8FF", "muted": "D4BDEB", "accent": "FF5C8A", "accent2": "5B7CFF", "decor": "vivid"},
    "clean-business": {"bg": "F7F9FC", "fg": "122033", "muted": "506176", "accent": "2563EB", "accent2": "0F9F8E", "decor": "swiss"},
}

STYLE_KEYWORDS: dict[str, tuple[str, ...]] = {
    "state-briefing": ("国企", "政府", "政务", "稳重", "沉稳", "汇报", "蓝白", "机场", "state-owned", "briefing"),
    "aviation-blue": ("机场蓝白", "航空蓝白", "航站楼专项", "机场专项", "aviation briefing", "airport blue"),
    "aqua-planning": ("浅青规划", "年度规划", "年度工作计划", "未来规划", "aqua planning", "annual planning"),
    "security-report": ("安护深蓝", "安保总结", "安全年度总结", "安检护卫", "security annual report", "security report"),
    "dark-tech": ("科技", "未来", "人工智能", "ai", "赛博", "暗色", "深色", "科技感", "tech", "cyber", "futuristic"),
    "swiss-minimal": ("极简", "留白", "网格", "瑞士", "minimal", "swiss", "minimalist"),
    "glassmorphism": ("玻璃", "毛玻璃", "半透明", "glassmorphism", "glass"),
    "data-journalism": ("数据新闻", "数据报告", "数据分析", "金融", "投研", "仪表盘", "dashboard", "data", "财务"),
    "editorial": ("杂志", "社论", "出版物", "编辑", "刊物", "editorial", "magazine"),
    "blueprint": ("蓝图", "工程图", "工程", "架构图", "系统架构", "技术图纸", "schematic", "blueprint"),
    "ink-notes": ("手绘", "白板", "墨迹", "草图", "涂鸦", "ink", "sketch"),
    "photo-editorial": ("摄影", "照片", "图片主导", "影像", "photo", "photography"),
    "soft-rounded": ("圆润", "友好", "柔和", "亲和", "saas", "soft-rounded"),
    "vivid-launch": ("发布会", "发布", "高能", "鲜明", "霓虹", "vivid", "launch", "发布活动"),
}

MODE_GUIDANCE = {
    "pyramid": "结论先行，标题写成判断，正文用证据支撑决策。",
    "narrative": "用情境、张力、转折和解决方案推进，标题像故事节拍。",
    "instructional": "按先简单后复杂的步骤拆解，平行概念保持同样层级。",
    "showcase": "视觉和数字先行，控制文字密度，用强弱节奏制造记忆点。",
    "briefing": "中性、完整、便于查阅，主题标题和同层级信息保持一致。",
}

FEATURE_GUIDANCE = {
    "assertion_titles": "结论式标题",
    "kicker_summary": "眉题和页面一句话总结",
    "layout_variety": "多版式轮换",
    "visual_decor": "与风格匹配的装饰元素",
    "metrics": "关键指标页",
    "roadmap": "路线图或阶段时间轴",
    "comparison": "对比决策页",
    "source_notes": "来源备注",
    "data_story": "数据叙事与图表",
    "template_fidelity": "原生模板保真填充",
}


@dataclass
class Canvas:
    prs: Presentation

    def x(self, value: float) -> int:
        return int(self.prs.slide_width * value / CANVAS_W)

    def y(self, value: float) -> int:
        return int(self.prs.slide_height * value / CANVAS_H)

    def w(self, value: float) -> int:
        return int(self.prs.slide_width * value / CANVAS_W)

    def h(self, value: float) -> int:
        return int(self.prs.slide_height * value / CANVAS_H)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _clip(text: str, max_chars: int) -> str:
    text = " ".join(str(text or "").split())
    return text if len(text) <= max_chars else text[: max_chars - 1].rstrip() + "…"


def _clip_multiline(text: str, max_chars: int) -> str:
    """Bound extracted document text without destroying row/page structure."""
    lines = [" ".join(line.split()) for line in str(text or "").replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    normalized = "\n".join(line for line in lines if line).strip()
    return normalized if len(normalized) <= max_chars else normalized[: max_chars - 1].rstrip() + "…"


def _safe_json(text: str) -> dict[str, Any] | None:
    candidate = text.strip()
    if "```" in candidate:
        blocks = re.findall(r"```(?:json)?\s*(.*?)```", candidate, flags=re.S | re.I)
        candidate = blocks[0].strip() if blocks else candidate
    start, end = candidate.find("{"), candidate.rfind("}")
    if start < 0 or end <= start:
        return None
    try:
        parsed = json.loads(candidate[start : end + 1])
    except (TypeError, ValueError):
        return None
    return parsed if isinstance(parsed, dict) else None


def _job_options(job: PresentationJob) -> tuple[str, set[str]]:
    try:
        raw = json.loads(job.metadata_json or "{}")
    except (TypeError, ValueError):
        raw = {}
    mode = str(raw.get("mode") or "pyramid")
    if "features" in raw:
        features = {str(item) for item in raw.get("features", []) if str(item)}
    else:
        features = set(FEATURE_GUIDANCE)
    return mode, features


def _style_from_text(text: str) -> str | None:
    """Find an explicit visual direction in the user's own brief."""
    haystack = str(text or "").casefold()
    style_markers = ("风格", "视觉", "样式", "设计", "配色", "主题", "采用", "style", "visual", "design", "theme")
    explicit_phrases = {
        "国企蓝白": "state-briefing", "政务蓝白": "state-briefing", "暗夜科技": "dark-tech",
        "机场蓝白": "aviation-blue", "航空蓝白": "aviation-blue", "浅青规划": "aqua-planning",
        "年度规划风格": "aqua-planning", "安护深蓝": "security-report", "安保总结风格": "security-report",
        "瑞士极简": "swiss-minimal", "玻璃拟态": "glassmorphism", "数据新闻": "data-journalism",
        "杂志社论": "editorial", "工程蓝图": "blueprint", "墨迹笔记": "ink-notes",
        "摄影社论": "photo-editorial", "柔和圆角": "soft-rounded", "鲜明发布": "vivid-launch",
    }
    for phrase, style_id in explicit_phrases.items():
        if phrase in haystack:
            return style_id
    generic_keywords = {
        "ai", "data", "tech", "cyber", "futuristic", "photo", "launch", "glass", "minimal", "swiss", "blueprint", "ink", "sketch", "magazine", "saas", "vivid",
        "国企", "政府", "政务", "稳重", "沉稳", "汇报", "机场", "年度", "规划", "安保", "安检", "护卫", "总结", "科技", "未来", "人工智能", "暗色", "深色", "极简", "留白", "网格", "玻璃", "工程", "架构图", "手绘", "摄影", "照片", "发布", "发布会", "圆润", "柔和",
    }
    has_style_marker = any(marker in haystack for marker in style_markers)
    for style_id, keywords in STYLE_KEYWORDS.items():
        for keyword in keywords:
            normalized = keyword.casefold()
            if normalized not in haystack:
                continue
            if normalized in generic_keywords and not has_style_marker:
                continue
            return style_id
    return None


def _choose_intelligent_style(text: str, *, has_images: bool, has_template: bool) -> str:
    """Choose a deterministic style when the user asks for a random direction."""
    explicit = _style_from_text(text)
    if explicit:
        return explicit
    compact = re.sub(r"\s+", "", text.casefold())
    if any(token in compact for token in ("安检护卫", "安保年度总结", "安全年度总结", "安护总结")):
        return "security-report"
    if any(token in compact for token in ("年度工作计划", "年度规划", "明年规划", "工作展望", "年度部署")):
        return "aqua-planning"
    if any(token in compact for token in ("航站楼外包", "机场专项汇报", "机场岗位", "航站楼岗位")):
        return "aviation-blue"
    # Random is intentionally conservative for management reporting. Images
    # become evidence panels inside a calm briefing palette rather than a
    # reason to switch the whole deck to an editorial/photo treatment.
    if has_images and any(token in text.casefold() for token in ("摄影", "影像", "图片主导", "photo", "photography")):
        return "photo-editorial"
    if any(token in text.casefold() for token in ("数据", "指标", "增长", "营收", "市场", "财务", "报告", "趋势", "data", "metric")):
        return "data-journalism"
    if any(token in text.casefold() for token in ("发布", "产品", "品牌", "路演", "launch", "campaign")):
        return "vivid-launch"
    if any(token in text.casefold() for token in ("架构", "工程", "技术", "系统", "api", "研发")):
        return "blueprint"
    if any(token in text.casefold() for token in ("培训", "课程", "教程", "课堂", "方法论")):
        return "soft-rounded"
    if has_template:
        return "state-briefing"
    return "state-briefing"


def _resolve_style(job: PresentationJob, reference_text: str, *, has_images: bool, has_template: bool) -> tuple[str, str | None, str]:
    """Resolve style precedence: explicit brief > template context > manual/random choice."""
    try:
        metadata = json.loads(job.metadata_json or "{}")
    except (TypeError, ValueError):
        metadata = {}
    user_text = "\n".join((job.title or "", job.purpose or "", job.brief or ""))
    prompt_style = _style_from_text(user_text)
    if prompt_style:
        return prompt_style, prompt_style, "prompt"
    if str(job.style or "") == "random":
        resolved = _choose_intelligent_style(f"{user_text}\n{reference_text}", has_images=has_images, has_template=has_template)
        return resolved, None, "intelligent"
    resolved = str(job.style or "state-briefing")
    if resolved not in STYLE_PRESETS:
        resolved = "state-briefing"
    return resolved, None, "manual"


def _prompt_disables_native_template(text: str) -> bool:
    """Recognize only explicit requests to abandon an uploaded visual shell."""
    compact = re.sub(r"\s+", "", str(text or "")).casefold()
    return any(
        phrase in compact
        for phrase in (
            "不要使用模板", "不用模板", "忽略模板", "不按模板", "放弃模板",
            "模板仅作内容参考", "模板只作内容参考", "redesignwithouttemplate",
            "ignorethetemplate", "donotusethetemplate",
        )
    )


def _prompt_allows_template_media(text: str) -> bool:
    """Only reuse photos embedded in a template when the user asks for it."""
    compact = re.sub(r"\s+", "", str(text or "")).casefold()
    return any(
        phrase in compact
        for phrase in (
            "沿用模板图片", "使用模板图片", "保留模板图片", "复用模板图片",
            "沿用模板中的图", "使用模板中的图", "保留模板中的图", "复用模板中的图",
            "沿用上传模板图片", "使用上传模板图片", "保留上传模板图片", "复用上传模板图片",
            "沿用上传的模板图片", "使用上传的模板图片", "保留上传的模板图片", "复用上传的模板图片",
            "保留原图", "沿用原图", "复用原图",
            "usetemplatemedia", "reusetemplatemedia", "keeptemplateimages",
        )
    )


def _safe_hex_color(color: Any) -> str | None:
    try:
        rgb = color.rgb
        if rgb is not None:
            return str(rgb).upper()
    except (AttributeError, ValueError, TypeError):
        return None
    return None


def _shape_fill_color(shape: Any) -> str | None:
    try:
        if shape.fill.type is None:
            return None
        return _safe_hex_color(shape.fill.fore_color)
    except (AttributeError, ValueError, TypeError):
        return None


def _relative_luminance(color: str) -> float:
    value = str(color or "").strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        return 0.0
    channels = []
    for index in (0, 2, 4):
        channel = int(value[index : index + 2], 16) / 255
        channels.append(channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _color_contrast(first: str, second: str) -> float:
    light, dark = sorted((_relative_luminance(first), _relative_luminance(second)), reverse=True)
    return (light + 0.05) / (dark + 0.05)


def _contrasting_color(color: str, background: str, *, minimum: float, muted: bool = False) -> str:
    if _color_contrast(color, background) >= minimum:
        return color
    light_background = _relative_luminance(background) > 0.45
    if muted:
        return "4B5563" if light_background else "CBD5E1"
    return "111827" if light_background else "F8FAFC"


def _extract_template_palette(prs: Presentation) -> dict[str, str]:
    """Read only stable visual facts from a native template before clearing slides."""
    palette: dict[str, str] = {}
    candidates = list(prs.slides)
    if candidates:
        try:
            palette["bg"] = _safe_hex_color(candidates[0].background.fill.fore_color) or ""
        except (AttributeError, ValueError, TypeError):
            pass
    if not palette.get("bg"):
        for master in getattr(prs, "slide_masters", []):
            try:
                value = _safe_hex_color(master.background.fill.fore_color)
            except (AttributeError, ValueError, TypeError):
                value = None
            if value:
                palette["bg"] = value
                break
    if not palette.get("bg") and candidates:
        # Many native decks use a full-bleed rectangle instead of a slide fill.
        slide = candidates[0]
        canvas_area = max(int(prs.slide_width) * int(prs.slide_height), 1)
        for shape in slide.shapes:
            if int(getattr(shape, "width", 0)) * int(getattr(shape, "height", 0)) < canvas_area * 0.72:
                continue
            value = _shape_fill_color(shape)
            if value:
                palette["bg"] = value
                break
    for slide in candidates:
        for shape in slide.shapes:
            fill_color = _shape_fill_color(shape)
            if fill_color and fill_color != palette.get("bg") and "accent" not in palette:
                palette["accent"] = fill_color
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_name = str(run.font.name or "").strip()
                    if font_name and "font" not in palette:
                        palette["font"] = font_name
                    text_color = _safe_hex_color(run.font.color)
                    if text_color and "fg" not in palette:
                        palette["fg"] = text_color
                    if text_color and text_color != palette.get("fg") and "accent2" not in palette:
                        palette["accent2"] = text_color
                    if palette.get("font") and palette.get("fg"):
                        break
                if palette.get("font") and palette.get("fg"):
                    break
            if palette.get("font") and palette.get("fg"):
                break
        if palette.get("font") and palette.get("fg") and palette.get("accent"):
            break
    bg = palette.get("bg")
    if bg:
        try:
            red, green, blue = (int(bg[index : index + 2], 16) for index in (0, 2, 4))
            luminance = (0.2126 * red + 0.7152 * green + 0.0722 * blue) / 255
            contrast_fg = "111827" if luminance > 0.58 else "F8FAFC"
            if not palette.get("fg"):
                palette["fg"] = contrast_fg
            else:
                try:
                    fg_luminance = sum(weight * int(palette["fg"][index : index + 2], 16) for index, weight in ((0, 0.2126), (2, 0.7152), (4, 0.0722))) / 255
                    if abs(fg_luminance - luminance) < 0.28:
                        palette["fg"] = contrast_fg
                except (TypeError, ValueError):
                    palette["fg"] = contrast_fg
            if not palette.get("muted"):
                palette["muted"] = "596A7F" if luminance > 0.58 else "B8C2D9"
            else:
                try:
                    muted_luminance = sum(weight * int(palette["muted"][index : index + 2], 16) for index, weight in ((0, 0.2126), (2, 0.7152), (4, 0.0722))) / 255
                    if abs(muted_luminance - luminance) < 0.16:
                        palette["muted"] = "596A7F" if luminance > 0.58 else "B8C2D9"
                except (TypeError, ValueError):
                    palette["muted"] = "596A7F" if luminance > 0.58 else "B8C2D9"
        except (TypeError, ValueError):
            pass
    palette.setdefault("font", "Aptos")
    return {key: value for key, value in palette.items() if value}


def _analyse_native_template(template_path: Path) -> dict[str, Any] | None:
    """Read a PPTX as a native slide library using the project-local skill.

    Importing lazily keeps the normal free-design path independent from the
    skill package, while still making the deployed application self-contained.
    """
    if not template_path.exists() or not PPT_MASTER_SCRIPTS.exists():
        return None
    try:
        script_root = str(PPT_MASTER_SCRIPTS)
        if script_root not in sys.path:
            sys.path.insert(0, script_root)
        from template_fill_pptx.analyzer import analyze_pptx

        library = analyze_pptx(template_path)
        return library if isinstance(library, dict) else None
    except Exception:
        logger.warning("Unable to analyse native PPTX template %s", template_path, exc_info=True)
        return None


def _template_context(library: dict[str, Any] | None) -> str:
    """Create compact, model-readable facts from a native template library."""
    if not library:
        return ""
    chunks: list[str] = []
    for slide in library.get("slides", []):
        if not isinstance(slide, dict):
            continue
        index = slide.get("slide_index")
        page_type = slide.get("page_type") or "content_candidate"
        layout_family = _native_layout_family(slide)
        affordance = _native_layout_affordance(slide)
        slot_roles = Counter(str(slot.get("role") or "") for slot in (slide.get("slots") or []) if isinstance(slot, dict))
        table_bits: list[str] = []
        for table in (slide.get("tables") or [])[:3]:
            if not isinstance(table, dict):
                continue
            rows = table.get("rows") or []
            header_count = 0
            if rows and isinstance(rows[0], dict):
                for cell in (rows[0].get("cells") or [])[:12]:
                    if isinstance(cell, dict) and not cell.get("is_merge_slave"):
                        value = str(cell.get("text") or "").strip()
                        if value:
                            header_count += 1
            table_bits.append(f"table={table.get('table_id')} {table.get('row_count')}x{table.get('column_count')} header_slots={header_count}")
        chunks.append(
            f"P{index}: type={page_type}; family={layout_family}; affordance={affordance}; "
            f"large_pictures={_native_large_picture_count(slide)}; slots={dict(slot_roles)}; "
            f"text_shapes={len(slide.get('slots') or [])}"
            f"{'; ' + '; '.join(table_bits) if table_bits else ''}"
        )
    return _clip_multiline("\n".join(chunks), MAX_TEMPLATE_CONTEXT)


def _normalised_template_text(value: Any) -> str:
    return " ".join(str(value or "").replace("\n", " ").split()).strip().casefold()


def _native_static_texts(library: dict[str, Any]) -> Counter[str]:
    """Find recurring chrome (headers, page numbers, brand labels).

    Recurring text is intentionally left untouched during fill. One-off text
    slots remain available for generated titles and evidence, which prevents a
    template's running header from being overwritten as slide content.
    """
    counts: Counter[str] = Counter()
    for slide in library.get("slides", []):
        if not isinstance(slide, dict):
            continue
        seen_on_slide: set[str] = set()
        for slot in slide.get("slots", []):
            if not isinstance(slot, dict):
                continue
            text = _normalised_template_text(slot.get("text"))
            if text and text not in seen_on_slide:
                counts[text] += 1
                seen_on_slide.add(text)
    return counts


_NATIVE_PLACEHOLDER_WORDS = {
    "title", "subtitle", "heading", "body", "text", "内容", "标题", "副标题",
    "请替换", "点击此处添加文本", "单击此处添加文本", "your title", "lorem ipsum",
}


def _native_is_chrome(slot: dict[str, Any], static_texts: Counter[str], slide_count: int) -> bool:
    text = _normalised_template_text(slot.get("text"))
    if not text:
        return False
    if text in _NATIVE_PLACEHOLDER_WORDS:
        return False
    # Tiny text is normally a running header/page number. Treating every
    # repeated string as chrome is unsafe: real templates often repeat a large
    # section title on several pages, and that title must be replaced when the
    # same visual shell carries a new story.
    geometry = slot.get("geometry") if isinstance(slot.get("geometry"), dict) else {}
    metrics = slot.get("text_metrics") if isinstance(slot.get("text_metrics"), dict) else {}
    font_size = float(metrics.get("font_size_px") or 0)
    width = float(geometry.get("width") or 0)
    height = float(geometry.get("height") or 0)
    y = float(geometry.get("y") or 0)
    if font_size and font_size <= 5:
        return True
    if height and height <= 20 and (font_size <= 12 or width < 260):
        return True
    # A label repeated on several pages is usually a running header or brand
    # marker only when it is small/edge-aligned. Large title-sized text in the
    # content band is a replaceable section heading, not chrome.
    threshold = 2 if slide_count <= 6 else 3
    if static_texts.get(text, 0) >= threshold:
        if font_size >= 24 and y < 180 and width >= 300:
            return False
        if font_size >= 20 and 300 <= y < 590 and width >= 400:
            # Large centered brand/chapter wording is part of the native
            # visual identity; chapter-specific labels sit above it and are
            # the safe replacement target.
            return True
        if font_size >= 20 and 120 < y < 300 and width >= 260:
            return False
        return True
    if re.fullmatch(r"(?:part|section|chapter)?\s*[-/]?\s*\d{1,3}", text):
        return True
    if re.fullmatch(r"part\s+(?:one|two|three|four|five|six|seven|eight|nine|ten)", text):
        return True
    return False


def _native_slot_geometry(slot: dict[str, Any]) -> tuple[float, float, float, float]:
    geometry = slot.get("geometry") if isinstance(slot.get("geometry"), dict) else {}
    return (
        float(geometry.get("x") or 0),
        float(geometry.get("y") or 0),
        float(geometry.get("width") or 0),
        float(geometry.get("height") or 0),
    )


def _native_slot_font_size(slot: dict[str, Any]) -> float:
    metrics = slot.get("text_metrics") if isinstance(slot.get("text_metrics"), dict) else {}
    return float(metrics.get("font_size_px") or 0)


def _native_slot_text(slot: dict[str, Any]) -> str:
    return " ".join(str(slot.get("text") or "").replace("\n", " ").split()).strip()


def _native_slot_is_empty_decor(slot: dict[str, Any]) -> bool:
    """Return true for a blank text container that is probably just a mask.

    The analyzer intentionally exposes empty text containers because they can
    be real placeholders. In native decks, however, most empty containers are
    decorative shapes with no font metrics. They must not receive generated
    labels merely because a semantic field happened to be available.
    """
    text = _native_slot_text(slot)
    if text:
        return False
    _x, _y, width, height = _native_slot_geometry(slot)
    font_size = _native_slot_font_size(slot)
    role = str(slot.get("role") or "")
    if font_size > 0:
        return False
    shape_name = str(slot.get("shape_name") or "").casefold()
    if int(slot.get("group_depth") or 0) > 0:
        return True
    if (width and width <= 24) or (height and height <= 20):
        return True
    if shape_name and not any(token in shape_name for token in ("文本框", "text box", "placeholder", "占位")):
        return True
    # Keep a genuinely small blank title/label placeholder available, but do
    # not write into full-slide polygons or image masks.
    if "title" in role or "label" in role:
        return width >= 720 or height >= 240
    return True


def _native_slot_is_body_like(slot: dict[str, Any]) -> bool:
    """Recognize visual body frames even when the analyzer labels them badly.

    PowerPoint authors often use ordinary text boxes for both headings and
    paragraphs. The intake analyzer therefore has to infer a role and can mark
    a wide multi-line explanation as ``title_candidate``. Source geometry and
    visible copy are stronger signals for those frames than the inferred role.
    """
    role = str(slot.get("role") or "")
    if "body" in role:
        return True
    if not any(token in role for token in ("title", "label")):
        return False
    text = _native_slot_text(slot)
    if not text:
        return False
    _x, y, width, height = _native_slot_geometry(slot)
    if y < 100 or width < 420 or height < 48:
        return False
    paragraph_count = int(slot.get("paragraph_count") or 0)
    return len(text) >= 12 or paragraph_count >= 2 or height >= 120


def _native_slot_is_fillable(slot: dict[str, Any], static_texts: Counter[str], slide_count: int) -> bool:
    if _native_slot_is_empty_decor(slot):
        return False
    if _native_is_chrome(slot, static_texts, slide_count):
        return False
    role = str(slot.get("role") or "")
    return any(token in role for token in ("title", "body", "label"))


def _native_slots_in_reading_order(slots: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(
        slots,
        key=lambda slot: (
            _native_slot_geometry(slot)[1],
            _native_slot_geometry(slot)[0],
            str(slot.get("slot_id") or ""),
        ),
    )


def _native_slots_by_column(slots: list[dict[str, Any]]) -> list[list[dict[str, Any]]]:
    """Cluster slots into left/right/column groups using source geometry."""
    ordered = _native_slots_in_reading_order(slots)
    columns: list[list[dict[str, Any]]] = []
    for slot in ordered:
        x = _native_slot_geometry(slot)[0]
        placed = False
        for column in columns:
            anchor = _native_slot_geometry(column[0])[0]
            if abs(x - anchor) <= 180:
                column.append(slot)
                placed = True
                break
        if not placed:
            columns.append([slot])
    return sorted(columns, key=lambda column: _native_slot_geometry(column[0])[0])


def _native_row_pairs(slots: list[dict[str, Any]]) -> list[tuple[dict[str, Any], dict[str, Any]]]:
    """Pair short left labels with wide explanations on the same visual row."""
    labels = [
        slot
        for slot in slots
        if any(token in str(slot.get("role") or "") for token in ("title", "label"))
        and not _native_slot_is_body_like(slot)
        and (_native_slot_text(slot) or _native_slot_font_size(slot) > 0)
        and _native_slot_geometry(slot)[1] >= 100
        and 0 < _native_slot_geometry(slot)[2] <= 420
    ]
    bodies = [
        slot
        for slot in slots
        if (
            _native_slot_is_body_like(slot)
            or (
                _native_slot_text(slot)
                and any(token in str(slot.get("role") or "") for token in ("title", "label"))
                and _native_slot_geometry(slot)[1] >= 100
                and _native_slot_geometry(slot)[2] >= 420
                and _native_slot_geometry(slot)[3] >= 48
            )
        )
        and _native_slot_geometry(slot)[1] >= 100
        and _native_slot_geometry(slot)[2] >= 420
    ]
    remaining = list(bodies)
    pairs: list[tuple[dict[str, Any], dict[str, Any]]] = []
    for label in _native_slots_in_reading_order(labels):
        lx, ly, lw, lh = _native_slot_geometry(label)
        label_center = ly + lh / 2
        matches: list[tuple[float, float, dict[str, Any]]] = []
        for body in remaining:
            bx, by, bw, bh = _native_slot_geometry(body)
            if bx <= lx + max(40, lw * 0.55) or bw < max(420, lw * 1.5):
                continue
            center_gap = abs((by + bh / 2) - label_center)
            tolerance = max(42.0, min(78.0, max(lh, bh) * 0.62))
            if center_gap <= tolerance:
                matches.append((center_gap, bx - (lx + lw), body))
        if not matches:
            continue
        body = min(matches, key=lambda item: (item[0], item[1]))[2]
        remaining.remove(body)
        pairs.append((label, body))
    return pairs


def _native_slot_score(slot: dict[str, Any], *, title: bool = False) -> float:
    geometry = slot.get("geometry") if isinstance(slot.get("geometry"), dict) else {}
    width = float(geometry.get("width") or 0)
    height = float(geometry.get("height") or 0)
    metrics = slot.get("text_metrics") if isinstance(slot.get("text_metrics"), dict) else {}
    font_size = float(metrics.get("font_size_px") or 0)
    role = str(slot.get("role") or "")
    has_source_text = bool(str(slot.get("text") or "").strip())
    score = width * height / 1000.0 + font_size * (30.0 if title else 12.0)
    score += 2000 if has_source_text else -400
    if title:
        if "title" in role:
            score += 400
        if "label" in role:
            score += 100
        if float(geometry.get("y") or 0) < 110:
            score += 150
    else:
        if "body" in role:
            score += 800
        if "label" in role:
            score += 180
    return score


def _fit_native_text(text: str, slot: dict[str, Any]) -> str:
    """Fit copy to a source slot using geometry instead of shrinking fonts."""
    raw = str(text or "").strip()
    if not raw:
        return ""
    geometry = slot.get("geometry") if isinstance(slot.get("geometry"), dict) else {}
    metrics = slot.get("text_metrics") if isinstance(slot.get("text_metrics"), dict) else {}
    width = max(float(geometry.get("width") or 0), 80.0)
    height = max(float(geometry.get("height") or 0), 24.0)
    font_size = float(metrics.get("font_size_px") or 16.0)
    # Analyzer geometry is in pixels while PowerPoint's East Asian glyphs are
    # close to one em wide. A minimum of six characters overestimates narrow
    # chapter labels and lets text overflow the native frame; use the actual
    # measured width with a small floor instead.
    chars_per_line = max(3, int(width / max(font_size * 0.95, 8.0)))
    line_count = max(1, int(height / max(font_size * 1.35, 12.0)))
    capacity = max(4, min(900, chars_per_line * line_count))
    if len(raw) <= capacity:
        return raw
    return raw[: max(1, capacity - 1)].rstrip() + "…"


def _native_replacement_texts(slide_spec: dict[str, Any]) -> list[str]:
    """Flatten semantic slide fields into short, template-friendly strings."""
    values: list[str] = []
    for key in ("body", "summary"):
        value = str(slide_spec.get(key) or "").strip()
        if value:
            values.append(_clip(value, 520))
    bullets = slide_spec.get("bullets") if isinstance(slide_spec.get("bullets"), list) else []
    values.extend(_clip(str(item), 160) for item in bullets if str(item).strip())
    for key in ("left_bullets", "right_bullets"):
        items = slide_spec.get(key) if isinstance(slide_spec.get(key), list) else []
        values.extend(_clip(str(item), 160) for item in items if str(item).strip())
    steps = slide_spec.get("steps") if isinstance(slide_spec.get("steps"), list) else []
    for step in steps:
        if not isinstance(step, dict):
            continue
        title = str(step.get("title") or "").strip()
        body = str(step.get("body") or "").strip()
        if title:
            values.append(_clip(title, 100))
        if body:
            values.append(_clip(body, 160))
    return values


_GENERIC_NATIVE_TITLE_RE = re.compile(
    r"^(?:第\s*\d+\s*页?(?:核心判断|要点|内容)?|核心判断|页面(?:标题|内容)?|正文|内容|未命名(?:页面)?|slide\s*\d+)$",
    flags=re.IGNORECASE,
)


def _native_title_candidate(slide_spec: dict[str, Any]) -> str:
    """Return a concise semantic title for native section shells.

    Model plans sometimes contain an indexing placeholder such as ``第 12 页
    核心判断``. That text is useful as an internal fallback label but looks
    unfinished when enlarged on a chapter divider. Prefer the page's summary,
    body, or first bullet only for those obvious placeholders; authored titles
    always win.
    """
    title = _clip(str(slide_spec.get("title") or ""), 90).strip()
    compact = re.sub(r"\s+", "", title)
    if title and not _GENERIC_NATIVE_TITLE_RE.fullmatch(compact):
        return title
    for key in ("summary", "body"):
        value = _clip(str(slide_spec.get(key) or ""), 90).strip()
        if value and not _GENERIC_NATIVE_TITLE_RE.fullmatch(re.sub(r"\s+", "", value)):
            return value
    bullets = slide_spec.get("bullets") if isinstance(slide_spec.get("bullets"), list) else []
    for value in bullets:
        candidate = _clip(str(value or ""), 90).strip()
        if candidate and not _GENERIC_NATIVE_TITLE_RE.fullmatch(re.sub(r"\s+", "", candidate)):
            return candidate
    return title


def _native_chapter_title(slides: list[dict[str, Any]], index: int) -> str:
    """Choose a nearby real topic when a chapter title is an index placeholder."""
    current = _native_title_candidate(slides[index]) if 0 <= index < len(slides) else ""
    original = _clip(str(slides[index].get("title") or ""), 90) if 0 <= index < len(slides) else ""
    if original and not _GENERIC_NATIVE_TITLE_RE.fullmatch(re.sub(r"\s+", "", original)):
        return original
    if current and current != original and not _GENERIC_NATIVE_TITLE_RE.fullmatch(re.sub(r"\s+", "", current)):
        return current
    for distance in range(1, 5):
        for candidate_index in (index + distance, index - distance):
            if not 0 <= candidate_index < len(slides):
                continue
            candidate = _native_title_candidate(slides[candidate_index])
            if candidate and not _GENERIC_NATIVE_TITLE_RE.fullmatch(re.sub(r"\s+", "", candidate)):
                return candidate
    return current or original


def _native_compact_chapter_title(value: str, max_chars: int = 7) -> str:
    """Fit a chapter label into the narrow native divider title slot.

    Native chapter dividers in the uploaded deck use a large display face in a
    roughly seven-CJK-character box.  Appending an ellipsis to a long model
    title makes the page look unfinished, so prefer the first complete semantic
    phrase (usually the text before ``与`` / ``及``) and only then fall back to
    a hard, punctuation-free trim.
    """
    text = " ".join(str(value or "").split()).strip()
    if not text:
        return ""
    if len(text) <= max_chars:
        return text
    for separator in ("与", "及", "：", ":", "—", "-", "/", "、"):
        head = text.split(separator, 1)[0].strip()
        if 2 <= len(head) <= max_chars:
            return head
    # Keep a whole-word prefix for Latin headings where character count is a
    # poor proxy for visual width.
    words = text.split(" ")
    if len(words) > 1:
        prefix = " ".join(words[:2]).strip()
        if len(prefix) <= max_chars:
            return prefix
    return text[:max_chars].rstrip(" ，,。；;:：-/—")


def _native_compact_side_label(value: Any, fallback: str = "内容提要", max_chars: int = 6) -> str:
    """Turn a sentence-like heading into a short native side-tab label."""
    text = " ".join(str(value or "").split()).strip()
    if not text:
        return fallback
    clause = re.split(r"[，,。；;：:、/]|(?:需要|应当|应该|集中|按照|按|支持|形成|实现|覆盖)", text, maxsplit=1)[0].strip()
    candidate = clause if 2 <= len(clause) <= max_chars else text
    return _native_compact_chapter_title(candidate, max_chars=max_chars) or fallback


def _native_closing_text(value: str, slot: dict[str, Any] | None) -> str:
    """Keep a native closing label on one line at the source font size.

    Closing pages often use display type inside a deliberately narrow frame.
    A generic four-character phrase can technically fit the frame height while
    wrapping into two vertical-looking lines. Preserve authored short closing
    copy when it fits; otherwise use a compact two-character sign-off instead
    of shrinking the template typography.
    """
    text = " ".join(str(value or "").split()).strip()
    if not slot:
        return text if 0 < len(text) <= 6 else "谢谢"
    _, _, width, _ = _native_slot_geometry(slot)
    font_size = max(_native_slot_font_size(slot), 16.0)
    one_line_capacity = max(2, int(width / max(font_size * 0.95, 8.0)))
    if text and len(text) <= one_line_capacity:
        return text
    return "谢谢"


def _native_side_label_text(slide_spec: dict[str, Any], title_text: str) -> str:
    kind = str(slide_spec.get("kind") or "content").casefold()
    layout = str(slide_spec.get("layout_id") or "").casefold()
    chart = slide_spec.get("chart") if isinstance(slide_spec.get("chart"), dict) else {}
    if kind == "chart" or "chart" in layout:
        return _native_compact_side_label(chart.get("series_name"), "数据分析", max_chars=7)
    kicker = str(slide_spec.get("kicker") or "").strip()
    if kicker and not re.fullmatch(r"(?:第\s*\d+\s*章|part\s+\w+)", kicker, flags=re.I):
        return _native_compact_side_label(kicker, max_chars=7)
    return _native_compact_side_label(title_text)


def _native_replacements(
    slide_entry: dict[str, Any],
    slide_spec: dict[str, Any],
    static_texts: Counter[str],
    slide_count: int,
) -> list[dict[str, Any]]:
    """Map semantic fields to the *right* native slots for each page family.

    A generic ``title -> body -> remaining values`` pass looks reasonable on a
    blank template, but real decks contain empty masks, repeated section
    headers, card labels, and table graphic frames. The family-aware pass below
    keeps those structures intact and only writes into slots whose geometry
    matches the intended role. This is the main guard against a filled deck
    looking like the same control panel on every page.
    """
    slots = [slot for slot in (slide_entry.get("slots") or []) if isinstance(slot, dict)]
    usable = [slot for slot in slots if _native_slot_is_fillable(slot, static_texts, slide_count)]
    if not usable:
        return []

    family = _native_layout_family(slide_entry)
    page_role = _native_page_role(slide_entry)
    kind = str(slide_spec.get("kind") or "content").casefold()
    used: set[str] = set()
    replacements: list[dict[str, Any]] = []

    def slot_id(slot: dict[str, Any]) -> str:
        return str(slot.get("slot_id") or "")

    def available(predicate: Any) -> list[dict[str, Any]]:
        return [slot for slot in usable if slot_id(slot) not in used and predicate(slot)]

    def assign(slot: dict[str, Any] | None, text: Any, *, limit: int = 520) -> bool:
        if not slot:
            return False
        raw = str(text or "").strip()
        if not raw:
            return False
        fitted = _fit_native_text(_clip(raw, limit), slot)
        if not fitted:
            return False
        target = slot_id(slot)
        if not target or target in used:
            return False
        used.add(target)
        replacements.append({"slot_id": target, "text": fitted})
        return True

    def take(predicate: Any, *, title: bool = False) -> dict[str, Any] | None:
        options = available(predicate)
        return max(options, key=lambda value: _native_slot_score(value, title=title), default=None)

    def role_is(slot: dict[str, Any], *roles: str) -> bool:
        role = str(slot.get("role") or "")
        return any(token in role for token in roles)

    def y_of(slot: dict[str, Any]) -> float:
        return _native_slot_geometry(slot)[1]

    def x_of(slot: dict[str, Any]) -> float:
        return _native_slot_geometry(slot)[0]

    title_text = _clip(_native_title_candidate(slide_spec), 120)
    subtitle = _clip(str(slide_spec.get("subtitle") or ""), 220)
    summary = _clip(str(slide_spec.get("summary") or ""), 260)
    body = _clip(str(slide_spec.get("body") or ""), 520)
    bullets = [
        _clip(str(value), 150)
        for value in (slide_spec.get("bullets") or [])
        if str(value).strip()
    ][:8]
    steps = [
        {
            "title": _clip(str(value.get("title") or ""), 100),
            "body": _clip(str(value.get("body") or ""), 180),
        }
        for value in (slide_spec.get("steps") or [])
        if isinstance(value, dict)
        and (str(value.get("title") or "").strip() or str(value.get("body") or "").strip())
    ][:5]
    is_timeline = (
        kind == "timeline"
        or str(slide_spec.get("layout_id") or "").casefold() in {"timeline", "roadmap"}
        or str(slide_spec.get("visual_intent") or "").casefold() == "process"
    )

    # Cover / chapter / ending shells have intentionally sparse copy. Never
    # feed their empty decorative containers with comparison or metric fields.
    if family == "cover" or kind == "cover":
        cover_title = take(lambda slot: role_is(slot, "title", "label") and y_of(slot) > 120, title=True)
        assign(cover_title, title_text, limit=90)
        subtitle_slot = take(lambda slot: role_is(slot, "body", "label") and y_of(slot) > 350)
        cover_support = subtitle or summary or body
        cover_support = re.sub(r"^面向(.{2,8})的", r"\1", cover_support)
        assign(subtitle_slot, cover_support, limit=180)
    elif family == "chapter" or page_role == "chapter" or kind in {"section", "chapter", "section-divider"}:
        # The large centered airport/brand heading is usually repeated chrome;
        # the shorter label above it is the replaceable chapter title.
        chapter_slots = available(lambda slot: role_is(slot, "title", "label") and y_of(slot) >= 180)
        chapter_slot = max(chapter_slots, key=lambda value: _native_slot_score(value, title=True), default=None)
        assign(chapter_slot, _native_compact_chapter_title(title_text), limit=24)
        if chapter_slot:
            cx, cy, cw, ch = _native_slot_geometry(chapter_slot)
            for candidate in slots:
                if candidate is chapter_slot or not _native_slot_text(candidate):
                    continue
                sx, sy, sw, sh = _native_slot_geometry(candidate)
                overlaps_title = (
                    min(cx + cw, sx + sw) - max(cx, sx) > 40
                    and min(cy + ch, sy + sh) - max(cy, sy) > 20
                )
                if (
                    overlaps_title
                    and sw >= 400
                    and _native_slot_font_size(candidate) >= 20
                    and _native_is_chrome(candidate, static_texts, slide_count)
                ):
                    target = slot_id(candidate)
                    if target and target not in used:
                        used.add(target)
                        replacements.append({"slot_id": target, "text": ""})
        # Chapter counters such as “Part Three” are deliberately classified as
        # recurring chrome, so they are absent from ``usable``. They still
        # belong to the chapter's semantic identity and must be rewritten to
        # avoid leaking the source deck's section numbering.
        part_candidates = [
            slot for slot in slots
            if role_is(slot, "label", "title")
            and re.search(r"(?:part|section|chapter|第|部分)", _native_slot_text(slot), re.I)
        ]
        part_slot = max(part_candidates, key=lambda value: _native_slot_score(value, title=False), default=None)
        part_text = str(slide_spec.get("kicker") or "").strip()
        if part_text and part_slot:
            assign(part_slot, part_text, limit=32)
    elif _native_is_agenda(slide_entry) or kind in {"agenda", "toc"}:
        agenda_title = take(lambda slot: role_is(slot, "title", "label") and y_of(slot) < 180, title=True)
        agenda_heading = title_text or "目录"
        if "目录" in agenda_heading and len(agenda_heading) > 2:
            agenda_heading = "目录"
        assign(agenda_title, agenda_heading, limit=40)
        agenda_items = _native_slots_in_reading_order(
            available(lambda slot: role_is(slot, "label") and y_of(slot) >= 130)
        )
        # Native agenda shells often carry four or six visual bars. Fill an
        # unprovided bar with a neutral management heading rather than leave
        # a conspicuous empty control in the exported deck.
        neutral_items = ["背景与范围", "方案与风险", "测算与建议", "实施计划与请示事项"]
        # When the model supplies fewer headings than the native shell has
        # bars, append only the *remaining* canonical headings.  Starting at
        # ``len(bullets)`` avoids repeating a near-identical first heading such
        # as “项目背景与范围” in the fourth bar.
        agenda_values = bullets + [
            value for value in neutral_items[len(bullets):]
            if value not in bullets
        ]
        for slot, value in zip(agenda_items, agenda_values[: len(agenda_items)]):
            assign(slot, value, limit=48)
    elif family == "ending" or kind in {"closing", "ending"}:
        closing_slot = take(lambda slot: role_is(slot, "title", "label") and y_of(slot) >= 100, title=True)
        closing_text = _native_closing_text(title_text, closing_slot)
        assign(closing_slot, closing_text, limit=16)
    else:
        # Ordinary page heading. Prefer the top band; this prevents a large
        # body/caption slot from stealing the title on photo and process pages.
        title_options = available(
            lambda slot: role_is(slot, "title", "label")
            and not _native_slot_is_body_like(slot)
            and y_of(slot) < 180
        )
        def title_rank(value: dict[str, Any]) -> float:
            y = y_of(value)
            # Native templates put the page heading in a narrow top band. A
            # large body/caption frame can otherwise outscore that heading by
            # area alone (especially on hero/photo pages).
            band_bonus = 5200 if y < 80 else 1200 if y < 125 else 0
            return _native_slot_score(value, title=True) + band_bonus
        title_slot = max(title_options, key=title_rank, default=None)
        if title_slot is None:
            title_slot = take(
                lambda slot: role_is(slot, "title", "label") and not _native_slot_is_body_like(slot),
                title=True,
            )
        if title_slot is None:
            title_slot = take(lambda slot: role_is(slot, "title", "label"), title=True)
        assign(title_slot, title_text, limit=96)

        # Table graphic frames own their text nodes. Replacing the frame as a
        # normal body slot corrupts the table or duplicates its contents; the
        # dedicated table editor below handles those cells instead.
        table_slots = bool(slide_entry.get("tables"))
        if family.startswith("table-") or table_slots:
            if not table_slots:
                body_slot = take(lambda slot: role_is(slot, "body"), title=False)
                assign(body_slot, body or summary, limit=520)
        elif is_timeline and steps:
            remaining = [slot for slot in usable if slot_id(slot) not in used]
            row_pairs = _native_row_pairs(remaining)
            for step, (label_slot, body_slot) in zip(steps, row_pairs):
                heading = str(step.get("title") or "").strip()
                detail = str(step.get("body") or "").strip()
                assign(label_slot, heading or detail, limit=72)
                assign(body_slot, detail or heading, limit=180)

            remaining_steps = steps[len(row_pairs) :]
            remaining = [slot for slot in usable if slot_id(slot) not in used]
            label_slots = _native_slots_in_reading_order(
                [
                    slot for slot in remaining
                    if role_is(slot, "label", "title") and not _native_slot_is_body_like(slot)
                ]
            )
            body_slots = _native_slots_in_reading_order(
                [slot for slot in remaining if _native_slot_is_body_like(slot)]
            )
            # Native timelines vary between title/body pairs, body-only rows,
            # and compact labels. Keep each generated step together rather
            # than flattening all titles before all explanations.
            for step_index, step in enumerate(remaining_steps):
                heading = str(step.get("title") or "").strip()
                detail = str(step.get("body") or "").strip()
                label_slot = label_slots[step_index] if step_index < len(label_slots) else None
                body_slot = body_slots[step_index] if step_index < len(body_slots) else None
                if label_slot and body_slot:
                    assign(label_slot, heading or detail, limit=72)
                    assign(body_slot, detail or heading, limit=180)
                elif body_slot:
                    assign(body_slot, "\n".join(value for value in (heading, detail) if value), limit=220)
                elif label_slot:
                    assign(label_slot, "\n".join(value for value in (heading, detail) if value), limit=140)
        elif family == "split" or kind in {"comparison", "split"}:
            remaining = [slot for slot in usable if slot_id(slot) not in used]
            left_heading = str(slide_spec.get("left_title") or "现状").strip()
            right_heading = str(slide_spec.get("right_title") or "目标").strip()
            left_values = [str(item) for item in (slide_spec.get("left_bullets") or bullets[:2]) if str(item).strip()]
            right_values = [str(item) for item in (slide_spec.get("right_bullets") or bullets[2:4]) if str(item).strip()]
            row_pairs = _native_row_pairs(remaining)
            if len(row_pairs) >= 3:
                units: list[tuple[str, str]] = []
                for heading, values in ((left_heading, left_values), (right_heading, right_values)):
                    for item_index, value in enumerate(values, start=1):
                        label = heading if item_index == 1 else _native_compact_side_label(value, f"{heading}{item_index}")
                        units.append((label, value))
                for (label, value), (label_slot, body_slot) in zip(units, row_pairs):
                    assign(label_slot, label, limit=48)
                    assign(body_slot, value, limit=180)
            else:
                columns = _native_slots_by_column(remaining)
                left_slots = columns[0] if columns else []
                right_slots = columns[1] if len(columns) > 1 else []
                for column, heading, values in (
                    (left_slots, left_heading, left_values),
                    (right_slots, right_heading, right_values),
                ):
                    # A native column may mix label and body frames. Reading
                    # order is the reliable semantic signal; splitting by role
                    # can put a heading halfway down the column and leave the
                    # actual body frame blank.
                    ordered_column = _native_slots_in_reading_order(column)
                    assign(ordered_column[0] if ordered_column else None, heading, limit=48)
                    for slot, value in zip(ordered_column[1:], values[:4]):
                        assign(slot, value, limit=160)
        elif kind == "cards":
            remaining = [slot for slot in usable if slot_id(slot) not in used]
            row_pairs = _native_row_pairs(remaining)
            for item_index, (value, (label_slot, body_slot)) in enumerate(zip(bullets, row_pairs), start=1):
                assign(label_slot, _native_compact_side_label(value, f"要点{item_index}"), limit=48)
                assign(body_slot, value, limit=180)
            remaining_values = bullets[len(row_pairs) :]
            remaining = [slot for slot in usable if slot_id(slot) not in used]
            # A card bullet is one semantic unit. Prefer the larger body frames;
            # filling every label before every body mismatches headings and
            # evidence and creates half-empty controls.
            body_slots = _native_slots_in_reading_order(
                [slot for slot in remaining if _native_slot_is_body_like(slot)]
            )
            label_slots = _native_slots_in_reading_order(
                [
                    slot for slot in remaining
                    if role_is(slot, "label", "title") and not _native_slot_is_body_like(slot)
                ]
            )
            card_slots = body_slots if len(body_slots) >= len(remaining_values) else label_slots
            for slot, value in zip(card_slots, remaining_values):
                assign(slot, value, limit=180)
        elif family.startswith("grid-") or family in {"tri-column", "stack"}:
            ordered = _native_slots_in_reading_order([slot for slot in usable if slot_id(slot) not in used])
            values = bullets or [summary or body]
            row_pairs = _native_row_pairs(ordered)
            for item_index, (value, (label_slot, body_slot)) in enumerate(zip(values, row_pairs), start=1):
                assign(label_slot, _native_compact_side_label(value, f"要点{item_index}"), limit=48)
                assign(body_slot, value, limit=220)
            ordered = _native_slots_in_reading_order([slot for slot in usable if slot_id(slot) not in used])
            labels = [
                slot for slot in ordered
                if role_is(slot, "label", "title") and not _native_slot_is_body_like(slot)
            ]
            bodies = [slot for slot in ordered if _native_slot_is_body_like(slot)]
            remaining_values = values[len(row_pairs) :]
            # Card headings consume short labels; body copy consumes body
            # containers. Do not write left/right/metric fields unless the
            # source geometry actually offers a corresponding slot.
            for slot, value in zip(labels, remaining_values[: len(labels)]):
                assign(slot, value, limit=100)
            body_values = ([body] if body else []) + ([summary] if summary else []) + remaining_values[len(labels) :]
            for slot, value in zip(bodies, [item for item in body_values if item][: len(bodies)]):
                assign(slot, value, limit=220)
        else:
            body_slots = _native_slots_in_reading_order(
                available(_native_slot_is_body_like)
            )
            label_slots = _native_slots_in_reading_order(
                available(lambda slot: role_is(slot, "label", "title") and not _native_slot_is_body_like(slot))
            )
            values = ([body] if body else []) + ([summary] if summary else []) + bullets
            # Hero/single pages usually have one body frame; stack only as many
            # values as there are real source frames, dropping overflow instead
            # of concatenating unrelated bullets into a tiny slot.
            for slot, value in zip(body_slots, [item for item in values if item][: len(body_slots)]):
                assign(slot, value, limit=520 if slot is body_slots[0] else 180)
            if not body_slots:
                for slot, value in zip(label_slots, [item for item in values if item][: len(label_slots)]):
                    assign(slot, value, limit=140)
            elif family in {"hero-body", "single"} and (
                kind == "chart"
                or "chart" in str(slide_spec.get("layout_id") or "").casefold()
                or not (slide_spec.get("metric") or slide_spec.get("metric_label"))
            ):
                side_slot = next((slot for slot in label_slots if slot_id(slot) not in used and y_of(slot) >= 100), None)
                assign(side_slot, _native_side_label_text(slide_spec, title_text), limit=56)

        # Metrics are only written when the native shell exposes a real value
        # slot. This avoids putting a number in an arbitrary empty decoration.
        metric = str(slide_spec.get("metric") or "").strip()
        metric_label = str(slide_spec.get("metric_label") or "").strip()
        if metric or metric_label:
            metric_slots = _native_slots_in_reading_order(
                available(
                    lambda slot: role_is(slot, "label", "title")
                    and not _native_slot_is_body_like(slot)
                    and _native_slot_font_size(slot) >= 20
                )
            )
            if metric_slots:
                assign(metric_slots[0], metric or metric_label, limit=32)
            if metric_label and len(metric_slots) > 1:
                assign(metric_slots[1], metric_label, limit=70)

    page_marker = str(slide_spec.get("_native_page_marker") or "").strip()
    if page_marker:
        marker_candidates = [
            slot
            for slot in slots
            if re.fullmatch(r"\d{1,2}", _native_slot_text(slot))
            and _native_slot_geometry(slot)[1] < 90
            and 40 <= _native_slot_geometry(slot)[2] <= 140
            and 35 <= _native_slot_geometry(slot)[3] <= 100
            and _native_slot_font_size(slot) >= 28
        ]
        marker_slot = max(
            marker_candidates,
            key=lambda slot: _native_slot_font_size(slot),
            default=None,
        )
        assign(marker_slot, page_marker, limit=3)

    # A raw PPTX template normally contains example copy. Clear only the
    # replaceable text frames that were not used; chrome and non-text artwork
    # remain untouched. Table/SmartArt text is handled by its own safety path.
    for slot in usable:
        target = slot_id(slot)
        if target in used or not _native_slot_text(slot):
            continue
        if slide_entry.get("tables") and role_is(slot, "body"):
            continue
        if slide_entry.get("diagrams") and role_is(slot, "body"):
            continue
        replacements.append({"slot_id": target, "text": ""})
    return replacements


def _native_layout_family(item: dict[str, Any]) -> str:
    """Return a compact geometry family for one analyzed source slide."""
    page_type = str(item.get("page_type") or "").casefold()
    if "cover" in page_type:
        return "cover"
    if "ending" in page_type or "closing" in page_type:
        return "ending"
    summary = re.sub(r"\s+", "", str(item.get("text_summary") or "").casefold())
    if "toc" in page_type or "目录" in summary or "agenda" in summary or "contents" in summary:
        return "agenda"
    tables = [table for table in (item.get("tables") or []) if isinstance(table, dict)]
    if tables:
        table = max(tables, key=lambda value: int(value.get("row_count") or 0) * int(value.get("column_count") or 0))
        return f"table-{int(table.get('column_count') or 0)}-{min(int(table.get('row_count') or 0), 8)}"
    if "chapter" in page_type:
        return "chapter"
    slots = [slot for slot in (item.get("slots") or []) if isinstance(slot, dict)]
    content_slots = [
        slot
        for slot in slots
        if _native_slot_geometry(slot)[1] >= 100 and not _native_slot_is_empty_decor(slot)
    ]
    body_slots = [slot for slot in content_slots if _native_slot_is_body_like(slot)]
    label_slots = [
        slot
        for slot in content_slots
        if any(token in str(slot.get("role") or "") for token in ("label", "title"))
        and not _native_slot_is_body_like(slot)
    ]
    # Formal report templates often express a process as repeated short tabs
    # on the left with full explanations on the right. Three or more aligned
    # pairs form a vertical stack even if the analyzer called most boxes labels.
    if len(_native_row_pairs(content_slots)) >= 3:
        return "stack"
    # A full-width caption plus one visual label is an image/hero shell, not a
    # true two-column comparison. Treating it as ``split`` sends left/right
    # copy into unrelated caption slots and makes the result look broken.
    if len(body_slots) == 1:
        body_width = _native_slot_geometry(body_slots[0])[2]
        if _native_large_picture_count(item):
            return "photo-body"
        if body_width > 500 and len(label_slots) <= 2:
            return "hero-body"
    # Cluster content columns by their x origin. The analyzer reports pixels,
    # so a 180px gap is a stable threshold across 16:9 decks.
    x_origins = sorted(
        float((slot.get("geometry") or {}).get("x") or 0)
        for slot in body_slots + label_slots
        if slot.get("text") or "body" in str(slot.get("role") or "")
    )
    clusters: list[list[float]] = []
    for x_origin in x_origins:
        if not clusters or x_origin - clusters[-1][-1] > 180:
            clusters.append([x_origin])
        else:
            clusters[-1].append(x_origin)
    if len(label_slots) >= 6 and len(body_slots) >= 3:
        return f"grid-{min(max(len(clusters), 3), 6)}"
    if len(clusters) >= 4:
        return f"grid-{min(len(clusters), 6)}"
    if len(clusters) == 3:
        return "tri-column"
    if len(clusters) == 2:
        return "split"
    if len(body_slots) >= 3:
        return "stack"
    if len(body_slots) == 1:
        geometry = body_slots[0].get("geometry") if isinstance(body_slots[0].get("geometry"), dict) else {}
        if _native_large_picture_count(item):
            return "photo-body"
        if float(geometry.get("width") or 0) > 750:
            return "hero-body"
    return "single"


def _native_layout_unit_capacity(item: dict[str, Any]) -> int:
    """Estimate how many independent content units a native shell exposes.

    Geometry families alone are not enough for layout matching: a two-row,
    six-card page and a one-row, three-card page can both be ``grid-3``. Count
    the visible content frames below the running header so card/timeline pages
    prefer a shell with the same number of units.
    """
    family = _native_layout_family(item)
    if family in {"cover", "agenda", "chapter", "ending", "photo-body", "hero-body", "single"}:
        return 1
    if family.startswith("table-"):
        return 1

    slots = [
        slot
        for slot in (item.get("slots") or [])
        if isinstance(slot, dict)
        and _native_slot_geometry(slot)[1] >= 100
        and not _native_slot_is_empty_decor(slot)
    ]
    label_slots = [
        slot
        for slot in slots
        if any(token in str(slot.get("role") or "") for token in ("label", "title"))
        and not _native_slot_is_body_like(slot)
    ]
    body_slots = [slot for slot in slots if _native_slot_is_body_like(slot)]
    row_pairs = _native_row_pairs(slots)

    if family == "split":
        if len(row_pairs) >= 2:
            return len(row_pairs)
        return max(1, min(len(_native_slots_by_column(label_slots + body_slots)), 4))
    if family.startswith("grid-"):
        try:
            minimum = int(family.split("-", 1)[1])
        except (TypeError, ValueError):
            minimum = 1
        return max(minimum, len(label_slots), len(body_slots), 1)
    if family == "tri-column":
        return max(3, len(label_slots), len(body_slots))
    if family == "stack":
        return len(row_pairs) if row_pairs else max(len(label_slots), len(body_slots), 1)
    return 1


def _native_large_picture_count(item: dict[str, Any]) -> int:
    """Count meaningful page imagery while ignoring small logos and badges."""
    pictures = item.get("pictures") if isinstance(item.get("pictures"), list) else []
    count = 0
    for picture in pictures:
        if not isinstance(picture, dict):
            continue
        geometry = picture.get("geometry") if isinstance(picture.get("geometry"), dict) else {}
        width = float(geometry.get("width") or 0)
        height = float(geometry.get("height") or 0)
        if width >= 260 and height >= 150 and width * height >= 60_000:
            count += 1
    return count


def _native_layout_affordance(item: dict[str, Any]) -> str:
    """Describe what a native page can communicate, not its example copy."""
    family = _native_layout_family(item)
    labels = {
        "cover": "航拍或品牌封面",
        "agenda": "四段式汇报目录",
        "chapter": "全宽图片章节分隔",
        "ending": "品牌收束与致谢",
        "photo-body": "图片证据加短结论",
        "hero-body": "单一结论或图文说明",
        "split": "双栏对比或两组责任",
        "tri-column": "三类并列信息",
        "stack": "纵向步骤或分层说明",
        "single": "简洁正文或单点判断",
    }
    if family.startswith("table-"):
        return "原生数据表或明细清单"
    if family.startswith("grid-"):
        return "多指标卡片或多对象横向比较"
    return labels.get(family, "通用内容页")


def _native_page_role(item: dict[str, Any]) -> str:
    """Collapse template analysis into the roles that control deck rhythm."""
    family = _native_layout_family(item)
    if family in {"cover", "agenda", "chapter", "ending"}:
        return family
    if family.startswith("table-") or item.get("charts"):
        return "data"
    return "content"


def _native_is_agenda(item: dict[str, Any]) -> bool:
    summary = re.sub(r"\s+", "", str(item.get("text_summary") or "").casefold())
    return "目录" in summary or "agenda" in summary or "contents" in summary


def _native_sequence_sources(library: dict[str, Any], target_count: int) -> list[dict[str, Any]]:
    """Scale the source deck's visual sequence to the requested page count.

    Ranking every generated page independently made one flexible source slide
    win repeatedly. A template is more than a palette, so retain its actual
    cover/chapter/content/ending cadence and geometry families. Long decks may
    repeat pages, while short decks keep a bounded subset of chapter dividers.
    """
    slides = sorted(
        (item for item in library.get("slides", []) if isinstance(item, dict)),
        key=lambda item: int(item.get("slide_index") or 0),
    )
    target_count = max(0, int(target_count or 0))
    if not slides or target_count <= 0:
        return []
    if len(slides) == 1:
        return [slides[0]] * target_count
    if target_count == 1:
        return [slides[0]]

    source_last = len(slides) - 1
    target_last = target_count - 1
    positions = [round(index * source_last / target_last) for index in range(target_count)]

    chapter_positions = [
        index for index, slide in enumerate(slides)
        if _native_page_role(slide) == "chapter"
    ]
    # A 10-12 page management briefing should not spend a third of its pages
    # on dividers, while a long report should preserve every source section.
    chapter_limit = min(len(chapter_positions), max(1, math.ceil(target_count / 8)))
    if len(chapter_positions) > chapter_limit:
        if chapter_limit == 1:
            chapter_positions = [chapter_positions[0]]
        else:
            chapter_positions = [
                chapter_positions[round(index * (len(chapter_positions) - 1) / (chapter_limit - 1))]
                for index in range(chapter_limit)
            ]
    allowed_chapters = set(chapter_positions)

    content_positions = [
        index for index, slide in enumerate(slides)
        if _native_page_role(slide) in {"content", "data"}
    ]
    for target_index, source_index in enumerate(positions):
        if _native_page_role(slides[source_index]) != "chapter" or source_index in allowed_chapters:
            continue
        positions[target_index] = min(content_positions, key=lambda value: abs(value - source_index))

    for chapter_index in chapter_positions:
        target_index = round(chapter_index * target_last / source_last)
        target_index = min(max(target_index, 1), max(1, target_last - 1))
        positions[target_index] = chapter_index

    cover_index = next((index for index, slide in enumerate(slides) if _native_page_role(slide) == "cover"), 0)
    ending_index = next((index for index, slide in reversed(list(enumerate(slides))) if _native_page_role(slide) == "ending"), source_last)
    positions[0] = cover_index
    positions[-1] = ending_index
    agenda_index = next((index for index, slide in enumerate(slides) if _native_page_role(slide) == "agenda"), None)
    if agenda_index is not None and target_count >= 7:
        positions[1] = agenda_index
    for target_index in range(1, target_last):
        if _native_page_role(slides[positions[target_index]]) in {"cover", "ending"}:
            positions[target_index] = min(
                content_positions,
                key=lambda value: abs(value - positions[target_index]),
            )
    return [slides[index] for index in positions]


def _nearest_plain_native_slide(
    library: dict[str, Any],
    source_slide: int,
    used_counts: Counter[int] | None = None,
    preferred_family: str | None = None,
    exclude_slide: int | None = None,
    exclude_media: bool = False,
) -> dict[str, Any] | None:
    """Find a nearby non-data content page when a template table has no data."""
    counts = used_counts or Counter()
    candidates = [
        item for item in library.get("slides", [])
        if isinstance(item, dict)
        and _native_page_role(item) == "content"
        and not item.get("diagrams")
        and (not exclude_media or not _native_large_picture_count(item))
        and int(item.get("slide_index") or 0) != exclude_slide
    ]
    return min(
        candidates,
        key=lambda item: (
            0 if not preferred_family or _native_layout_family(item) == preferred_family else 200,
            counts[int(item.get("slide_index") or 0)] * 40,
            abs(int(item.get("slide_index") or 0) - source_slide),
        ),
        default=None,
    )


def _native_chart_host_usable(slide_entry: dict[str, Any], canvas_px: dict[str, Any] | None = None) -> bool:
    """Return whether a native shell has enough body area for an overlay chart."""
    body_slots = [
        slot for slot in (slide_entry.get("slots") or [])
        if isinstance(slot, dict) and _native_slot_is_body_like(slot)
    ]
    if not body_slots:
        return False
    _left, _top, width, height = _native_chart_bounds(slide_entry, canvas_px)
    # A chart needs room for marks, labels, and a short takeaway. Smaller
    # shells are still useful for text, but forcing a chart into them produces
    # the tiny donut/overlapping-controls failure seen in compact split pages.
    return width >= 6.0 and height >= 2.2


def _nearest_chart_native_slide(
    library: dict[str, Any],
    source_slide: int,
    used_counts: Counter[int] | None = None,
    used_families: Counter[str] | None = None,
) -> dict[str, Any] | None:
    """Find a clean, spacious native shell for an editable shape chart."""
    counts = used_counts or Counter()
    families = used_families or Counter()
    canvas_px = library.get("canvas_px") if isinstance(library.get("canvas_px"), dict) else None
    candidates: list[dict[str, Any]] = []
    for item in library.get("slides", []):
        if not isinstance(item, dict) or _native_page_role(item) != "content":
            continue
        if item.get("tables") or item.get("charts") or item.get("diagrams"):
            continue
        if _native_large_picture_count(item) or not _native_chart_host_usable(item, canvas_px):
            continue
        candidates.append(item)
    return min(
        candidates,
        key=lambda item: (
            counts[int(item.get("slide_index") or 0)] * 1000,
            families[_native_layout_family(item)] * 180,
            0 if _native_layout_family(item) == "hero-body" else 40,
            abs(int(item.get("slide_index") or 0) - source_slide),
        ),
        default=None,
    )


def _native_supported_charts(item: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        chart
        for chart in (item.get("charts") or [])
        if isinstance(chart, dict)
        and isinstance(chart.get("edit_capability"), dict)
        and chart.get("edit_capability", {}).get("supported") is True
    ]


def _native_table_shape(table: dict[str, Any]) -> tuple[int, int, int, bool]:
    source_rows = table.get("rows") if isinstance(table.get("rows"), list) else []
    first_row = source_rows[0] if source_rows and isinstance(source_rows[0], dict) else {}
    first_cells = [
        cell
        for cell in (first_row.get("cells") or [])
        if isinstance(cell, dict) and not cell.get("is_merge_slave")
    ]
    merged_title = (
        len([cell for cell in first_cells if str(cell.get("text") or "").strip()]) <= 1
        and any(int(cell.get("col_span") or 1) > 1 for cell in first_cells)
    )
    row_count = int(table.get("row_count") or 0)
    data_start = 2 if merged_title and row_count > 2 else 1
    return row_count, int(table.get("column_count") or 0), data_start, merged_title


def _native_best_fit_table(
    tables: list[dict[str, Any]],
    *,
    data_rows: int,
    data_cols: int,
) -> dict[str, Any] | None:
    if not tables:
        return None

    def mismatch(table: dict[str, Any]) -> tuple[int, int, int]:
        row_count, col_count, data_start, _merged_title = _native_table_shape(table)
        row_capacity = max(0, row_count - data_start)
        shortage = max(0, data_rows - row_capacity) + max(0, data_cols - col_count)
        distance = abs(row_capacity - data_rows) + abs(col_count - data_cols)
        return shortage, distance, row_count * col_count

    return min(tables, key=mismatch)


def _native_pick_slide(
    library: dict[str, Any],
    slide_spec: dict[str, Any],
    used_counts: Counter[int],
    last_slide: int | None,
    used_families: Counter[str] | None = None,
    last_family: str | None = None,
    sequence_hint: int | None = None,
) -> dict[str, Any] | None:
    """Pick a source page while keeping native layouts visually diverse.

    Most real-world decks label every ordinary page ``content_candidate``.
    Ranking only by slot count therefore tends to clone one attractive page
    over and over.  The filler still preserves the source slide exactly, but
    uses a lightweight geometry fingerprint and semantic preferences to spread
    output across the source deck's actual page families.
    """
    kind = str(slide_spec.get("kind") or "content").casefold()
    layout = str(slide_spec.get("layout_id") or "").casefold()
    visual_intent = str(slide_spec.get("visual_intent") or "").casefold()
    chart_payload = slide_spec.get("chart") if isinstance(slide_spec.get("chart"), dict) else {}
    chart_labels = [str(value).strip() for value in (chart_payload.get("labels") or []) if str(value).strip()]
    chart_values = _chart_values(chart_payload.get("values"))
    chart_data_ready = len(chart_labels) >= 2 and len(chart_values) >= 2
    chart_intent = kind == "chart" or "chart" in layout
    # SmartArt cannot be rewritten by the native filler. Native charts are safe
    # only when every chart on the page is editable and verified replacement
    # data is available; otherwise source semantics would survive unchanged.
    slides: list[dict[str, Any]] = []
    for item in library.get("slides", []):
        if not isinstance(item, dict) or item.get("diagrams"):
            continue
        charts = [chart for chart in (item.get("charts") or []) if isinstance(chart, dict)]
        if charts and (
            not chart_intent
            or not chart_data_ready
            or len(_native_supported_charts(item)) != len(charts)
        ):
            continue
        slides.append(item)
    if not slides:
        return None
    family_counts = used_families or Counter()
    density = str(slide_spec.get("density") or "medium").casefold()
    allow_template_media = bool(slide_spec.get("_allow_template_media"))
    requested_raw = str(slide_spec.get("template_slide") or "").strip()
    requested_raw = re.sub(r"^[Pp]\s*", "", requested_raw)
    try:
        requested_slide = int(requested_raw or 0)
    except (TypeError, ValueError):
        requested_slide = 0

    preferred_families: tuple[str, ...]
    if kind == "cover" or layout == "cover":
        preferred_families = ("cover",)
    elif kind in {"section", "chapter", "section-divider"} or layout == "section-divider" or visual_intent == "chapter":
        preferred_families = ("chapter",)
    elif kind == "closing" or layout == "closing" or visual_intent == "closing":
        preferred_families = ("ending",)
    elif kind in {"agenda", "toc"} or layout in {"agenda", "toc"}:
        preferred_families = ("agenda",)
    elif kind == "table" or layout in {"table", "data-table"} or visual_intent == "table":
        preferred_families = ("table-",)
    elif kind == "chart" or "chart" in layout:
        # A table-shaped page is not a chart. Prefer a native chart/data page
        # when the template actually contains one; otherwise choose an
        # ordinary visual shell and let the caller degrade the chart to a
        # readable evidence block instead of exposing an empty example table.
        has_native_chart = any(_native_supported_charts(item) for item in slides)
        preferred_families = (
            ("data", "table-", "hero-body", "single")
            if has_native_chart
            else ("hero-body", "single", "split", "tri-column", "stack", "grid-")
        )
    elif kind in {"comparison", "split"} or layout in {"comparison", "split"} or visual_intent == "comparison":
        preferred_families = ("split", "tri-column", "grid-")
    elif kind == "cards" or layout in {"cards", "card-grid", "six-card"}:
        preferred_families = ("grid-", "tri-column", "split")
    elif kind == "timeline" or layout in {"timeline", "roadmap"} or visual_intent == "process":
        preferred_families = ("stack", "split", "hero-body")
    elif kind in {"image", "photo"} or "photo" in layout or visual_intent in {"photo", "evidence", "site", "image"}:
        preferred_families = ("photo-body", "hero-body", "single", "split")
    elif kind in {"metric", "statement", "quote"} or visual_intent == "metrics":
        preferred_families = ("single", "hero-body", "split", "tri-column", "grid-")
    elif visual_intent == "hero":
        preferred_families = ("hero-body", "single", "photo-body", "split")
    else:
        preferred_families = ("single", "split", "tri-column", "stack", "hero-body", "grid-")

    family = _native_layout_family
    table_payload = slide_spec.get("table") if isinstance(slide_spec.get("table"), dict) else {}
    table_rows = [row for row in (table_payload.get("rows") or []) if isinstance(row, list) and row]
    table_data_ready = bool(table_rows)
    table_columns = [value for value in (table_payload.get("columns") or []) if str(value).strip()]
    target_table_cols = max([len(table_columns), *(len(row) for row in table_rows)], default=0)
    target_table_rows = len(table_rows)
    requested_units = 0
    if kind == "cards" or layout in {"cards", "card-grid", "six-card"}:
        requested_units = len([
            value for value in (slide_spec.get("bullets") or []) if str(value).strip()
        ])
    elif kind in {"comparison", "split"} or layout in {"comparison", "split"} or visual_intent == "comparison":
        requested_units = sum(
            len([value for value in (slide_spec.get(key) or []) if str(value).strip()])
            for key in ("left_bullets", "right_bullets")
        )
    elif kind == "timeline" or layout in {"timeline", "roadmap"} or visual_intent == "process":
        requested_units = len([
            value
            for value in (slide_spec.get("steps") or [])
            if isinstance(value, dict)
            and (str(value.get("title") or "").strip() or str(value.get("body") or "").strip())
        ])

    def rank(item: dict[str, Any]) -> float:
        page_type = str(item.get("page_type") or "").casefold()
        has_table = bool(item.get("tables"))
        has_chart = bool(item.get("charts"))
        item_family = family(item)
        item_role = _native_page_role(item)
        item_index = int(item.get("slide_index") or 0)
        item_capacity = _native_layout_unit_capacity(item)
        summary = re.sub(r"\s+", "", str(item.get("text_summary") or "").casefold())
        # A source slide can be reused for very long decks, but it should not
        # win again merely because it has many editable slots.
        score = -used_counts[item_index] * 140.0
        if item_index == last_slide:
            score -= 1000
        if last_family and item_family == last_family:
            score -= 260
        # Reuse is allowed for long (up to 100 page) decks, but a fresh family
        # wins until every family has appeared once. This is the main guard
        # against the "same controls on every page" failure mode.
        score -= family_counts[item_family] * 170
        family_rank = next(
            (
                index
                for index, preferred in enumerate(preferred_families)
                if item_family == preferred or (preferred.endswith("-") and item_family.startswith(preferred))
            ),
            None,
        )
        if family_rank is not None:
            score += max(170, 600 - family_rank * 125)
        else:
            score -= 25
        if requested_slide and item_index == requested_slide:
            # An explicit page hint is advisory: it must still match the
            # requested visual intent, and repeated hints rotate once the
            # same shell has already been used nearby.
            primary_hint_compatible = family_rank == 0
            if visual_intent in {"photo", "evidence", "site", "image"} or kind in {"image", "photo"}:
                primary_hint_compatible = _native_large_picture_count(item) > 0
            elif visual_intent == "table" or kind == "table":
                primary_hint_compatible = item_family.startswith("table-")
            elif visual_intent == "chapter" or kind in {"section", "chapter", "section-divider"}:
                primary_hint_compatible = item_family == "chapter"
            elif visual_intent == "closing" or kind == "closing":
                primary_hint_compatible = item_family == "ending"
            if primary_hint_compatible:
                score += 2350 if used_counts[item_index] == 0 and item_index != last_slide else -1700
            else:
                score += 40
        if sequence_hint and item_index == sequence_hint:
            score += 260
        elif sequence_hint:
            hinted = next(
                (slide for slide in slides if int(slide.get("slide_index") or 0) == sequence_hint),
                None,
            )
            if hinted is not None and family(hinted) == item_family:
                score += 90

        expected_roles = {
            "cover": {"cover"},
            "closing": {"ending"},
            "section": {"chapter"},
            "chapter": {"chapter"},
            "section-divider": {"chapter"},
            "agenda": {"agenda"},
            "toc": {"agenda"},
            "table": {"data"},
        }
        role_kind = visual_intent if kind in {"content", "statement"} and visual_intent in expected_roles else kind
        expected = expected_roles.get(role_kind)
        if expected is not None and item_role not in expected:
            score -= 5000
        elif expected is None and item_role in {"cover", "agenda", "chapter", "ending"}:
            score -= 5000
        if role_kind not in {"table", "chart"} and item_role == "data":
            score -= 3200
        template_image_allowed = (
            kind in {"cover", "section", "chapter", "section-divider", "closing", "image", "photo"}
            or visual_intent in {"photo", "evidence", "site", "image", "chapter", "closing"}
        )
        if kind in {"image", "photo"} or visual_intent in {"photo", "evidence", "site", "image"}:
            template_image_allowed = allow_template_media
        if _native_large_picture_count(item) and not template_image_allowed:
            # Native template-fill cannot replace an arbitrary source photo.
            # Never expose an old template image as evidence for a new story
            # merely to gain visual variety.
            score -= 4200
        if item.get("diagrams"):
            # ppt-master intentionally preserves SmartArt but cannot rewrite
            # its nodes through a text replacement plan. Prefer ordinary
            # shape-based layouts unless the user explicitly requested a
            # diagram-like page.
            score -= 1400 if kind not in {"timeline", "comparison", "process", "diagram"} else 40
        if has_chart:
            # A native chart carries its source data even when its text slots
            # are replaced. Only select it when the generated page has a
            # verified chart payload; otherwise an example chart would leak
            # stale numbers into the new deck.
            if kind == "chart" or "chart" in layout:
                score += 450 if chart_data_ready else -1200
            else:
                score -= 1600
        if kind == "chart" and not has_chart and has_table:
            # The source table may contain example numbers unrelated to the
            # generated chart. Never select it merely because it has many
            # editable cells; that produces a misleading half-empty table.
            score -= 2200
        if kind in {"metric", "statement", "quote"} and item_family.startswith("grid-"):
            # Six-card research grids are visually distinctive but usually
            # carry dense source-specific labels. Using one for a single KPI
            # leaves half the cards empty and makes the generated page look
            # like a damaged template; reserve grids for explicit card pages.
            score -= 5200
        if kind in {"metric", "statement", "quote"} and item_family in {"tri-column", "stack"}:
            score -= 1800
        if kind in {"metric", "statement", "quote"} and item_family == "split":
            score -= 2200
        if requested_units:
            # Match semantic item count rather than only x-axis columns. This
            # distinguishes a six-card grid from a three-card row and avoids
            # conspicuous empty source controls.
            overflow = max(0, item_capacity - requested_units)
            shortage = max(0, requested_units - item_capacity)
            score += 900 - overflow * 380 - shortage * 560
        if kind == "table" and has_table and target_table_rows and target_table_cols:
            candidate = _native_best_fit_table(
                [table for table in item.get("tables") or [] if isinstance(table, dict)],
                data_rows=target_table_rows,
                data_cols=target_table_cols,
            )
            if candidate:
                rows, cols, data_start, _merged_title = _native_table_shape(candidate)
                row_capacity = max(0, rows - data_start)
                if row_capacity >= target_table_rows and cols >= target_table_cols:
                    score += 620 - min((row_capacity - target_table_rows) * 12 + (cols - target_table_cols) * 18, 430)
                else:
                    score -= 2600
        if has_table and kind in {"table", "chart"} and not table_data_ready and not chart_data_ready:
            # Table shells are still useful for an explicitly requested table,
            # but a data-free generic page should not win over ordinary text
            # layouts when the source table contains example values.
            score -= 500
        # Contents/agenda pages are useful only when the generated slide is
        # explicitly a section map. Reusing one as ordinary body content is a
        # common source of empty labels and a visibly wrong hierarchy.
        if "目录" in summary and kind not in {"agenda", "toc"}:
            score -= 5000
        if kind == "cover":
            score += 1000 if "cover" in page_type else 0
        elif kind == "closing":
            score += 1000 if any(token in page_type for token in ("ending", "closing")) else 0
        elif kind in {"section", "chapter", "section-divider"}:
            score += 900 if "chapter" in page_type else 0
        elif kind in {"chart", "table"} or "chart" in layout:
            score += 450 if has_chart and chart_data_ready else 0
            score += 220 if has_table else 0
            score += 100 if "content" in page_type else 0
        elif kind in {"image", "photo"} or "photo" in layout or visual_intent in {"photo", "evidence", "site", "image"}:
            picture_count = _native_large_picture_count(item)
            score += 760 if picture_count else -420
            score += 120 if (has_table is False and len(item.get("slots") or []) < 14) else 0
        else:
            score += 180 if "content" in page_type else 0
            score += 80 if not has_table else 0
        slot_count = len(item.get("slots") or [])
        if density == "high":
            score += min(slot_count, 20) * 5
        elif density == "low":
            score += max(0, 14 - slot_count) * 4
        # Prefer pages with enough editable slots for substantive content, but
        # cap this contribution so a dense page cannot dominate every choice.
        score += min(slot_count, 18) * 2
        # Stable tie-breaker: move through the source deck instead of always
        # taking the first equally ranked page.
        score += (int(item.get("slide_index") or 0) % 17) * 0.01
        return score

    return max(slides, key=rank)


def _native_table_edits(slide_entry: dict[str, Any], slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    """Edit one best-fit native table and clear every other source table."""
    tables = [table for table in (slide_entry.get("tables") or []) if isinstance(table, dict)]
    data = slide_spec.get("table") if isinstance(slide_spec.get("table"), dict) else None
    chart = slide_spec.get("chart") if isinstance(slide_spec.get("chart"), dict) else None
    if not tables:
        return []
    rows = data.get("rows") if isinstance(data, dict) and isinstance(data.get("rows"), list) else []
    columns = data.get("columns") if isinstance(data, dict) and isinstance(data.get("columns"), list) else []
    labels = chart.get("labels") if isinstance(chart, dict) else []
    values = chart.get("values") if isinstance(chart, dict) else []
    if not rows and labels and values:
        rows = [[label, value] for label, value in zip(labels, values)]
    has_payload = bool(rows or columns)

    target: dict[str, Any] | None = None
    if has_payload:
        desired_cols = max(
            [len(columns), *(len(row) for row in rows if isinstance(row, list))],
            default=0,
        )
        target = _native_best_fit_table(tables, data_rows=len(rows), data_cols=desired_cols)

    edits: list[dict[str, Any]] = []
    for table in tables:
        table_id = str(table.get("table_id") or "")
        if not table_id:
            continue
        source_rows = table.get("rows") if isinstance(table.get("rows"), list) else []
        _row_count, max_cols, data_start, merged_title = _native_table_shape(table)
        cells: list[dict[str, Any]] = []
        # Blank every visible cell first, including merge anchors. Secondary
        # tables and unused rows can therefore never expose source figures.
        for source_row in source_rows:
            if not isinstance(source_row, dict):
                continue
            for source_cell in source_row.get("cells") or []:
                if not isinstance(source_cell, dict) or source_cell.get("is_merge_slave"):
                    continue
                row_index = int(source_cell.get("row", -1))
                col_index = int(source_cell.get("col", -1))
                if row_index < 0 or col_index < 0:
                    continue
                cells.append({"row": row_index, "col": col_index, "text": ""})

        if table is target:
            value_by_cell: dict[tuple[int, int], str] = {}
            header_index = data_start - 1
            if merged_title:
                value_by_cell[(0, 0)] = _clip(str(slide_spec.get("title") or ""), 100)
            for col_index in range(max_cols):
                value_by_cell[(header_index, col_index)] = _clip(
                    str(columns[col_index]) if col_index < len(columns) else "",
                    80,
                )
            for offset, values_row in enumerate(rows):
                if not isinstance(values_row, list):
                    continue
                row_index = data_start + offset
                for col_index in range(max_cols):
                    value_by_cell[(row_index, col_index)] = _clip(
                        str(values_row[col_index]) if col_index < len(values_row) else "",
                        100,
                    )
            for cell in cells:
                key = (int(cell["row"]), int(cell["col"]))
                if key in value_by_cell:
                    cell["text"] = value_by_cell[key]
        if cells:
            edits.append({"table_id": table_id, "cells": cells})
    return edits


def _native_chart_edits(slide_entry: dict[str, Any], slide_spec: dict[str, Any]) -> list[dict[str, Any]]:
    """Feed verified model data into an editable classic native chart."""
    charts = [chart for chart in (slide_entry.get("charts") or []) if isinstance(chart, dict)]
    chart = slide_spec.get("chart") if isinstance(slide_spec.get("chart"), dict) else None
    if not charts or not chart:
        return []
    labels = [str(value) for value in (chart.get("labels") or []) if str(value).strip()]
    values = _chart_values(chart.get("values"))
    count = min(len(labels), len(values), 12)
    if count < 2:
        return []
    targets = _native_supported_charts(slide_entry)
    if len(targets) != len(charts):
        return []
    edits: list[dict[str, Any]] = []
    for target in targets:
        chart_id = str(target.get("chart_id") or "")
        if not chart_id:
            return []
        edits.append({
            "chart_id": chart_id,
            "categories": labels[:count],
            "series": [{"name": str(chart.get("series_name") or "系列1"), "values": values[:count]}],
        })
    return edits


def _pick_reference_image(slide_spec: dict[str, Any], image_assets: list[PresentationJobAsset], used: Counter[int]) -> Path | None:
    """Choose the most relevant uploaded image instead of reusing image 1."""
    if not image_assets:
        return None
    haystack = " ".join(
        str(slide_spec.get(key) or "")
        for key in ("title", "summary", "body", "image_slot", "kicker")
    ).casefold()
    tokens = {token for token in re.findall(r"[\w\u4e00-\u9fff]{2,}", haystack) if token not in {"参考", "图片", "视觉", "证据"}}
    ranked: list[tuple[float, PresentationJobAsset]] = []
    for asset in image_assets:
        filename = str(asset.filename or "").casefold()
        overlap = sum(1 for token in tokens if token in filename)
        score = overlap * 80.0 - used.get(int(asset.id), 0) * 8.0
        ranked.append((score, asset))
    ranked.sort(key=lambda item: (item[0], -int(item[1].id)), reverse=True)
    selected = ranked[0][1]
    used[int(selected.id)] += 1
    path = Path(selected.file_path)
    return path if path.exists() else None


def _presentation_vision_attachments(image_assets: list[PresentationJobAsset]) -> list[dict[str, Any]]:
    """Build a bounded multimodal context for the story-planning request.

    Images remain stored on disk and are still used as editable visual
    evidence later. Only a small, bounded subset is sent to a vision-capable
    text model so a screenshot or field photo can influence the outline rather
    than being treated as an opaque filename. ``OpenAIService`` gracefully
    falls back to text-only completions when the configured gateway has no
    vision support.
    """
    attachments: list[dict[str, Any]] = []
    total_bytes = 0
    for asset in image_assets:
        if len(attachments) >= MAX_VISION_IMAGES or total_bytes >= MAX_VISION_TOTAL_BYTES:
            break
        path = Path(asset.file_path)
        if not path.exists() or not path.is_file():
            continue
        try:
            file_size = path.stat().st_size
        except OSError:
            continue
        if file_size <= 0 or file_size > MAX_VISION_IMAGE_BYTES or total_bytes + file_size > MAX_VISION_TOTAL_BYTES:
            continue
        content_type = str(asset.content_type or "image/png").split(";", 1)[0].strip().lower()
        if not content_type.startswith("image/"):
            continue
        try:
            encoded = base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            continue
        attachments.append(
            {
                "filename": str(asset.filename or path.name),
                "content_type": content_type,
                "file_size": file_size,
                "data_url": f"data:{content_type};base64,{encoded}",
            }
        )
        total_bytes += file_size
    return attachments


def _build_native_fill_plan(
    library: dict[str, Any],
    slides: list[dict[str, Any]],
    *,
    allow_template_media: bool = False,
) -> dict[str, Any]:
    """Build a confirmed internal plan for the project-local native filler."""
    static_texts = _native_static_texts(library)
    slide_count = len(library.get("slides") or [])
    used_counts: Counter[int] = Counter()
    used_families: Counter[str] = Counter()
    plan_slides: list[dict[str, Any]] = []
    last_slide: int | None = None
    last_family: str | None = None
    chapter_number = 0
    section_item_number = 0
    sequence = _native_sequence_sources(library, len(slides))
    for index, slide_spec in enumerate(slides, start=1):
        effective_spec = dict(slide_spec)
        effective_spec["_allow_template_media"] = allow_template_media
        source = sequence[index - 1] if index <= len(sequence) else None
        kind = str(effective_spec.get("kind") or "content").casefold()
        layout = str(effective_spec.get("layout_id") or "").casefold()
        chart_overlay: dict[str, Any] | None = None
        has_chart_data = isinstance(effective_spec.get("chart"), dict) and len(_chart_values(effective_spec["chart"].get("values"))) >= 2
        table_payload = effective_spec.get("table") if isinstance(effective_spec.get("table"), dict) else {}
        has_table_data = bool([row for row in (table_payload.get("rows") or []) if isinstance(row, list) and row])

        if index == 1:
            effective_spec["kind"], effective_spec["layout_id"] = "cover", "cover"
            effective_spec["subtitle"] = str(
                effective_spec.get("subtitle")
                or effective_spec.get("summary")
                or effective_spec.get("body")
                or ""
            )
            effective_spec["body"] = ""
            effective_spec["summary"] = ""
            effective_spec["bullets"] = []
        elif index == len(slides):
            effective_spec["kind"], effective_spec["layout_id"] = "closing", "closing"
            effective_spec["subtitle"] = str(
                effective_spec.get("subtitle")
                or effective_spec.get("summary")
                or effective_spec.get("body")
                or ""
            )
            effective_spec["body"] = ""
            effective_spec["summary"] = ""
            effective_spec["bullets"] = []
        kind = str(effective_spec.get("kind") or "content").casefold()
        layout = str(effective_spec.get("layout_id") or "").casefold()
        # The source sequence provides formal report rhythm. Semantic ranking
        # then selects the best page inside that rhythm, with an explicit
        # model-selected template page taking priority only when it is safe.
        semantic_source = _native_pick_slide(
            library,
            effective_spec,
            used_counts,
            last_slide,
            used_families,
            last_family,
            int(source.get("slide_index") or 0) if source else None,
        )
        source = semantic_source
        if not source:
            raise RuntimeError(
                f"模板中没有可安全填充第 {index} 页（{kind or 'content'}）的原生版式；"
                "请更换模板或移除无法编辑的 SmartArt/图表页。"
            )
        # A native template may have no editable chart object. Move to a plain
        # shell and retain the verified chart as an editable shape overlay;
        # this keeps the uploaded background/typography while still giving
        # the user a real visual data explanation instead of a stale example
        # table or a paragraph of numbers.
        if source and (kind == "chart" or "chart" in layout) and not source.get("charts"):
            source_index_hint = int(source.get("slide_index") or 0)
            if (
                source.get("tables")
                or _native_large_picture_count(source)
                or not _native_chart_host_usable(source, library.get("canvas_px"))
            ):
                source = _nearest_chart_native_slide(
                    library,
                    source_index_hint,
                    used_counts,
                    used_families,
                ) or _nearest_plain_native_slide(
                    library,
                    source_index_hint,
                    used_counts,
                    preferred_family="hero-body",
                    exclude_media=True,
                ) or source
            chart_payload = effective_spec.get("chart") if isinstance(effective_spec.get("chart"), dict) else {}
            labels = [str(value).strip() for value in (chart_payload.get("labels") or []) if str(value).strip()]
            values = _chart_values(chart_payload.get("values"))
            unit = str(chart_payload.get("unit") or "").strip()
            if labels and values:
                chart_overlay = {
                    "type": str(chart_payload.get("type") or "bar"),
                    "labels": labels[:8],
                    "values": values[:8],
                    "unit": unit,
                    "series_name": str(chart_payload.get("series_name") or "系列1"),
                    "metric": str(effective_spec.get("metric") or "构成"),
                    "summary": str(effective_spec.get("summary") or "图表只呈现资料中可验证的变化。"),
                    "bounds": _native_chart_bounds(source, library.get("canvas_px")),
                }
                effective_spec["body"] = ""
                effective_spec["summary"] = ""
                effective_spec["bullets"] = []
        source_index = int(source.get("slide_index") or 0)
        source_family = _native_layout_family(source)
        if source_family == "chapter":
            chapter_number += 1
            section_item_number = 0
            kicker = str(effective_spec.get("kicker") or "").strip()
            # Model output can omit a chapter counter or echo the source
            # template's English “Part Three”.  Number the cloned chapter in
            # output order while preserving an explicit user-authored label.
            if not kicker or re.fullmatch(r"(?:part|section|chapter)\s+\w+", kicker, flags=re.I):
                effective_spec["kicker"] = f"第{chapter_number:02d}章"
        elif source_family not in {"cover", "agenda", "ending"}:
            section_item_number += 1
            effective_spec["_native_page_marker"] = f"{min(section_item_number, 99):02d}"
        used_counts[source_index] += 1
        used_families[source_family] += 1
        last_slide = source_index
        last_family = source_family
        replacements = _native_replacements(source, effective_spec, static_texts, slide_count)
        plan_item: dict[str, Any] = {
            "source_slide": source_index,
            "purpose": str(effective_spec.get("kind") or "content"),
            "layout_rationale": {
                "layout_pattern": str(source.get("page_type") or "content_candidate"),
                "layout_family": source_family,
                "why_fit": f"{_native_layout_affordance(source)}；按语义版式选择并保留母版、主题和媒体关系。",
                "risk": "文本会受源页面槽位容量约束，超长内容已压缩。",
            },
            "replacements": replacements,
            "table_edits": _native_table_edits(source, effective_spec),
            "chart_edits": _native_chart_edits(source, effective_spec),
        }
        hidden_shapes = _native_hidden_empty_shape_ids(source, replacements)
        if hidden_shapes:
            plan_item["hide_shapes"] = hidden_shapes
        if chart_overlay:
            plan_item["shape_chart"] = chart_overlay
            plan_item["source_analysis"] = source
        plan_slides.append(plan_item)
    return {"schema": "template_fill_pptx_plan.v1", "status": "confirmed", "slides": plan_slides}


def _apply_native_fill(template_path: Path, plan: dict[str, Any], output_path: Path) -> None:
    """Clone and patch native slides through the bundled ppt-master runtime."""
    script_root = str(PPT_MASTER_SCRIPTS)
    if script_root not in sys.path:
        sys.path.insert(0, script_root)
    from template_fill_pptx.applier import apply_plan

    apply_plan(template_path, plan, output_path, transition="keep")


def _font_for(palette: dict[str, Any], requested: str) -> str:
    """Use the template's first detected font for the normal UI font family."""
    if requested in {"Aptos", "Aptos Mono"} and palette.get("font"):
        return str(palette["font"])
    return requested


def _fallback_reference_lines(references: str) -> list[str]:
    """Extract conservative, human-readable facts for model-free fallback.

    The fallback is used when a provider is unavailable or returns malformed
    JSON. It should still reflect the uploaded material instead of filling a
    deck with generic lorem ipsum.
    """
    facts: list[str] = []
    seen: set[str] = set()
    for raw_line in str(references or "").splitlines():
        line = " ".join(str(raw_line).replace("\u3000", " ").split()).strip()
        line = re.sub(r"^\[(?:slide|page|sheet)[^\]]*\]\s*", "", line, flags=re.I)
        line = re.sub(r"^【[^】]+】\s*", "", line)
        if not line or len(line) < 6 or re.fullmatch(r"[-\d\s./]+", line):
            continue
        if re.fullmatch(r"(?:指标|年份|年度|日期|时间|year)\s*[|,:：]\s*(?:19|20)\d{2}", line, flags=re.I):
            continue
        normalized = line.casefold()
        if normalized in seen:
            continue
        seen.add(normalized)
        facts.append(_clip(line, 180))
        if len(facts) >= 80:
            break
    return facts


def _fallback_spec(job: PresentationJob, references: str) -> dict[str, Any]:
    title = job.title.strip() or "专题演示"
    audience = job.audience.strip() or "团队与决策者"
    purpose = job.purpose.strip() or "对齐认知并推动下一步行动"
    count = max(5, min(100, int(job.slide_count or 10)))
    topic = _clip(job.brief, 90) or title
    facts = _fallback_reference_lines(references)
    slides: list[dict[str, Any]] = [
        {"kind": "cover", "layout_id": "cover", "title": title, "subtitle": f"面向{audience} · {purpose}", "kicker": "专项汇报 / SPECIAL BRIEFING"},
    ]

    if facts:
        first = facts[0]
        slides.append({
            "kind": "statement", "layout_id": "statement", "kicker": "EXECUTIVE TAKEAWAY",
            "title": _clip(f"资料首先指向：{first}", 72),
            "summary": "以下页面只保留上传资料中出现的事实，未补造缺失数据。",
            "body": first,
            "bullets": facts[1:4],
        })
        cursor = 1
        page_index = 1
        while cursor < len(facts) and len(slides) < count - 1:
            group = facts[cursor : cursor + 4]
            if not group:
                break
            if len(group) >= 3 and page_index % 3 == 0:
                slides.append({
                    "kind": "cards", "layout_id": "cards", "kicker": f"EVIDENCE / {page_index:02d}",
                    "title": _clip(group[0], 72),
                    "summary": "把同一主题下的资料要点拆成可快速扫描的判断单元。",
                    "bullets": group[:4],
                })
            elif len(group) >= 3 and page_index % 3 == 1:
                slides.append({
                    "kind": "comparison", "layout_id": "comparison", "kicker": f"DECISION / {page_index:02d}",
                    "title": _clip(group[0], 72),
                    "summary": "左右两栏保留原始事实，方便在汇报现场直接比较。",
                    "left_title": "资料现状", "left_bullets": group[:2],
                    "right_title": "需要关注", "right_bullets": group[2:4] or ["请结合业务口径进一步核验"],
                })
            else:
                slides.append({
                    "kind": "content", "layout_id": "content", "kicker": f"INSIGHT / {page_index:02d}",
                    "title": _clip(group[0], 72),
                    "summary": "围绕上传资料提炼的一条可核验信息。",
                    "body": group[0], "bullets": group[1:4],
                })
            cursor += len(group)
            page_index += 1
    else:
        slides.extend([
            {"kind": "statement", "layout_id": "statement", "kicker": "EXECUTIVE TAKEAWAY", "title": "先把关键判断说清楚", "summary": "复杂资料需要被压缩成一个可以继续讨论和执行的判断。", "body": f"围绕“{topic}”建立一条结论先行的汇报主线。"},
            {"kind": "metric", "layout_id": "metric", "kicker": "SIGNAL / 01", "title": "最值得关注的是影响最大的变量", "summary": "用少量可验证指标判断优先级，而不是堆积所有数据。", "metric": "01", "metric_label": "核心信号", "bullets": ["明确当前状态", "找到变化来源", "标记下一步动作"]},
            {"kind": "comparison", "layout_id": "comparison", "kicker": "DECISION FRAME", "title": "现状与目标之间存在一条可缩短的路径", "summary": "把选项放入同一判断框架，方便管理层快速取舍。", "left_title": "当前状态", "left_bullets": ["事实和约束可见", "问题集中在关键环节"], "right_title": "目标状态", "right_bullets": ["优先级统一", "动作可以被追踪"]},
        ])

    inferred_chart = _infer_chart_from_reference(references)
    if inferred_chart and len(slides) < count - 1:
        slides.append({
            "kind": "chart", "layout_id": f"{inferred_chart.get('type') or 'bar'}-chart",
            "kicker": "DATA / EXTRACTED", "title": "资料中的量化变化值得单独看",
            "summary": "图表直接取自参考资料中的标签和数字，未补造缺失数据。", "chart": inferred_chart,
        })

    # Keep a deterministic minimum narrative spine when the source is short.
    titles = ["关键洞察", "方案框架", "风险与边界", "衡量成功", "下一步行动"]
    while len(slides) < count - 1:
        index = len(slides) - 1
        slides.append({
            "kind": "content", "layout_id": "content", "kicker": f"INSIGHT / {index:02d}", "title": titles[index % len(titles)],
            "summary": "围绕核心目标提炼一条可验证、可执行的结论。", "body": f"围绕“{topic}”把事实、判断和动作放在同一页中。",
            "bullets": ["用事实描述当前状态", "指出影响最大的变量", "给出本周可以启动的动作"],
            "_generated_filler": True,
        })
    slides.append({"kind": "closing", "layout_id": "closing", "kicker": "NEXT MOVE", "title": "把共识变成下一次可见的交付", "subtitle": "今天决定方向，下一次带回结果。"})
    return {"title": title, "slides": slides[:count]}


def _chart_values(raw: Any) -> list[float]:
    values: list[float] = []
    if not isinstance(raw, list):
        return values
    for value in raw[:8]:
        try:
            number = float(value)
        except (TypeError, ValueError):
            continue
        if math.isfinite(number):
            values.append(number)
    return values


def _infer_chart_from_reference(reference_text: str) -> dict[str, Any] | None:
    """Extract conservative label-number pairs from office/CSV-like text.

    ``document_extract`` renders spreadsheets as ``a | b`` rows, while plain
    notes commonly use ``label: number``. Accept both forms and choose a line
    chart for ordered time labels; never manufacture missing values.
    """
    candidates: list[tuple[int, str, str, float, str]] = []
    section = 0
    numeric_pattern = re.compile(r"^-?\d+(?:\.\d+)?\s*(%|％|万元|万|亿|人|户|件)?$")
    for raw_line in str(reference_text or "").splitlines():
        line = " ".join(raw_line.strip().split())
        if not line:
            section += 1
            continue
        if line.startswith("【") or re.match(r"^\[(?:slide|page|sheet|table)\b", line, flags=re.I):
            section += 1
            continue
        parts = [part.strip() for part in re.split(r"\s*(?:\||,|\t|;|：|:)\s*", line) if part.strip()]
        if len(parts) < 2:
            # A space-delimited note: take the final numeric token as value.
            match = re.match(r"^(.+?)\s+(-?\d+(?:\.\d+)?)\s*(%|％|万元|万|亿|人|户|件)?$", line)
            if not match:
                continue
            parts = [match.group(1).strip(), f"{match.group(2)}{match.group(3) or ''}"]
        numeric_index = next((index for index, part in enumerate(parts[1:], start=1) if numeric_pattern.match(part)), None)
        if numeric_index is None:
            continue
        value_token = parts[numeric_index]
        number_match = re.match(r"^(-?\d+(?:\.\d+)?)\s*(%|％|万元|万|亿|人|户|件)?$", value_token)
        if not number_match:
            continue
        label = _clip(" / ".join(parts[:numeric_index]), 24)
        try:
            value = float(number_match.group(1))
        except (TypeError, ValueError):
            continue
        if not math.isfinite(value):
            continue
        normalized_label = re.sub(r"\s+", "", label).casefold()
        # Do not chart a table heading such as "年份 | 2025" as a measure.
        if 1900 <= value <= 2100 and normalized_label in {"年份", "年度", "日期", "时间", "指标", "year"}:
            continue
        if normalized_label in {"序号", "编号", "index", "no.", "no"}:
            continue
        unit = (number_match.group(2) or "").replace("％", "%")
        candidates.append((section, unit, label, value, line))

    if len(candidates) < 2:
        return None

    # Prefer a coherent table/slide block with one unit. If source extraction
    # split a plain-text list into separate blocks, fall back to the largest
    # same-unit group across the document.
    grouped: dict[tuple[int, str], list[tuple[int, str, str, float, str]]] = {}
    for candidate in candidates:
        grouped.setdefault((candidate[0], candidate[1]), []).append(candidate)
    coherent = [values for values in grouped.values() if len(values) >= 2]
    if coherent:
        selected = max(coherent, key=lambda values: (len(values), bool(values[0][1])))
    else:
        by_unit: dict[str, list[tuple[int, str, str, float, str]]] = {}
        for candidate in candidates:
            by_unit.setdefault(candidate[1], []).append(candidate)
        eligible = [values for values in by_unit.values() if len(values) >= 2]
        if not eligible:
            return None
        selected = max(eligible, key=lambda values: (len(values), bool(values[0][1])))

    labels: list[str] = []
    values: list[float] = []
    source_lines: list[str] = []
    unit = selected[0][1]
    for _section, _unit, label, value, source_line in selected:
        if label in labels:
            continue
        labels.append(label)
        values.append(value)
        source_lines.append(source_line)
        if len(labels) >= 8:
            break
    if len(labels) < 2:
        return None

    ordered_time = all(bool(re.search(r"(?:19|20)\d{2}|年|季度|Q[1-4]|月", label, flags=re.I)) for label in labels)
    chart_type = "line" if ordered_time else "bar"
    selected_context = " ".join(source_lines)
    looks_like_composition = any(token in selected_context for token in ("占比", "构成", "比例")) or (unit == "%" and 95 <= sum(values) <= 105)
    if looks_like_composition and all(value >= 0 for value in values) and sum(values) > 0:
        chart_type = "donut"
    return {"type": chart_type, "labels": labels, "values": values, "unit": unit}


def _chart_is_supported_by_source(chart: Any, evidence_text: str) -> bool:
    """Return true only when every label/value pair is traceable to a source row."""
    if not isinstance(chart, dict):
        return False
    values = _chart_values(chart.get("values"))
    labels = [str(value).strip() for value in (chart.get("labels") or []) if str(value).strip()]
    if len(values) < 2 or len(labels) < len(values):
        return False
    evidence = str(evidence_text or "")
    fragments = [
        fragment for fragment in re.split(r"[\n。；;]+", evidence)
        if fragment.strip()
    ]

    def numbers(fragment: str) -> list[float]:
        result: list[float] = []
        for raw in re.findall(r"(?<![\d.])(-?\d[\d,]*(?:\.\d+)?)\s*(%|％)?", fragment):
            try:
                number = float(raw[0].replace(",", ""))
            except ValueError:
                continue
            result.append(number)
            if raw[1] and number != 0:
                result.append(number / 100.0)
        return result

    def supported(value: float, candidates: list[float]) -> bool:
        tolerance = max(1e-6, abs(value) * 1e-6)
        return any(abs(source - value) <= tolerance for source in candidates)

    for label, value in zip(labels, values):
        label_key = re.sub(r"[\W_]+", "", label, flags=re.UNICODE).casefold()
        matching = [
            fragment for fragment in fragments
            if label_key and label_key in re.sub(r"[\W_]+", "", fragment, flags=re.UNICODE).casefold()
        ]
        if not matching or not any(supported(value, numbers(fragment)) for fragment in matching):
            return False
    return True


def _table_is_supported_by_source(table: Any, evidence_text: str) -> bool:
    """Fail closed when a model adds table cells that do not occur in source."""
    if not isinstance(table, dict):
        return False
    rows = [row for row in (table.get("rows") or []) if isinstance(row, list) and row]
    if not rows:
        return False
    evidence_key = re.sub(r"[\W_]+", "", str(evidence_text or ""), flags=re.UNICODE).casefold()
    structural_headers = {
        "序号", "编号", "指标", "项目", "类别", "名称", "类型", "单位", "数量", "金额",
        "说明", "备注", "状态", "时间", "日期", "合计", "占比", "变化", "责任人", "完成情况",
        "index", "no", "item", "name", "type", "unit", "value", "amount", "note", "status",
    }

    for column in (table.get("columns") or []):
        header = str(column or "").strip()
        if not header:
            continue
        header_key = re.sub(r"[\W_]+", "", header, flags=re.UNICODE).casefold()
        if header_key in structural_headers:
            continue
        if not header_key or header_key not in evidence_key:
            return False
    evidence_numbers = _chart_values([
        token.replace(",", "")
        for token in re.findall(r"(?<![\d.])-?\d[\d,]*(?:\.\d+)?", str(evidence_text or ""))
    ])
    checked = 0
    for row in rows:
        for cell in row:
            value = str(cell or "").strip()
            if not value:
                continue
            checked += 1
            cell_key = re.sub(r"[\W_]+", "", value, flags=re.UNICODE).casefold()
            if cell_key and cell_key in evidence_key:
                continue
            try:
                number = float(value.replace(",", "").rstrip("%％"))
            except ValueError:
                return False
            tolerance = max(1e-6, abs(number) * 1e-6)
            if not any(abs(source - number) <= tolerance for source in evidence_numbers):
                return False
    return checked > 0

def _guard_model_data_slides(spec: dict[str, Any], evidence_text: str) -> dict[str, Any]:
    """Downgrade untraceable model charts/tables before they reach PowerPoint."""
    slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        chart_valid = not isinstance(slide.get("chart"), dict) or _chart_is_supported_by_source(slide.get("chart"), evidence_text)
        table_valid = not isinstance(slide.get("table"), dict) or _table_is_supported_by_source(slide.get("table"), evidence_text)
        if not chart_valid:
            slide["chart"] = None
        if not table_valid:
            slide["table"] = None
        kind = str(slide.get("kind") or "")
        layout = str(slide.get("layout_id") or "")
        invalid_visual = (kind == "chart" or "chart" in layout) and not chart_valid
        invalid_visual = invalid_visual or (kind == "table" or layout in {"table", "data-table"}) and not table_valid
        if invalid_visual:
            slide["kind"] = "content"
            slide["layout_id"] = "content"
            slide["body"] = slide.get("body") or "当前资料不足以核验该图表口径，已保留文字结论供人工确认。"
    return spec


def _normalise_spec(raw: dict[str, Any] | None, job: PresentationJob, references: str) -> dict[str, Any]:
    fallback = _fallback_spec(job, references)
    if not raw:
        return fallback
    source_slides = raw.get("slides") if isinstance(raw.get("slides"), list) else []
    target_count = max(5, min(100, int(job.slide_count or 10)))
    slides: list[dict[str, Any]] = []
    for item in source_slides[:target_count]:
        if not isinstance(item, dict):
            continue
        title = _clip(str(item.get("title") or ""), 90)
        if not title:
            continue
        bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
        steps = item.get("steps") if isinstance(item.get("steps"), list) else []
        metric = item.get("metric")
        if isinstance(metric, dict):
            metric_value = _clip(str(metric.get("value") or ""), 24)
            metric_label = _clip(str(metric.get("label") or ""), 80)
        else:
            metric_value = _clip(str(metric or ""), 24)
            metric_label = _clip(str(item.get("metric_label") or ""), 80)
        template_slide_raw = str(item.get("template_slide") or "").strip()
        template_slide_raw = re.sub(r"^[Pp]\s*", "", template_slide_raw)
        try:
            template_slide = int(template_slide_raw or 0)
        except (TypeError, ValueError):
            template_slide = 0
        visual_intent = str(item.get("visual_intent") or "").strip().casefold()
        if visual_intent not in {"hero", "evidence", "comparison", "process", "metrics", "photo", "site", "image", "table", "chapter", "closing"}:
            visual_intent = ""
        density = str(item.get("density") or "medium").strip().casefold()
        if density not in {"low", "medium", "high"}:
            density = "medium"
        slides.append({
            "kind": str(item.get("kind") or "content"), "layout_id": str(item.get("layout_id") or ""),
            "page_type": str(item.get("page_type") or ""),
            "visual_intent": visual_intent, "density": density,
            "template_slide": template_slide if template_slide > 0 else None,
            "kicker": _clip(str(item.get("kicker") or ""), 60), "title": title,
            "subtitle": _clip(str(item.get("subtitle") or ""), 180),
            "summary": _clip(str(item.get("summary") or item.get("takeaway") or ""), 220),
            "body": _clip(str(item.get("body") or ""), 420),
            "bullets": [_clip(str(value), 110) for value in bullets if str(value).strip()][:5],
            "metric": metric_value, "metric_label": metric_label,
            "left_title": _clip(str(item.get("left_title") or ""), 80),
            "left_bullets": [_clip(str(value), 90) for value in (item.get("left_bullets") or []) if str(value).strip()][:4],
            "right_title": _clip(str(item.get("right_title") or ""), 80),
            "right_bullets": [_clip(str(value), 90) for value in (item.get("right_bullets") or []) if str(value).strip()][:4],
            "steps": [{"title": _clip(str(value.get("title") or ""), 60), "body": _clip(str(value.get("body") or ""), 100)} for value in steps if isinstance(value, dict)][:5],
            "image_slot": _clip(str(item.get("image_slot") or ""), 80),
            "chart": {
                "type": str((item.get("chart") or {}).get("type") or "bar"),
                "labels": [_clip(str(value), 32) for value in ((item.get("chart") or {}).get("labels") or [])][:8],
                "values": _chart_values((item.get("chart") or {}).get("values")),
                "unit": _clip(str((item.get("chart") or {}).get("unit") or ""), 20),
            } if isinstance(item.get("chart"), dict) else None,
            "table": {
                "columns": [_clip(str(value), 36) for value in ((item.get("table") or {}).get("columns") or [])][:12],
                "rows": [
                    [_clip(str(cell), 90) for cell in row[:12]]
                    for row in ((item.get("table") or {}).get("rows") or [])[:30]
                    if isinstance(row, list)
                ],
            } if isinstance(item.get("table"), dict) else None,
        })
    if len(slides) < 3:
        return fallback
    slides[0]["kind"], slides[0]["layout_id"] = "cover", "cover"
    if len(slides) > 4:
        slides[-1]["kind"], slides[-1]["layout_id"] = "closing", "closing"
    while len(slides) < target_count:
        filler = dict(fallback["slides"][min(len(slides), len(fallback["slides"]) - 2)])
        filler["_generated_filler"] = True
        slides.insert(-1, filler)
    return {"title": _clip(str(raw.get("title") or job.title), 180), "slides": slides[:target_count]}


def _enforce_layout_variety(spec: dict[str, Any], features: set[str]) -> dict[str, Any]:
    """Turn generic model output into a restrained but varied page rhythm."""
    if "layout_variety" not in features:
        return spec
    slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []
    generic_layouts = ("content", "content-stack", "content-rail", "content-emphasis")
    generic_cursor = 0
    recent: list[str] = []
    for index, item in enumerate(slides):
        if not isinstance(item, dict) or index == 0 or index == len(slides) - 1:
            continue
        kind = str(item.get("kind") or "content")
        current = str(item.get("layout_id") or kind or "content")
        if kind == "content" and current in {"", "content", "body", "default"}:
            bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
            if len(bullets) >= 3 and index % 5 == 0:
                item["kind"], item["layout_id"] = "cards", "cards"
                current = "cards"
            else:
                current = generic_layouts[generic_cursor % len(generic_layouts)]
                generic_cursor += 1
                item["kind"], item["layout_id"] = "content", current
        recent.append(current)
        if len(recent) < 3 or len(set(recent[-3:])) > 1:
            continue
        # Preserve explicit chart/table/image choices; rotate only generic
        # content pages where a repeated composition is visibly monotonous.
        if current not in {"content", "statement", "cards"}:
            continue
        if item.get("steps"):
            next_kind, next_layout = "timeline", "timeline"
        elif item.get("left_bullets") or item.get("right_bullets"):
            next_kind, next_layout = "comparison", "comparison"
        elif item.get("bullets") and len(item.get("bullets")) >= 3:
            next_kind, next_layout = "cards", "cards"
        else:
            next_kind, next_layout = "statement", "statement"
        item["kind"], item["layout_id"] = next_kind, next_layout
        recent[-1] = next_layout
    return spec


def _ensure_report_structure(
    spec: dict[str, Any],
    target_count: int,
    *,
    native_agenda: bool = False,
) -> tuple[dict[str, Any], list[str]]:
    """Repair missing report landmarks without rewriting authored evidence.

    Models occasionally return a valid list of content pages but omit the
    navigation landmarks that make a long management deck easy to follow.
    We repurpose only generic content shells, keep the exact page count, and
    leave authored charts/tables/images untouched.
    """
    slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []
    if len(slides) < 4:
        return spec, []
    repairs: list[str] = []
    section_indexes = [
        index for index, item in enumerate(slides)
        if isinstance(item, dict) and str(item.get("kind") or "").casefold() in {"section", "chapter", "section-divider"}
    ]

    def generic_candidate(preferred: int) -> int | None:
        order = sorted(range(1, max(1, len(slides) - 1)), key=lambda index: abs(index - preferred))
        for index in order:
            item = slides[index]
            if not isinstance(item, dict):
                continue
            kind = str(item.get("kind") or "content").casefold()
            if kind in {"cover", "closing", "section", "chapter", "section-divider", "chart", "table", "image", "photo"}:
                continue
            title = str(item.get("title") or "").strip()
            body = str(item.get("body") or "").strip()
            bullets = [str(value).strip() for value in (item.get("bullets") or []) if str(value).strip()]
            safe_marker = bool(item.get("_generated_filler"))
            if not safe_marker and (body or bullets):
                continue
            return index
        return None

    desired_sections = 2 if target_count > 12 else 1 if target_count >= 8 else 0
    for ordinal in range(len(section_indexes), desired_sections):
        preferred = round((ordinal + 1) * (len(slides) - 1) / (desired_sections + 1))
        index = generic_candidate(preferred)
        if index is None:
            break
        original = slides[index]
        title = _clip(str(original.get("title") or original.get("summary") or "重点议题"), 72)
        slides[index] = {
            "kind": "section",
            "layout_id": "section-divider",
            "visual_intent": "chapter",
            "density": "low",
            "template_slide": None,
            "kicker": f"PART {ordinal + 1:02d}",
            "title": title,
            "subtitle": _clip(str(original.get("summary") or original.get("body") or ""), 150),
            "summary": "",
            "body": "",
            "bullets": [],
            "left_bullets": [],
            "right_bullets": [],
            "steps": [],
            "chart": None,
            "table": None,
        }
        section_indexes.append(index)
        repairs.append(f"section@{index + 1}")

    if native_agenda and target_count >= 7 and not any(
        isinstance(item, dict) and str(item.get("kind") or "").casefold() in {"agenda", "toc"}
        for item in slides
    ):
        index = generic_candidate(1)
        if index is not None:
            chapter_titles = [
                _clip(str(slides[item_index].get("title") or slides[item_index].get("summary") or "重点议题"), 48)
                for item_index in sorted(section_indexes)
            ]
            chapter_titles.extend(
                _clip(str(item.get("title") or item.get("summary") or "重点议题"), 48)
                for item in slides[1:]
                if isinstance(item, dict)
                and str(item.get("kind") or "content").casefold() not in {"closing", "section", "chapter", "section-divider"}
            )
            slides[index] = {
                "kind": "agenda",
                "layout_id": "agenda",
                "visual_intent": "hero",
                "density": "low",
                "template_slide": None,
                "kicker": "",
                "title": "汇报目录",
                "subtitle": "按管理决策逻辑展开本次汇报。",
                "summary": "",
                "body": "",
                "bullets": list(dict.fromkeys(chapter_titles))[:6],
                "chart": None,
                "table": None,
            }
            repairs.append(f"agenda@{index + 1}")
    spec["slides"] = slides
    return spec, repairs


def _add_text(canvas: Canvas, slide: Any, text: str, left: float, top: float, width: float, height: float, *, size: int, color: str, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT, font: str = "Aptos") -> Any:
    box = slide.shapes.add_textbox(canvas.x(left), canvas.y(top), canvas.w(width), canvas.h(height))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = canvas.w(0.04); frame.margin_right = canvas.w(0.04); frame.margin_top = canvas.h(0.02); frame.margin_bottom = canvas.h(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]; paragraph.alignment = align
    run = paragraph.add_run(); run.text = _clip(text, 560); run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = _rgb(color)
    return box


def _shape(canvas: Canvas, slide: Any, shape_type: MSO_SHAPE, left: float, top: float, width: float, height: float, fill: str | None, *, transparency: int = 0, line: str | None = None) -> Any:
    shape = slide.shapes.add_shape(shape_type, canvas.x(left), canvas.y(top), canvas.w(width), canvas.h(height))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(fill); shape.fill.transparency = transparency
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line); shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def _add_picture_cover(canvas: Canvas, slide: Any, image_path: Path, left: float, top: float, width: float, height: float) -> Any:
    """Place a picture with a centered crop, preserving its aspect ratio."""
    picture = slide.shapes.add_picture(str(image_path), canvas.x(left), canvas.y(top), width=canvas.w(width), height=canvas.h(height))
    try:
        image_width, image_height = picture.image.size
        source_ratio = image_width / max(image_height, 1)
        target_ratio = width / max(height, 0.01)
        if source_ratio > target_ratio:
            visible = target_ratio / source_ratio
            picture.crop_left = picture.crop_right = (1.0 - visible) / 2.0
        elif source_ratio < target_ratio:
            visible = source_ratio / target_ratio
            picture.crop_top = picture.crop_bottom = (1.0 - visible) / 2.0
    except Exception:
        logger.debug("Unable to calculate image crop", exc_info=True)
    return picture


def _add_notes(slide: Any, job: PresentationJob, enabled: bool) -> None:
    if not enabled:
        return
    try:
        slide.notes_slide.notes_text_frame.text = f"[Sources]\nAIWeb user brief: {job.title}\nUploaded references are user-provided materials."
    except Exception:
        return


def _clear_slides(prs: Presentation) -> None:
    """Remove slides and relationships while preserving template masters."""
    for slide_id in list(prs.slides._sldIdLst):
        rel_id = getattr(slide_id, "rId", None)
        if rel_id:
            try:
                prs.part.drop_rel(rel_id)
            except KeyError:
                pass
        prs.slides._sldIdLst.remove(slide_id)


def _blank_layout(prs: Presentation) -> Any:
    layouts = sorted(prs.slide_layouts, key=lambda item: len(item.placeholders))
    return layouts[0] if layouts else prs.slide_layouts[0]


def _slide_palette(palette: dict[str, Any], *, kind: str, index: int, layout: str = "") -> dict[str, Any]:
    """Apply restrained per-page color variants for free-form decks.

    A coherent deck still benefits from small changes in stage color and
    accent weight. These variants stay inside the selected palette and are
    deliberately quiet for the state-owned-enterprise preset.
    """
    result = dict(palette)
    decor = str(result.get("decor") or "")
    if decor == "state":
        variants = (
            ("FFFFFF", "00479D", "5B9BD5"),
            ("F5F8FC", "0B5A9A", "78A8C9"),
            ("FAFCFE", "1E4E79", "6C9BC0"),
            ("F3F7FB", "135A86", "84B4D1"),
        )
        bg, accent, accent2 = variants[(max(index, 1) - 1) % len(variants)]
        # Keep the formal palette, but let page roles carry a different visual
        # temperature. This prevents a long state-owned-enterprise deck from
        # becoming a stack of identical white cards and blue headers.
        if kind in {"section", "chapter", "section-divider"}:
            bg, accent, accent2 = "0B3D70", "00479D", "8DC5E8"
        elif kind in {"metric", "chart", "table"}:
            bg = "EEF6FC"
        elif kind in {"comparison", "cards"}:
            bg = "F7FAFD"
        elif kind in {"image", "photo"}:
            bg = "F2F6F9"
        elif layout == "content-emphasis":
            bg = "EAF3FA"
        elif layout == "content-rail":
            bg = "F3F7FB"
        elif layout == "content-stack":
            bg = "FAFCFE"
        elif kind == "closing":
            bg, accent, accent2 = "F4F8FC", "0B5A9A", "7AA8C5"
        result.update({"bg": bg, "accent": accent, "accent2": accent2})
    elif decor == "aviation":
        variants = (
            ("FFFFFF", "005BAC", "4FA3D1"),
            ("F6F9FC", "07569F", "70AED2"),
            ("FBFDFE", "164F7A", "82B9D5"),
        )
        bg, accent, accent2 = variants[(max(index, 1) - 1) % len(variants)]
        if kind in {"section", "chapter", "section-divider"}:
            bg, accent, accent2 = "0A477F", "07569F", "88C2E5"
        elif kind in {"metric", "chart", "table"}:
            bg = "EEF6FC"
        elif kind in {"comparison", "cards", "timeline"}:
            bg = "F5F9FC"
        result.update({"bg": bg, "accent": accent, "accent2": accent2})
    elif decor == "aqua":
        variants = (
            ("F3FBFC", "32B8C7", "80D6DD"),
            ("EAF8F9", "239EAC", "64C8D2"),
            ("F8FDFD", "2AA8B5", "97DDE2"),
        )
        bg, accent, accent2 = variants[(max(index, 1) - 1) % len(variants)]
        if kind in {"section", "chapter", "section-divider"}:
            bg, accent, accent2 = "DFF5F7", "24A8B7", "8AD7DE"
        elif kind in {"metric", "chart", "table"}:
            bg = "E8F7F8"
        result.update({"bg": bg, "accent": accent, "accent2": accent2})
    elif decor == "security":
        variants = (
            ("FFFFFF", "07569F", "4874CB"),
            ("F5F8FC", "0B4C8A", "6A8FD5"),
            ("FAFCFE", "124F84", "7596D6"),
        )
        bg, accent, accent2 = variants[(max(index, 1) - 1) % len(variants)]
        if kind in {"section", "chapter", "section-divider"}:
            bg, accent, accent2 = "0A477F", "07569F", "6F98D8"
        elif kind in {"metric", "chart", "table"}:
            bg = "EEF3FB"
        elif kind in {"comparison", "cards", "timeline"}:
            bg = "F6F8FC"
        result.update({"bg": bg, "accent": accent, "accent2": accent2})
    elif decor == "tech":
        if kind in {"chart", "table", "metric"}:
            result["bg"] = "101A30"
        elif index % 3 == 0:
            result["bg"] = "0D1628"
    elif decor in {"editorial", "data", "ink", "swiss"} and index % 3 == 0:
        result["bg"] = "F8F6F1" if decor in {"editorial", "ink"} else "F1F5F6"
    return result


def _add_decor(canvas: Canvas, slide: Any, palette: dict[str, Any], *, index: int, kind: str = "content", layout: str = "") -> None:
    decor = palette.get("decor")
    if decor == "state":
        # Formal blue/white language with role-specific chrome. The chapter,
        # cover and closing pages intentionally do not share the ordinary
        # header so a deck has a visible opening, middle and ending rhythm.
        if kind in {"section", "chapter", "section-divider"}:
            return
        if kind in {"cover", "closing"}:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.72, 0.08, 5.62, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.72, 2.05, 0.012, palette["accent2"], transparency=5)
            _add_text(canvas, slide, f"{index:02d} / BRIEF", 10.6, 0.72, 1.9, 0.22, size=8, color=palette["muted"], align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 6.78, 11.7, 0.012, palette["accent2"], transparency=15)
            return

        quiet_header = kind in {"statement", "quote", "timeline", "cards", "comparison", "image", "photo"} or layout in {
            "content-stack", "content-rail", "content-emphasis",
        }
        if quiet_header:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.45, 1.1, 0.05, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.98, 0.45, 0.42, 0.05, palette["accent2"])
            _add_text(canvas, slide, f"{index:02d}", 11.62, 0.38, 0.72, 0.28, size=11, color=palette["accent"], bold=True, align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 6.78, 11.7, 0.012, palette["accent2"], transparency=15)
            _add_text(canvas, slide, "汇报材料 / INTERNAL BRIEFING", 0.82, 6.9, 4.2, 0.2, size=8, color=palette["muted"], font=_font_for(palette, "Aptos Mono"))
            return

        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, 0.86, palette["accent"])
        if kind in {"chart", "table"}:
            # Data pages use a square index marker and a two-tone data tab.
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.62, 0.06, 0.76, 0.74, palette["accent2"], line="FFFFFF")
            _add_text(canvas, slide, f"{index:02d}", 1.62, 0.19, 0.76, 0.28, size=18, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 11.68, 0.18, 0.28, 0.16, palette["accent2"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 12.02, 0.18, 0.58, 0.16, "FFFFFF", transparency=25)
        elif kind in {"metric", "comparison"}:
            # Decision pages get a paired marker that echoes their two-level
            # reading order without adding a dashboard-like widget.
            _shape(canvas, slide, MSO_SHAPE.OVAL, 1.63, 0.06, 0.74, 0.74, palette["accent"], line="FFFFFF")
            _add_text(canvas, slide, f"{index:02d}", 1.63, 0.19, 0.74, 0.28, size=18, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 11.68, 0.18, 0.34, 0.16, palette["accent2"], transparency=5)
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 12.1, 0.18, 0.5, 0.16, "FFFFFF", transparency=18)
        else:
            _shape(canvas, slide, MSO_SHAPE.OVAL, 1.63, 0.06, 0.74, 0.74, palette["accent"], line="FFFFFF")
            _add_text(canvas, slide, f"{index:02d}", 1.63, 0.19, 0.74, 0.28, size=18, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            if kind in {"image", "photo"}:
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 11.68, 0.18, 0.92, 0.16, palette["accent2"], transparency=10)
            elif index % 3 == 0:
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 11.68, 0.18, 0.92, 0.16, "FFFFFF", transparency=20)

        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 6.78, 11.7, 0.012, palette["accent2"], transparency=15)
        _add_text(canvas, slide, "汇报材料 / INTERNAL BRIEFING", 0.82, 6.9, 4.2, 0.2, size=8, color=palette["muted"], font=_font_for(palette, "Aptos Mono"))
    elif decor == "aviation":
        if kind in {"section", "chapter", "section-divider"}:
            return
        if kind in {"cover", "closing"}:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.7, 0.08, 5.8, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.7, 1.42, 0.035, palette["accent2"])
        else:
            marker_shape = MSO_SHAPE.RECTANGLE if kind in {"chart", "table"} else MSO_SHAPE.OVAL
            _shape(canvas, slide, marker_shape, 0.84, 0.18, 0.58, 0.58, palette["accent"], line="FFFFFF")
            _add_text(canvas, slide, f"{index:02d}", 0.84, 0.29, 0.58, 0.24, size=13, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.56, 0.45, 1.32, 0.035, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 2.94, 0.45, 0.45, 0.035, palette["accent2"])
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 6.82, 11.7, 0.012, palette["accent2"], transparency=12)
        _add_text(canvas, slide, "专项汇报 / AVIATION BRIEF", 0.82, 6.94, 4.0, 0.18, size=8, color=palette["muted"], font=_font_for(palette, "Aptos Mono"))
    elif decor == "aqua":
        if kind in {"section", "chapter", "section-divider"}:
            return
        # These quiet bars echo the planning reference without imposing the
        # same panel composition on every page role.
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 10.45, 0, 2.4, 0.36, palette["accent2"], transparency=55)
        if kind not in {"cover", "closing"}:
            _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.55, 0.34, 0.72, 0.46, palette["accent"])
            _add_text(canvas, slide, f"{index:02d}", 11.55, 0.43, 0.72, 0.22, size=12, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.48, 0.92, 0.045, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.8, 0.48, 0.34, 0.045, palette["accent2"])
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 6.82, 1.65, 0.48, palette["accent2"], transparency=68)
        _add_text(canvas, slide, "年度规划 / ANNUAL PLAN", 9.18, 6.94, 3.1, 0.18, size=8, color=palette["muted"], align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos Mono"))
    elif decor == "security":
        if kind in {"section", "chapter", "section-divider"}:
            return
        if kind not in {"cover", "closing"}:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, 0.22, "DDE7F5")
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.54, 0.62, 0.06, palette["accent"])
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.52, 0.54, 0.26, 0.06, palette["accent2"])
            if kind in {"metric", "chart", "table"}:
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 10.82, 0.38, 1.32, 0.36, palette["accent2"])
                _add_text(canvas, slide, "数据要点", 10.82, 0.45, 1.32, 0.18, size=9, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 7.24, CANVAS_W, 0.26, "D9E5F3")
        _add_text(canvas, slide, f"安护工作汇报  {index:02d}", 10.06, 7.27, 2.2, 0.14, size=7, color=palette["muted"], align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos Mono"))
    elif decor == "tech":
        for x in (0.35, 0.64, 0.93): _shape(canvas, slide, MSO_SHAPE.RECTANGLE, x, 0.35, 0.17, 0.012, palette["accent2"], transparency=10)
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 10.95, 0.42, 1.6, 0.012, palette["accent"], transparency=25)
        _shape(canvas, slide, MSO_SHAPE.OVAL, 11.4, 5.3, 1.2, 1.2, None, line=palette["accent"])
        _shape(canvas, slide, MSO_SHAPE.OVAL, 11.7, 5.6, 0.6, 0.6, None, line=palette["accent2"])
    elif decor == "glass":
        _shape(canvas, slide, MSO_SHAPE.OVAL, 9.8, -0.8, 4.8, 3.4, palette["accent"], transparency=82)
        _shape(canvas, slide, MSO_SHAPE.OVAL, -1.2, 5.1, 3.5, 3.0, palette["accent2"], transparency=88)
        _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.85, 0.55, 3.8, 5.8, None, line=palette["accent"])
    elif decor == "blueprint":
        for x in range(1, 13): _shape(canvas, slide, MSO_SHAPE.RECTANGLE, x, 0.55, 0.006, 6.25, palette["accent2"], transparency=91)
        for y in range(1, 7): _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.55, y, 12.2, 0.006, palette["accent2"], transparency=91)
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 10.95, 0.58, 1.55, 0.9, None, line=palette["accent2"])
    elif decor == "swiss":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 12.63, 0, 0.7, 7.5, palette["accent"])
    elif decor in {"editorial", "data"}:
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 0.34, 11.65, 0.012, palette["accent"])
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, 6.72, 11.65, 0.008, palette["muted"], transparency=15)
    elif decor == "ink":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.62, 0.6, 1.1, 0.04, palette["accent"])
        _shape(canvas, slide, MSO_SHAPE.ARC, 11.45, 0.32, 1.0, 0.72, None, line=palette["accent2"])
    elif decor == "soft":
        _shape(canvas, slide, MSO_SHAPE.OVAL, 11.55, -0.55, 3.0, 2.2, palette["accent"], transparency=80)
        _shape(canvas, slide, MSO_SHAPE.OVAL, -1.1, 6.05, 2.2, 1.8, palette["accent2"], transparency=88)
    elif decor == "vivid":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 10.75, 0, 2.55, 0.2, palette["accent"])
        _shape(canvas, slide, MSO_SHAPE.DIAMOND, 11.65, 5.45, 0.75, 0.75, palette["accent2"], transparency=2)
    if decor not in {"swiss", "editorial", "data", "state", "aviation", "aqua", "security"}:
        _add_text(canvas, slide, f"{index:02d} / BRIEF", 10.65, 6.93, 1.85, 0.2, size=8, color=palette["muted"], align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos Mono"))


def _add_page_header(
    canvas: Canvas,
    slide: Any,
    spec: dict[str, Any],
    palette: dict[str, Any],
    index: int,
    *,
    serif: bool = False,
    kind: str = "content",
    layout: str = "",
) -> None:
    kicker = str(spec.get("kicker") or "")
    if palette.get("decor") == "state":
        quiet_header = kind in {"statement", "quote", "timeline", "cards", "comparison", "image", "photo"} or layout in {
            "content-stack", "content-rail", "content-emphasis",
        }
        if quiet_header:
            if kicker:
                _add_text(canvas, slide, kicker.upper(), 0.82, 0.58, 4.4, 0.24, size=9, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
            _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 0.82, 0.88, 10.9, 0.62, size=27, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
            summary = str(spec.get("summary") or "")
            if summary:
                _add_text(canvas, slide, summary, 0.86, 1.58, 11.0, 0.48, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
            return
        if kicker:
            _add_text(canvas, slide, kicker.upper(), 2.52, 0.16, 3.7, 0.2, size=9, color="D7E8FB", bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 2.52, 0.38, 8.15, 0.35, size=21, color="FFFFFF", bold=True, font=_font_for(palette, "Aptos"))
        summary = str(spec.get("summary") or "")
        if summary:
            _add_text(canvas, slide, summary, 0.88, 1.12, 11.2, 0.42, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
        return
    if palette.get("decor") == "aviation":
        if kicker:
            _add_text(canvas, slide, kicker.upper(), 1.58, 0.18, 3.7, 0.2, size=9, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 1.58, 0.72, 10.35, 0.55, size=25, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        summary = str(spec.get("summary") or "")
        if summary:
            _add_text(canvas, slide, summary, 1.6, 1.38, 10.15, 0.44, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
        return
    if palette.get("decor") == "aqua":
        if kicker:
            _add_text(canvas, slide, kicker.upper(), 0.82, 0.62, 4.1, 0.22, size=9, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 0.82, 0.9, 10.45, 0.58, size=26, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        summary = str(spec.get("summary") or "")
        if summary:
            _add_text(canvas, slide, summary, 0.84, 1.57, 10.65, 0.44, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
        return
    if palette.get("decor") == "security":
        if kicker:
            _add_text(canvas, slide, kicker.upper(), 0.82, 0.38, 4.1, 0.2, size=9, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 0.82, 0.74, 10.8, 0.58, size=25, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        summary = str(spec.get("summary") or "")
        if summary:
            _add_text(canvas, slide, summary, 0.84, 1.4, 10.75, 0.44, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
        return
    if kicker:
        _add_text(canvas, slide, kicker.upper(), 0.82, 0.56, 4.7, 0.28, size=10, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
    _add_text(canvas, slide, str(spec.get("title") or "未命名页面"), 0.82, 0.9, 10.7, 0.72, size=28, color=palette["fg"], bold=True, font="Georgia" if serif else _font_for(palette, "Aptos"))
    summary = str(spec.get("summary") or "")
    if summary: _add_text(canvas, slide, summary, 0.88, 1.66, 10.8, 0.5, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))


def _add_takeaway(canvas: Canvas, slide: Any, text: str, palette: dict[str, Any], *, top: float = 6.12, show: bool = True) -> None:
    if not show or not text: return
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.82, top, 0.06, 0.46, palette["accent2"])
    _add_text(canvas, slide, f"TAKEAWAY  {_clip(text, 190)}", 1.02, top + 0.01, 10.9, 0.42, size=11, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos Mono"))


def _add_cover(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], job: PresentationJob, cover_image: Path | None) -> None:
    decor = str(palette.get("decor") or "")
    if decor == "aviation":
        if cover_image and cover_image.exists() and job.include_images:
            try:
                _add_picture_cover(canvas, slide, cover_image, 7.15, 0, 6.18, CANVAS_H)
                _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 6.55, 0, 1.35, CANVAS_H, palette["accent"], transparency=4)
            except Exception:
                logger.warning("Unable to embed aviation cover image", exc_info=True)
        else:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 7.35, 0.72, 4.85, 5.8, "E7F2F9")
            _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 6.72, 0.72, 1.35, 5.8, palette["accent"], transparency=3)
        _add_text(canvas, slide, str(spec.get("kicker") or "专项工作汇报"), 1.04, 1.04, 4.8, 0.3, size=12, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or job.title), 1.0, 1.72, 5.85, 1.55, size=38, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or f"{job.audience or '管理层'} · {job.purpose or '对齐判断与行动'}"), 1.04, 3.62, 5.45, 0.72, size=17, color=palette["muted"], font=_font_for(palette, "Aptos"))
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.04, 5.54, 1.35, 0.05, palette["accent"])
        _add_text(canvas, slide, "AVIATION BRIEFING", 1.04, 5.72, 3.4, 0.24, size=10, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        return
    if decor == "aqua":
        if cover_image and cover_image.exists() and job.include_images:
            try:
                _add_picture_cover(canvas, slide, cover_image, 0, 0, CANVAS_W, CANVAS_H)
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "E9F9FA", transparency=18)
            except Exception:
                logger.warning("Unable to embed aqua cover image", exc_info=True)
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 8.35, 0, 4.0, 0.65, palette["accent2"], transparency=48)
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 6.76, 3.1, 0.55, palette["accent2"], transparency=55)
        _add_text(canvas, slide, str(spec.get("title") or job.title), 0.92, 1.62, 7.0, 1.72, size=40, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or f"{job.audience or '团队与决策者'} · {job.purpose or '明确年度方向'}"), 0.98, 3.72, 6.35, 0.68, size=18, color=palette["muted"], font=_font_for(palette, "Aptos"))
        _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.98, 5.28, 2.25, 0.48, palette["accent"])
        _add_text(canvas, slide, str(spec.get("kicker") or "ANNUAL PLAN"), 0.98, 5.39, 2.25, 0.2, size=10, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        return
    if decor == "security":
        if cover_image and cover_image.exists() and job.include_images:
            try:
                _add_picture_cover(canvas, slide, cover_image, 0, 0, CANVAS_W, CANVAS_H)
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "0A477F", transparency=42)
            except Exception:
                logger.warning("Unable to embed security cover image", exc_info=True)
        else:
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "EAF1F8")
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 0, 5.15, CANVAS_H, "FFFFFF", transparency=2)
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 4.82, 0, 1.42, CANVAS_H, palette["accent"], transparency=0)
        _add_text(canvas, slide, str(spec.get("title") or job.title), 0.92, 1.62, 3.72, 1.65, size=34, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or f"{job.audience or '管理层'} · {job.purpose or '总结成效与问题'}"), 0.96, 3.68, 3.62, 0.72, size=16, color=palette["muted"], font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("kicker") or "年度总结"), 0.96, 5.24, 2.7, 0.3, size=11, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        return
    if cover_image and cover_image.exists() and job.include_images:
        try:
            _add_picture_cover(canvas, slide, cover_image, 7.72, 0.72, 4.65, 5.75)
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 7.72, 0.72, 4.65, 5.75, palette["bg"], transparency=68)
        except Exception: logger.warning("Unable to embed cover reference image", exc_info=True)
    _add_text(canvas, slide, str(spec.get("kicker") or "专项汇报 / SPECIAL BRIEFING"), 0.82, 0.9, 6.2, 0.3, size=12, color=palette["accent2"], bold=True, font=_font_for(palette, "Aptos Mono"))
    _add_text(canvas, slide, str(spec.get("title") or job.title), 0.8, 1.48, 6.45, 1.75, size=40, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
    _add_text(canvas, slide, str(spec.get("subtitle") or f"{job.audience or '团队与决策者'} · {job.purpose or '对齐认知并推动行动'}"), 0.85, 3.48, 6.1, 0.8, size=19, color=palette["muted"], font=_font_for(palette, "Aptos"))
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.85, 5.66, 1.4, 0.05, palette["accent"])
    _add_text(canvas, slide, "内部资料 / INTERNAL BRIEFING", 0.85, 5.82, 4.6, 0.28, size=11, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))


def _add_metric(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    metric, label = str(spec.get("metric") or "01"), str(spec.get("metric_label") or "关键指标")
    _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 2.42, 3.25, 2.65, palette["accent"], transparency=8, line=palette["accent"])
    fg = "FFFFFF" if palette["bg"] != "F8F8F5" else palette["fg"]
    _add_text(canvas, slide, metric, 1.18, 2.75, 2.55, 0.95, size=44, color=fg, bold=True, align=PP_ALIGN.CENTER)
    _add_text(canvas, slide, label, 1.15, 3.88, 2.6, 0.45, size=14, color=fg, align=PP_ALIGN.CENTER)
    bullets = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    for item_index, bullet in enumerate(bullets[:3]):
        y = 2.55 + item_index * 0.92
        _shape(canvas, slide, MSO_SHAPE.OVAL, 4.72, y + 0.1, 0.17, 0.17, palette["accent2"])
        _add_text(canvas, slide, str(bullet), 5.1, y, 6.8, 0.52, size=17, color=palette["fg"])
    _add_takeaway(canvas, slide, str(spec.get("summary") or "数字只负责让判断更快被看见。"), palette, show=show_takeaway)


def _add_chart(
    canvas: Canvas,
    slide: Any,
    spec: dict[str, Any],
    palette: dict[str, Any],
    *,
    show_takeaway: bool = True,
    bounds: tuple[float, float, float, float] | None = None,
) -> None:
    """Draw an editable bar/line/donut chart from verified facts.

    These are native PowerPoint shapes rather than a raster screenshot, so a
    user can still restyle, relabel, or replace values after downloading.
    """
    chart = spec.get("chart") if isinstance(spec.get("chart"), dict) else {}
    labels = [str(value) for value in chart.get("labels", []) if str(value).strip()]
    values: list[float] = []
    for value in chart.get("values", []):
        try:
            number = float(value)
        except (TypeError, ValueError):
            continue
        if math.isfinite(number):
            values.append(number)
    count = min(len(labels), len(values), 8)
    if count < 2:
        _add_content(canvas, slide, spec, palette, show_takeaway=show_takeaway)
        return
    labels, values = labels[:count], values[:count]
    chart_type = str(chart.get("type") or "bar").casefold()
    bounded = bounds is not None
    left, top, width, height = bounds or (1.0, 2.38, 10.95, 3.65)
    # Native-template overlays receive the exact geometry of a replaceable
    # body slot. Do not inflate a narrow slot: every child shape below must
    # remain inside that host rectangle so it cannot cover adjacent chrome.
    if bounded:
        width = max(0.8, width)
        height = max(0.8, height)
    else:
        width = max(4.6, width)
        height = max(2.0, height)
    pad_x = min(0.18, max(0.04, width * 0.025))
    pad_y = min(0.16, max(0.04, height * 0.035))
    inner_left = left + pad_x
    inner_top = top + pad_y
    inner_right = left + width - pad_x
    inner_bottom = top + height - pad_y
    inner_width = max(0.4, inner_right - inner_left)
    inner_height = max(0.4, inner_bottom - inner_top)
    max_value = max(max(values), 1.0)
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, inner_left, inner_bottom - 0.012, inner_width, 0.012, palette["muted"], transparency=35)
    for index in range(4):
        y = inner_bottom - (index / 3) * inner_height
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, inner_left, y, inner_width, 0.008, palette["muted"], transparency=76)
    if chart_type == "line":
        points: list[tuple[float, float]] = []
        denom = max(count - 1, 1)
        low = min(min(values), 0.0)
        span = max(max_value - low, 1.0)
        label_height = min(0.34, max(0.18, inner_height * 0.18))
        plot_top = inner_top + min(0.28, max(0.08, inner_height * 0.12))
        plot_bottom = max(plot_top + 0.12, inner_bottom - label_height)
        plot_height = max(0.12, plot_bottom - plot_top)
        for index, (label, value) in enumerate(zip(labels, values)):
            x = inner_left + (inner_width * index / denom if denom else 0)
            y = plot_bottom - (value - low) / span * plot_height
            points.append((x, y))
            if index:
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, canvas.x(points[index - 1][0]), canvas.y(points[index - 1][1]), canvas.x(x), canvas.y(y))
                connector.line.color.rgb = _rgb(palette["accent"]); connector.line.width = Pt(2.2)
            _shape(canvas, slide, MSO_SHAPE.OVAL, x - 0.09, y - 0.09, 0.18, 0.18, palette["accent2"], line=palette["accent"])
            value_width = min(0.76, max(0.42, inner_width / max(count, 1)))
            value_left = max(left, min(x - value_width / 2, left + width - value_width))
            value_top = max(top, min(y - 0.28, top + height - 0.22))
            _add_text(canvas, slide, _clip(f"{value:g}{chart.get('unit') or ''}", 18), value_left, value_top, value_width, 0.22, size=9, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            label_width = min(1.1, max(0.42, inner_width / max(count, 1) * 1.35))
            label_left = max(left, min(x - label_width / 2, left + width - label_width))
            _add_text(canvas, slide, _clip(label, 14), label_left, inner_bottom - label_height, label_width, label_height, size=10, color=palette["muted"], align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
    elif chart_type == "donut":
        total = sum(abs(value) for value in values) or 1.0
        ring_size = min(3.0, inner_height, inner_width * 0.42)
        ring_size = max(0.42, ring_size)
        ring_left = inner_left
        ring_top = inner_top + max(0.0, (inner_height - ring_size) / 2)
        # PIE is used instead of BLOCK_ARC here.  PowerPoint interprets the
        # latter's adjustments as thickness/rotation, not start/end angles,
        # which makes multiple percentage slices overlap into one color.
        chart_colors = [
            str(palette.get("accent") or "00479D"),
            str(palette.get("accent2") or "5B9BD5"),
            str(palette.get("accent3") or "2E75B6"),
            str(palette.get("accent4") or "7FB3D5"),
        ]
        start = 0.0
        for index, value in enumerate(values):
            fraction = abs(value) / total
            end = min(360.0, start + fraction * 360.0)
            arc = _shape(canvas, slide, MSO_SHAPE.PIE, ring_left, ring_top, ring_size, ring_size, chart_colors[index % len(chart_colors)], line="FFFFFF")
            try:
                # python-pptx normalizes preset geometry adjustments to a
                # 100,000-unit scale, while PIE's OOXML angles use 60,000
                # units per degree.  Multiplying degrees by .6 gives the
                # correct editable start/end angles in the generated file.
                arc.adjustments[0] = start * 0.6
                arc.adjustments[1] = end * 0.6
            except Exception:
                pass
            start = end
        hole_size = max(0.22, ring_size * 0.42)
        hole_offset = (ring_size - hole_size) / 2
        _shape(canvas, slide, MSO_SHAPE.OVAL, ring_left + hole_offset, ring_top + hole_offset, hole_size, hole_size, palette["bg"])
        metric_width = max(0.2, hole_size - 0.06)
        _add_text(canvas, slide, _clip(str(spec.get("metric") or "构成"), 16), ring_left + (ring_size - metric_width) / 2, ring_top + (ring_size - 0.25) / 2, metric_width, 0.25, size=14, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER)
        legend_gap = min(0.42, max(0.08, inner_width * 0.04))
        legend_left = ring_left + ring_size + legend_gap
        legend_width = inner_right - legend_left
        legend_below = legend_width < 1.15
        if legend_below:
            legend_left = inner_left
            legend_top = min(inner_bottom - 0.2, ring_top + ring_size + 0.08)
            legend_width = inner_width
            available_legend_height = max(0.2, inner_bottom - legend_top)
        else:
            legend_top = inner_top
            available_legend_height = inner_height
        row_height = min(0.43, max(0.18, available_legend_height / max(count, 1)))
        for index, (label, value) in enumerate(zip(labels, values)):
            y = legend_top + index * row_height
            swatch = chart_colors[index % len(chart_colors)]
            swatch_size = min(0.14, max(0.06, row_height * 0.42))
            swatch_top = min(inner_bottom - swatch_size, y + (row_height - swatch_size) / 2)
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, legend_left, swatch_top, swatch_size, swatch_size, swatch)
            text_left = min(legend_left + swatch_size + 0.12, left + width - 0.2)
            text_width = max(0.2, left + width - text_left)
            text_top = max(top, min(y, top + height - row_height))
            _add_text(canvas, slide, f"{_clip(label, 18)}  {value:g}{chart.get('unit') or ''}", text_left, text_top, text_width, row_height, size=13, color=palette["fg"], font=_font_for(palette, "Aptos"))
    else:
        horizontal = count > 6 or any(len(label) > 7 for label in labels)
        if horizontal:
            row_height = min(0.46, max(0.16, inner_height / max(count, 1)))
            label_width = min(2.15, max(0.72, inner_width * 0.23))
            value_width = min(0.9, max(0.38, inner_width * 0.13))
            bar_left = min(inner_right - 0.24, inner_left + label_width + 0.18)
            bar_width = max(0.22, inner_right - bar_left - value_width - 0.08)
            max_abs = max(max(abs(value) for value in values), 1.0)
            for index, (label, value) in enumerate(zip(labels, values)):
                y = inner_top + index * row_height
                text_height = max(0.12, row_height - 0.04)
                _add_text(canvas, slide, _clip(label, 18), inner_left, y + 0.02, label_width, text_height, size=11, color=palette["muted"], align=PP_ALIGN.RIGHT, font=_font_for(palette, "Aptos"))
                bar_height = min(0.16, max(0.06, row_height * 0.42))
                bar_top = min(inner_bottom - bar_height, y + (row_height - bar_height) / 2)
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height, palette["muted"], transparency=86)
                actual_width = max(0.04, bar_width * abs(value) / max_abs)
                _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, min(bar_width, actual_width), bar_height, palette["accent"] if index % 2 == 0 else palette["accent2"], transparency=4)
                value_left = min(bar_left + bar_width + 0.04, inner_right - value_width)
                _add_text(canvas, slide, _clip(f"{value:g}{chart.get('unit') or ''}", 18), value_left, y, value_width, row_height, size=10, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos Mono"))
        else:
            label_height = min(0.42, max(0.18, inner_height * 0.18))
            plot_bottom = max(inner_top + 0.15, inner_bottom - label_height)
            plot_height = max(0.15, plot_bottom - inner_top)
            bar_width = min(0.72, max(0.12, (inner_width * 0.72) / max(count, 1)))
            gap = max(0.04, (inner_width - count * bar_width) / max(count + 1, 1))
            baseline = plot_bottom
            for index, (label, value) in enumerate(zip(labels, values)):
                x = inner_left + gap + index * (bar_width + gap)
                bar_height = max(0.04, plot_height * max(0.0, value) / max_value)
                _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, baseline - bar_height, bar_width, bar_height, palette["accent"] if index % 2 == 0 else palette["accent2"], transparency=4)
                value_height = min(0.25, max(0.14, label_height * 0.6))
                value_top = max(inner_top, baseline - bar_height - value_height)
                value_width = min(bar_width + 0.44, inner_width / max(count, 1) * 1.4)
                value_left = max(left, min(x + (bar_width - value_width) / 2, left + width - value_width))
                _add_text(canvas, slide, _clip(f"{value:g}{chart.get('unit') or ''}", 18), value_left, value_top, value_width, value_height, size=10, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
                label_width = min(bar_width + 0.84, inner_width / max(count, 1) * 1.55)
                label_left = max(left, min(x + (bar_width - label_width) / 2, left + width - label_width))
                _add_text(canvas, slide, _clip(label, 10), label_left, inner_bottom - label_height, label_width, label_height, size=10, color=palette["muted"], align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "图表只呈现资料中可验证的变化，结论仍需结合上下文判断。"), palette, show=show_takeaway)


def _native_chart_bounds(slide_entry: dict[str, Any], canvas_px: dict[str, Any] | None = None) -> tuple[float, float, float, float]:
    """Find a chart area wholly contained by the largest replaceable body slot.

    Analyzer geometry is expressed in source pixels while the renderer uses
    the standard 13.333 x 7.5 inch canvas.  Keep the conversion affine and add
    only a small inset; broad hard-coded clamps made overlays drift over
    neighboring template controls on non-16:9 decks.
    """
    candidates = [
        slot for slot in (slide_entry.get("slots") or [])
        if isinstance(slot, dict) and _native_slot_is_body_like(slot)
    ]
    slot = max(
        candidates,
        key=lambda item: _native_slot_geometry(item)[2] * _native_slot_geometry(item)[3],
        default=None,
    )
    if slot is None:
        return (0.95, 2.15, 11.35, 3.85)
    x, y, width, height = _native_slot_geometry(slot)
    canvas_px = canvas_px or {"width": 1280, "height": 720}
    source_width = max(float(canvas_px.get("width") or 1280), 1.0)
    source_height = max(float(canvas_px.get("height") or 720), 1.0)
    scale_x = CANVAS_W / source_width
    scale_y = CANVAS_H / source_height
    left = x * scale_x
    top = y * scale_y
    chart_width = max(0.8, width * scale_x)
    chart_height = max(0.8, height * scale_y)
    # Stay inside the physical slide even when a third-party analyzer reports
    # a slightly oversized text frame. The inset is applied symmetrically and
    # is included in the returned bounds contract.
    inset_x = min(0.16, chart_width * 0.04)
    inset_y = min(0.14, chart_height * 0.04)
    left = max(0.2, min(left + inset_x, CANVAS_W - 0.4))
    top = max(0.2, min(top + inset_y, CANVAS_H - 0.4))
    right = min(CANVAS_W - 0.2, left + max(0.25, chart_width - 2 * inset_x))
    bottom = min(CANVAS_H - 0.2, top + max(0.25, chart_height - 2 * inset_y))
    return (left, top, max(0.25, right - left), max(0.25, bottom - top))


def _native_hidden_empty_shape_ids(
    slide_entry: dict[str, Any],
    replacements: list[dict[str, Any]],
) -> list[int]:
    """Find grouped, textless artwork that would otherwise become a blank tile.

    Native templates frequently put a text box on top of a colored polygon.
    The analyzer records both as slots, but clearing the text box alone leaves
    the polygon visible as an unexplained empty control.  Only grouped shapes
    with no text nodes are eligible; ordinary agenda bars and page chrome stay
    intact.
    """
    replacement_map = {
        str(item.get("slot_id") or ""): str(item.get("text") or "")
        for item in replacements
        if isinstance(item, dict)
    }
    slots = [slot for slot in (slide_entry.get("slots") or []) if isinstance(slot, dict)]

    def geometry(slot: dict[str, Any]) -> tuple[float, float, float, float]:
        value = slot.get("geometry") or slot.get("local_geometry") or {}
        return (
            float(value.get("x") or 0),
            float(value.get("y") or 0),
            float(value.get("width") or 0),
            float(value.get("height") or 0),
        )

    def contains(outer: dict[str, Any], inner: dict[str, Any]) -> bool:
        ox, oy, ow, oh = geometry(outer)
        ix, iy, iw, ih = geometry(inner)
        return (
            ow > 0 and oh > 0 and iw > 0 and ih > 0
            and ix >= ox - 4 and iy >= oy - 4
            and ix + iw <= ox + ow + 4 and iy + ih <= oy + oh + 4
        )

    result: list[int] = []
    for slot in slots:
        if int(slot.get("group_depth") or 0) <= 0:
            continue
        if int(slot.get("text_node_count") or 0) > 0 or str(slot.get("text") or "").strip():
            continue
        # A polygon is a removable background only when a real text slot sits
        # inside it and that text slot was explicitly cleared by the fill
        # plan. Standalone grouped circles/icons remain template decoration.
        paired_text = [
            candidate for candidate in slots
            if candidate is not slot
            and int(candidate.get("text_node_count") or 0) > 0
            and str(candidate.get("text") or "").strip()
            and contains(slot, candidate)
            and str(candidate.get("slot_id") or "") in replacement_map
            and not replacement_map[str(candidate.get("slot_id") or "")].strip()
        ]
        if not paired_text:
            continue
        shape_id = str(slot.get("shape_id") or "").strip()
        if shape_id.isdigit():
            result.append(int(shape_id))
    return sorted(set(result))


def _iter_native_shapes(shapes: Any):
    """Yield top-level and nested group shapes from a python-pptx collection."""
    for shape in shapes:
        yield shape
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _iter_native_shapes(nested)


def _remove_native_empty_shapes(output_path: Path, plan: dict[str, Any]) -> int:
    """Remove unused grouped artwork after the native fill has been applied."""
    targets = {
        index: {
            int(shape_id)
            for shape_id in (item.get("hide_shapes") or [])
            if isinstance(shape_id, int) or str(shape_id).isdigit()
        }
        for index, item in enumerate(plan.get("slides") or [])
        if isinstance(item, dict) and item.get("hide_shapes")
    }
    if not targets:
        return 0
    temporary = output_path.with_name(f"{output_path.stem}.empty-shapes.tmp.pptx")
    removed = 0
    try:
        prs = Presentation(str(output_path))
        for slide_index, shape_ids in targets.items():
            if slide_index >= len(prs.slides):
                continue
            for shape in list(_iter_native_shapes(prs.slides[slide_index].shapes)):
                if int(getattr(shape, "shape_id", -1)) not in shape_ids:
                    continue
                element = getattr(shape, "_element", None)
                parent = element.getparent() if element is not None else None
                if parent is None:
                    continue
                parent.remove(element)
                removed += 1
        if removed:
            prs.save(str(temporary))
            temporary.replace(output_path)
        else:
            temporary.unlink(missing_ok=True)
    except Exception:
        logger.warning("Unable to remove unused native template artwork from %s", output_path, exc_info=True)
        temporary.unlink(missing_ok=True)
    return removed


def _native_overlay_palette(
    slide: Any,
    base_palette: dict[str, Any],
    bounds: tuple[float, float, float, float],
) -> dict[str, str]:
    """Resolve chart colors against the actual cloned slide surface."""
    background = ""
    try:
        background = _safe_hex_color(slide.background.fill.fore_color) or ""
    except (AttributeError, ValueError, TypeError):
        pass

    left, top, width, height = bounds
    center_x = int(Inches(left + width / 2))
    center_y = int(Inches(top + height / 2))
    target_area = max(int(Inches(width)) * int(Inches(height)), 1)
    # Shapes are stored in z-order. The last sufficiently large filled shape
    # covering the chart center is the visible local surface.
    for shape in getattr(slide, "shapes", []):
        shape_left = int(getattr(shape, "left", 0))
        shape_top = int(getattr(shape, "top", 0))
        shape_width = int(getattr(shape, "width", 0))
        shape_height = int(getattr(shape, "height", 0))
        if shape_width * shape_height < target_area * 0.55:
            continue
        if not (
            shape_left <= center_x <= shape_left + shape_width
            and shape_top <= center_y <= shape_top + shape_height
        ):
            continue
        fill = _shape_fill_color(shape)
        if fill:
            background = fill

    background = background or str(base_palette.get("bg") or "FFFFFF")
    light_background = _relative_luminance(background) > 0.45
    accent_fallback = "00479D" if light_background else "69B7FF"
    accent2_fallback = "3B82C4" if light_background else "7DD3FC"
    accent = str(base_palette.get("accent") or accent_fallback)
    accent2 = str(base_palette.get("accent2") or accent2_fallback)
    if _color_contrast(accent, background) < 2.4:
        accent = accent_fallback
    if _color_contrast(accent2, background) < 2.1:
        accent2 = accent2_fallback
    return {
        "bg": background,
        "fg": _contrasting_color(str(base_palette.get("fg") or "263238"), background, minimum=4.5),
        "muted": _contrasting_color(
            str(base_palette.get("muted") or "5B6472"),
            background,
            minimum=3.0,
            muted=True,
        ),
        "accent": accent,
        "accent2": accent2,
        "font": str(base_palette.get("font") or "微软雅黑"),
        "decor": "state",
    }


def _apply_native_chart_overlays(
    output_path: Path,
    plan: dict[str, Any],
    palette: dict[str, Any],
    *,
    canvas_px: dict[str, Any] | None = None,
) -> int:
    """Add editable shape charts when a native template has no chart object.

    The overlay is deliberately opt-in per planned slide and is written to a
    temporary PPTX before replacement. If a third-party package cannot survive
    python-pptx round-tripping, the original native-filled output remains
    intact and the caller can still deliver the verified text/table result.
    """
    overlays = [
        (index, item)
        for index, item in enumerate(plan.get("slides") or [])
        if isinstance(item, dict) and isinstance(item.get("shape_chart"), dict)
    ]
    if not overlays:
        return 0
    temporary = output_path.with_name(f"{output_path.stem}.chart-overlay.tmp.pptx")
    base_palette = {
        "bg": str(palette.get("bg") or "FFFFFF"),
        "fg": str(palette.get("fg") or "263238"),
        "muted": str(palette.get("muted") or "5B6472"),
        "accent": str(palette.get("accent") or "00479D"),
        "accent2": str(palette.get("accent2") or "5B9BD5"),
        "font": str(palette.get("font") or "微软雅黑"),
        "decor": "state",
    }
    try:
        prs = Presentation(str(output_path))
        canvas = Canvas(prs)
        applied = 0
        for index, item in overlays:
            if index >= len(prs.slides):
                continue
            payload = item.get("shape_chart") or {}
            chart_spec = {
                "chart": payload,
                "summary": "",
                "metric": str(payload.get("metric") or payload.get("unit") or "") if isinstance(payload, dict) else "",
            }
            bounds = tuple(
                payload.get("bounds")
                or _native_chart_bounds(item.get("source_analysis") or {}, canvas_px)
            )
            safe_palette = _native_overlay_palette(prs.slides[index], base_palette, bounds)
            _add_chart(
                canvas,
                prs.slides[index],
                chart_spec,
                safe_palette,
                show_takeaway=False,
                bounds=bounds,
            )
            summary = str(payload.get("summary") or "").strip()
            if summary:
                _add_text(canvas, prs.slides[index], _clip(summary, 150), 0.92, 1.55, 11.0, 0.38, size=11, color=safe_palette["muted"], font=safe_palette["font"])
            applied += 1
        if applied:
            prs.save(str(temporary))
            temporary.replace(output_path)
        else:
            temporary.unlink(missing_ok=True)
        return applied
    except Exception:
        logger.warning("Unable to add native shape-chart overlays to %s", output_path, exc_info=True)
        temporary.unlink(missing_ok=True)
        return 0


def _add_image_split(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], image_path: Path | None, *, show_takeaway: bool = True) -> None:
    """Use an uploaded image as a visual evidence panel while keeping text editable."""
    body = str(spec.get("body") or "")
    bullets = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    if body:
        _add_text(canvas, slide, body, 0.9, 2.38, 5.2, 0.88, size=17, color=palette["muted"], font=_font_for(palette, "Aptos"))
    for bullet_index, bullet in enumerate(bullets[:4]):
        y = 3.48 + bullet_index * 0.56
        _shape(canvas, slide, MSO_SHAPE.OVAL, 0.98, y + 0.11, 0.13, 0.13, palette["accent2"])
        _add_text(canvas, slide, str(bullet), 1.25, y, 4.8, 0.4, size=14, color=palette["fg"], font=_font_for(palette, "Aptos"))
    if image_path and image_path.exists():
        try:
            _add_picture_cover(canvas, slide, image_path, 6.55, 2.25, 5.55, 3.72)
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 6.55, 2.25, 5.55, 3.72, None, line=palette["accent"])
        except Exception:
            logger.warning("Unable to embed image evidence", exc_info=True)
    else:
        _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.55, 2.25, 5.55, 3.72, palette["accent"], transparency=86, line=palette["accent"])
        _add_text(canvas, slide, str(spec.get("image_slot") or "视觉证据位"), 7.0, 3.65, 4.6, 0.5, size=18, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER)
    _add_takeaway(canvas, slide, str(spec.get("summary") or "把资料中的视觉证据与结论放在同一页，减少来回切换。"), palette, show=show_takeaway)


def _add_statement(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """A calm conclusion page with generous whitespace and one evidence line."""
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.9, 2.35, 0.1, 2.45, palette["accent"])
    _add_text(canvas, slide, str(spec.get("body") or spec.get("title") or ""), 1.35, 2.35, 8.9, 1.45, size=29, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
    evidence = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    if evidence:
        _add_text(canvas, slide, "  ·  ".join(_clip(str(item), 90) for item in evidence[:3]), 1.38, 4.25, 9.3, 0.55, size=14, color=palette["muted"], font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "先统一判断，再讨论执行细节。"), palette, show=show_takeaway)


def _add_cards(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """Three/four compact information cards with alternating emphasis."""
    items = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    if not items:
        items = [spec.get("body") or "事实", "判断", "动作"]
    items = [str(item) for item in items if str(item).strip()][:4]
    count = max(1, len(items))
    gap = 0.22
    width = (11.4 - gap * (count - 1)) / count
    for index, item in enumerate(items):
        left = 0.86 + index * (width + gap)
        fill = palette["accent"] if index == 0 else palette["bg"]
        line = palette["accent"] if index == 0 else palette["accent2"]
        text_color = "FFFFFF" if index == 0 else palette["fg"]
        if palette.get("decor") in {"state", "aviation", "aqua", "security"}:
            # Alternate a highlighted first block with quiet paper-like blocks
            # so a formal card page has hierarchy without looking like a SaaS
            # dashboard.
            panel_fill = fill if index == 0 else "FFFFFF"
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 2.45, width, 2.9, panel_fill, line=line)
            if index > 0:
                _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 2.45, width, 0.08, palette["accent2"])
        else:
            _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.45, width, 2.9, fill, transparency=0 if index == 0 else 3, line=line)
        _add_text(canvas, slide, f"0{index + 1}", left + 0.24, 2.7, width - 0.48, 0.28, size=12, color=palette["accent2"] if index == 0 else palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, _clip(item, 150), left + 0.24, 3.18, width - 0.48, 1.45, size=17, color=text_color, bold=index == 0, font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "把信息拆成几个可快速扫描的判断单元。"), palette, show=show_takeaway)


def _add_table(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """Draw a small native editable table when the source provides rows."""
    data = spec.get("table") if isinstance(spec.get("table"), dict) else {}
    columns = [str(item) for item in (data.get("columns") or []) if str(item).strip()][:8]
    rows = [row for row in (data.get("rows") or []) if isinstance(row, list)][:8]
    if not columns and isinstance(spec.get("chart"), dict):
        chart = spec["chart"]
        columns = ["类别", str(chart.get("unit") or "数值")]
        rows = [[label, value] for label, value in zip(chart.get("labels") or [], chart.get("values") or [])][:8]
    if not columns or not rows:
        _add_content(canvas, slide, spec, palette, show_takeaway=show_takeaway)
        return
    col_count = len(columns)
    row_count = min(len(rows) + 1, 9)
    table_shape = slide.shapes.add_table(row_count, col_count, canvas.x(0.86), canvas.y(2.28), canvas.w(11.55), canvas.h(3.7))
    table = table_shape.table
    for col, heading in enumerate(columns):
        cell = table.cell(0, col)
        cell.text = _clip(heading, 36)
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(palette["accent"])
    for row_index, row in enumerate(rows[: row_count - 1], start=1):
        values = list(row)[:col_count]
        for col in range(col_count):
            cell = table.cell(row_index, col)
            cell.text = _clip(str(values[col]) if col < len(values) else "", 80)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb("F7F9FC" if row_index % 2 else "FFFFFF")
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = canvas.w(0.06); cell.margin_right = canvas.w(0.06)
            cell.margin_top = canvas.h(0.02); cell.margin_bottom = canvas.h(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = _font_for(palette, "Aptos")
                    run.font.size = Pt(11)
                    run.font.bold = cell in table.rows[0].cells
                    run.font.color.rgb = _rgb("FFFFFF" if cell in table.rows[0].cells else palette["fg"])
    _add_takeaway(canvas, slide, str(spec.get("summary") or "表格保留原始口径，便于会后继续核对和编辑。"), palette, show=show_takeaway)


def _add_agenda(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any]) -> None:
    """Render a restrained directory page using the selected report family."""
    items = [str(item) for item in (spec.get("bullets") or []) if str(item).strip()][:6]
    if not items:
        items = ["总体情况", "重点工作", "问题与风险", "下一步安排"]
    decor = str(palette.get("decor") or "")
    if decor == "aqua":
        count = min(4, len(items))
        width = 10.95 / max(count, 1)
        for item_index, item in enumerate(items[:count]):
            left = 0.9 + item_index * width
            _add_text(canvas, slide, f"{item_index + 1:02d}", left, 2.55, width - 0.2, 0.5, size=22, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 3.16, width - 0.32, 0.055, palette["accent2"])
            _add_text(canvas, slide, item, left, 3.5, width - 0.38, 1.05, size=17, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        return
    if decor == "security":
        for item_index, item in enumerate(items[:5]):
            top = 2.22 + item_index * 0.68
            _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 2.0, top, 8.6, 0.48, palette["accent"] if item_index == 0 else "E7EEF8", line=palette["accent2"])
            number_color = "FFFFFF" if item_index == 0 else palette["accent"]
            text_color = "FFFFFF" if item_index == 0 else palette["fg"]
            _add_text(canvas, slide, f"{item_index + 1:02d}", 2.24, top + 0.1, 0.72, 0.22, size=11, color=number_color, bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
            _add_text(canvas, slide, item, 3.15, top + 0.08, 6.7, 0.26, size=14, color=text_color, bold=True, font=_font_for(palette, "Aptos"))
        return
    for item_index, item in enumerate(items[:5]):
        top = 2.22 + item_index * 0.68
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.2, top, 0.74, 0.48, palette["accent"])
        _add_text(canvas, slide, f"{item_index + 1:02d}", 1.2, top + 0.1, 0.74, 0.22, size=11, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 2.04, top, 8.7, 0.48, "FFFFFF", line=palette["accent2"])
        _add_text(canvas, slide, item, 2.34, top + 0.08, 7.95, 0.26, size=14, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))


def _add_section_divider(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any]) -> None:
    decor = str(palette.get("decor") or "")
    if decor == "aqua":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "DFF5F7")
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 10.1, 0, 2.8, 0.5, palette["accent2"], transparency=50)
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 6.72, 2.3, 0.56, palette["accent2"], transparency=58)
        _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.83, 1.65, 1.15, 0.7, palette["accent"])
        _add_text(canvas, slide, str(spec.get("kicker") or "01"), 5.83, 1.79, 1.15, 0.32, size=20, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "下一部分"), 1.2, 2.78, 10.9, 0.95, size=39, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("summary") or spec.get("subtitle") or ""), 2.2, 4.02, 8.9, 0.58, size=17, color=palette["muted"], align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        return
    if decor == "security":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "0A477F")
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 0, 4.55, CANVAS_H, "FFFFFF", transparency=4)
        _add_text(canvas, slide, str(spec.get("kicker") or "01"), 1.0, 1.02, 2.7, 1.05, size=52, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.55, 2.45, 7.05, 1.08, palette["accent"])
        _add_text(canvas, slide, str(spec.get("title") or "下一部分"), 0.82, 2.68, 6.48, 0.56, size=28, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("summary") or spec.get("subtitle") or ""), 7.95, 2.7, 4.1, 1.0, size=17, color="E5EEF8", bold=True, font=_font_for(palette, "Aptos"))
        return
    if decor == "aviation":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "0A477F")
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 2.35, CANVAS_W, 2.02, "062F57", transparency=28)
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.86, 1.14, 0.1, 4.65, palette["accent2"])
        _add_text(canvas, slide, str(spec.get("kicker") or "PART 01"), 1.34, 1.38, 3.1, 0.3, size=12, color="CBE7F7", bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "下一部分"), 1.28, 2.68, 10.2, 0.92, size=39, color="FFFFFF", bold=True, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("summary") or spec.get("subtitle") or ""), 1.32, 4.02, 8.9, 0.58, size=18, color="D7E8F5", font=_font_for(palette, "Aptos"))
        return
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, palette["accent"])
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.86, 1.05, 0.12, 4.55, palette["accent2"])
    _add_text(canvas, slide, str(spec.get("kicker") or "SECTION"), 1.28, 1.2, 3.2, 0.3, size=12, color="D7E8FB", bold=True, font=_font_for(palette, "Aptos Mono"))
    _add_text(canvas, slide, str(spec.get("title") or "下一部分"), 1.24, 2.05, 10.3, 1.15, size=42, color="FFFFFF", bold=True, font=_font_for(palette, "Aptos"))
    _add_text(canvas, slide, str(spec.get("summary") or spec.get("subtitle") or ""), 1.28, 3.65, 8.7, 0.7, size=19, color="D7E8FB", font=_font_for(palette, "Aptos"))


def _add_content(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, serif: bool = False, show_takeaway: bool = True) -> None:
    body = str(spec.get("body") or "")
    bullets = spec.get("bullets") if isinstance(spec.get("bullets"), list) else []
    if body: _add_text(canvas, slide, body, 0.9, 2.42, 6.75, 0.85, size=18, color=palette["muted"], font="Georgia" if serif else _font_for(palette, "Aptos"))
    for bullet_index, bullet in enumerate(bullets[:4]):
        y = 3.36 + bullet_index * 0.62
        _shape(canvas, slide, MSO_SHAPE.OVAL, 0.96, y + 0.12, 0.14, 0.14, palette["accent2"])
        _add_text(canvas, slide, str(bullet), 1.25, y, 6.55, 0.43, size=16, color=palette["fg"], font=_font_for(palette, "Aptos"))
    # A thin evidence rail is less repetitive than the old full metric card;
    # dedicated metric pages still use the larger metric treatment.
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 8.72, 2.42, 0.06, 3.18, palette["accent"])
    value = str(spec.get("metric") or "")
    if value:
        _add_text(canvas, slide, value, 9.05, 2.72, 2.75, 0.72, size=31, color=palette["accent"], bold=True, align=PP_ALIGN.LEFT, font=_font_for(palette, "Aptos Mono"))
    _add_text(canvas, slide, str(spec.get("metric_label") or "重点内容"), 9.05, 3.52, 2.75, 0.42, size=13, color=palette["muted"], font=_font_for(palette, "Aptos"))
    _add_text(canvas, slide, "事实 · 判断 · 动作", 9.05, 4.72, 2.75, 0.35, size=11, color=palette["muted"], font=_font_for(palette, "Aptos Mono"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "把事实、判断和动作放在同一条主线上。"), palette, show=show_takeaway)


def _add_content_stack(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """Full-width briefing rows for pages that would otherwise repeat cards."""
    body = str(spec.get("body") or "")
    bullets = [str(item) for item in (spec.get("bullets") or []) if str(item).strip()][:4]
    if body:
        _add_text(canvas, slide, body, 0.9, 2.24, 11.1, 0.68, size=18, color=palette["muted"], font=_font_for(palette, "Aptos"))
    if not bullets:
        bullets = [str(spec.get("summary") or spec.get("title") or "重点内容")]
    row_height = min(0.72, 2.72 / max(len(bullets), 1))
    for item_index, item in enumerate(bullets):
        top = 3.08 + item_index * (row_height + 0.1)
        secondary_fill = "EAF7F8" if palette.get("decor") == "aqua" else "F2F5FA" if palette.get("decor") == "security" else "F2F7FB"
        fill = "FFFFFF" if item_index % 2 == 0 else secondary_fill
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.9, top, 11.15, row_height, fill, line="D7E5F0")
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.9, top, 0.08, row_height, palette["accent"] if item_index == 0 else palette["accent2"])
        _add_text(canvas, slide, f"{item_index + 1:02d}", 1.16, top + 0.12, 0.58, 0.28, size=11, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, item, 1.86, top + 0.06, 9.82, row_height - 0.1, size=15, color=palette["fg"], font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "按优先级逐项展开，便于现场快速扫描。"), palette, show=show_takeaway)


def _add_content_rail(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """Asymmetric editorial rail with a narrative left and evidence right."""
    body = str(spec.get("body") or spec.get("summary") or "")
    bullets = [str(item) for item in (spec.get("bullets") or []) if str(item).strip()][:4]
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.88, 2.28, 3.55, 3.28, palette["accent"], line=palette["accent"])
    _add_text(canvas, slide, str(spec.get("metric") or "核心判断"), 1.18, 2.58, 2.86, 0.55, size=23, color="FFFFFF", bold=True, font=_font_for(palette, "Aptos"))
    _add_text(canvas, slide, body, 1.16, 3.28, 2.92, 1.78, size=17, color="FFFFFF", bold=True, font=_font_for(palette, "Aptos"))
    for item_index, item in enumerate(bullets or [str(spec.get("metric_label") or "重点内容")]):
        top = 2.4 + item_index * 0.78
        _add_text(canvas, slide, f"{item_index + 1:02d}", 4.92, top + 0.08, 0.62, 0.28, size=11, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 5.62, top + 0.08, 0.04, 0.44, palette["accent2"])
        _add_text(canvas, slide, item, 5.9, top, 5.92, 0.58, size=16, color=palette["fg"], font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "左侧给结论，右侧只保留能够支撑判断的证据。"), palette, show=show_takeaway)


def _add_content_emphasis(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    """One large conclusion plus a low horizontal evidence strip."""
    body = str(spec.get("body") or spec.get("summary") or spec.get("title") or "")
    bullets = [str(item) for item in (spec.get("bullets") or []) if str(item).strip()][:3]
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0.88, 2.38, 0.1, 1.72, palette["accent"])
    _add_text(canvas, slide, body, 1.3, 2.3, 10.45, 1.82, size=28, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
    if bullets:
        width = 10.95 / len(bullets)
        for item_index, item in enumerate(bullets):
            left = 0.88 + item_index * width
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 4.55, width - 0.18, 1.0, "FFFFFF", line="D7E5F0")
            _add_text(canvas, slide, f"{item_index + 1:02d}", left + 0.2, 4.72, 0.52, 0.24, size=10, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
            _add_text(canvas, slide, item, left + 0.82, 4.64, width - 1.18, 0.62, size=13, color=palette["fg"], font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "先给出结论，再让少量证据完成支撑。"), palette, show=show_takeaway)


def _add_split(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    columns = [(0.86, palette["accent"]), (6.82, palette["accent2"])]
    values = [spec.get("left_title") or "现状", spec.get("right_title") or "目标"]
    bullets = [spec.get("left_bullets") or spec.get("bullets") or [], spec.get("right_bullets") or ["形成统一判断", "明确下一步动作"]]
    for left, color, heading, items in zip((item[0] for item in columns), (item[1] for item in columns), values, bullets):
        if palette.get("decor") in {"state", "aviation", "aqua", "security"}:
            # Formal comparison pages read better as two flat briefing panels;
            # reserve rounded cards for lighter product/editorial themes.
            panel_fill = "FFFFFF" if color == palette["accent"] else "EAF3FA"
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 2.45, 5.45, 3.15, panel_fill, line=color)
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left, 2.45, 0.08, 3.15, color)
        else:
            _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.45, 5.45, 3.15, color, transparency=88, line=color)
        _add_text(canvas, slide, str(heading), left + 0.34, 2.77, 4.7, 0.44, size=19, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
        for item_index, item in enumerate(items[:4]):
            y = 3.42 + item_index * 0.48
            _shape(canvas, slide, MSO_SHAPE.RECTANGLE, left + 0.36, y + 0.15, 0.12, 0.12, color)
            _add_text(canvas, slide, str(item), left + 0.66, y, 4.4, 0.38, size=14, color=palette["fg"], font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "对比不是罗列，而是为了帮助下一步取舍。"), palette, show=show_takeaway)


def _add_timeline(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    steps = spec.get("steps") if isinstance(spec.get("steps"), list) else []
    if not steps: steps = [{"title": "聚焦", "body": "统一问题定义"}, {"title": "验证", "body": "换取证据"}, {"title": "规模化", "body": "固化机制"}]
    count = min(5, len(steps))
    _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 1.1, 3.47, 10.75, 0.06, palette["muted"], transparency=30)
    for index, step in enumerate(steps[:count]):
        left = 0.92 + index * (10.95 / max(1, count - 1)) if count > 1 else 5.4
        _shape(canvas, slide, MSO_SHAPE.OVAL, left, 3.19, 0.6, 0.6, palette["accent"] if index == 0 else palette["bg"], line=palette["accent"])
        _add_text(canvas, slide, f"{index + 1:02d}", left, 3.28, 0.6, 0.3, size=12, color="FFFFFF" if index == 0 else palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(step.get("title") or f"阶段 {index + 1}"), left - 0.35, 4.08, 1.3, 0.38, size=15, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(step.get("body") or "交付与复盘"), left - 0.75, 4.55, 2.1, 0.65, size=12, color=palette["muted"], align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "每一步都有明确交付物，下一阶段才有入口。"), palette, show=show_takeaway)


def _add_quote(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any], *, show_takeaway: bool = True) -> None:
    _add_text(canvas, slide, "“", 0.85, 2.2, 0.7, 1.3, size=70, color=palette["accent"], bold=True, font="Georgia")
    _add_text(canvas, slide, str(spec.get("body") or spec.get("title") or "把复杂问题说成一句能被记住的话。"), 1.55, 2.55, 9.8, 1.45, size=28, color=palette["fg"], bold=True, font="Georgia")
    _add_text(canvas, slide, str(spec.get("subtitle") or "汇报现场 / 共识记录"), 1.6, 4.5, 7.5, 0.38, size=13, color=palette["muted"], font=_font_for(palette, "Aptos Mono"))
    _add_takeaway(canvas, slide, str(spec.get("summary") or "一句清晰的话，往往比一页完整的列表更能推动行动。"), palette, show=show_takeaway)


def _add_closing(canvas: Canvas, slide: Any, spec: dict[str, Any], palette: dict[str, Any]) -> None:
    decor = str(palette.get("decor") or "")
    if decor == "aviation":
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 8.7, 0, 4.62, CANVAS_H, palette["accent"], transparency=2)
        _add_text(canvas, slide, str(spec.get("kicker") or "CLOSING"), 1.02, 1.28, 2.8, 0.3, size=11, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "谢谢"), 1.0, 2.05, 6.7, 1.2, size=43, color=palette["accent"], bold=True, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or spec.get("summary") or "请审议并指导"), 1.04, 3.72, 6.25, 0.72, size=19, color=palette["muted"], font=_font_for(palette, "Aptos"))
        return
    if decor == "aqua":
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 9.1, 0, 3.3, 0.52, palette["accent2"], transparency=48)
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 6.8, 2.5, 0.52, palette["accent2"], transparency=58)
        _add_text(canvas, slide, str(spec.get("kicker") or "THE END"), 5.1, 2.18, 3.1, 0.3, size=11, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))
        _add_text(canvas, slide, str(spec.get("title") or "谢谢"), 3.1, 2.78, 7.1, 1.0, size=41, color=palette["fg"], bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or spec.get("summary") or "同心同行，落实到下一步"), 3.0, 4.02, 7.3, 0.58, size=17, color=palette["muted"], align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        return
    if decor == "security":
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W, CANVAS_H, "0A477F")
        _shape(canvas, slide, MSO_SHAPE.PARALLELOGRAM, 0, 0, 3.7, CANVAS_H, "FFFFFF", transparency=4)
        _add_text(canvas, slide, str(spec.get("title") or "谢谢聆听"), 4.0, 2.65, 8.1, 1.0, size=40, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        _add_text(canvas, slide, str(spec.get("subtitle") or spec.get("summary") or "请批评指正"), 4.35, 3.9, 7.45, 0.58, size=17, color="DCE8F5", align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos"))
        return
    _add_text(canvas, slide, str(spec.get("kicker") or "NEXT MOVE"), 0.82, 1.0, 3.0, 0.34, size=12, color=palette["accent2"], bold=True, font=_font_for(palette, "Aptos Mono"))
    _add_text(canvas, slide, str(spec.get("title") or "从今天开始推进"), 0.8, 1.65, 10.1, 1.25, size=38, color=palette["fg"], bold=True, font=_font_for(palette, "Aptos"))
    _add_text(canvas, slide, str(spec.get("subtitle") or spec.get("summary") or "把讨论转成下一次可见的交付"), 0.85, 3.3, 8.9, 0.75, size=21, color=palette["muted"], font=_font_for(palette, "Aptos"))
    _shape(canvas, slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.84, 5.0, 3.3, 0.65, palette["accent"], line=palette["accent"])
    _add_text(canvas, slide, "下一步 / NEXT STEP", 1.05, 5.16, 2.9, 0.28, size=11, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER, font=_font_for(palette, "Aptos Mono"))


def _add_slide(
    prs: Presentation,
    spec: dict[str, Any],
    palette: dict[str, Any],
    job: PresentationJob,
    index: int,
    cover_image: Path | None,
    features: set[str],
    *,
    preserve_template_background: bool = False,
) -> None:
    # Native template fills bypass this renderer. For free-form decks, apply a
    # quiet per-slide palette variant before drawing the background and chrome.
    raw_kind = str(spec.get("kind") or "content")
    raw_layout = str(spec.get("layout_id") or "")
    palette = _slide_palette(palette, kind=raw_kind, layout=raw_layout, index=index)
    canvas = Canvas(prs)
    slide = prs.slides.add_slide(_blank_layout(prs))
    if not preserve_template_background:
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(palette["bg"])
    elif palette.get("_template_bg"):
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(str(palette["_template_bg"]))
    if "visual_decor" in features: _add_decor(canvas, slide, palette, index=index, kind=raw_kind, layout=raw_layout)
    if palette.get("decor") not in {"state", "aviation", "aqua", "security"}:
        _shape(canvas, slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.12, CANVAS_H, palette["accent"])
    render_spec = dict(spec)
    if "kicker_summary" not in features:
        render_spec["kicker"] = ""
        render_spec["summary"] = ""
        render_spec["body"] = render_spec.get("body") or ""
    kind, layout = str(render_spec.get("kind") or "content"), str(render_spec.get("layout_id") or "")
    show_takeaway = "kicker_summary" in features
    if "layout_variety" not in features and index > 1 and kind != "closing":
        kind, layout = "content", "content"
    if kind == "metric" and "metrics" not in features:
        kind, layout = "content", "content"
    if kind == "timeline" and "roadmap" not in features:
        kind, layout = "content", "content"
    if kind == "comparison" and "comparison" not in features:
        kind, layout = "content", "content"
    if index == 1 or kind == "cover" or layout == "cover":
        _add_cover(canvas, slide, render_spec, palette, job, cover_image)
    elif kind in {"section", "chapter", "section-divider"} or layout == "section-divider":
        _add_section_divider(canvas, slide, render_spec, palette)
    elif kind == "closing" or layout == "closing":
        _add_closing(canvas, slide, render_spec, palette)
    else:
        _add_page_header(canvas, slide, render_spec, palette, index, serif=palette.get("decor") in {"editorial", "data", "ink"}, kind=kind, layout=layout)
        if layout in {"agenda", "toc", "contents"} or kind in {"agenda", "toc"}: _add_agenda(canvas, slide, render_spec, palette)
        elif layout in {"metric", "metrics"} or kind == "metric": _add_metric(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"chart", "bar-chart", "line-chart", "donut-chart"} or kind == "chart": _add_chart(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"table", "data-table"} or kind == "table": _add_table(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"image", "photo-split", "image-split"} or kind in {"image", "photo"}: _add_image_split(canvas, slide, render_spec, palette, cover_image, show_takeaway=show_takeaway)
        elif layout in {"timeline", "roadmap"} or kind == "timeline": _add_timeline(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"comparison", "split"} or kind in {"comparison", "split"}: _add_split(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout == "quote" or kind == "quote": _add_quote(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"cards", "card-grid", "six-card"} or kind == "cards": _add_cards(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout in {"statement", "hero-statement"} or kind == "statement": _add_statement(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout == "content-stack": _add_content_stack(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout == "content-rail": _add_content_rail(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        elif layout == "content-emphasis": _add_content_emphasis(canvas, slide, render_spec, palette, show_takeaway=show_takeaway)
        else: _add_content(canvas, slide, render_spec, palette, serif=palette.get("decor") in {"editorial", "data", "ink"}, show_takeaway=show_takeaway)
    _add_notes(slide, job, "source_notes" in features)


def _material_blueprint_prompt(job: PresentationJob, references: str, image_manifest: str) -> str:
    """Ask for a source-bound fact map before asking for visual composition.

    High-quality deck systems separate evidence extraction from slide styling.
    Keeping this pass design-neutral makes the second pass less likely to
    invent numbers merely to satisfy a visually attractive layout.
    """
    return (
        "你是汇报材料的事实分析师。只返回 JSON，不要 Markdown。先阅读用户简报和参考资料，建立一个可供另一位演示设计师使用的事实蓝图。"
        "不得补造数字、公司名称、时间、因果关系或效果；不确定的内容放入 uncertainties。把重复段落合并，保留每个关键数字的原文短引文。"
        "蓝图结构必须是："
        '{"executive_message":"...","audience_decisions":["..."],'
        '"sections":[{"title":"...","purpose":"...","key_points":["..."],"evidence_quotes":["..."]}],'
        '"verified_metrics":[{"label":"...","value":"...","unit":"...","context":"...","source_quote":"..."}],'
        '"comparison_candidates":[{"left":"...","right":"...","evidence_quotes":["..."]}],'
        '"process_candidates":[{"step":"...","evidence_quote":"..."}],'
        '"image_evidence":[{"filename":"...","what_is_visible":"...","safe_use":"cover|evidence|ignore"}],'
        '"uncertainties":["..."],"must_preserve_user_requirements":["..."]}. '
        f"主题：{job.title}\n目的：{job.purpose}\n简报：{job.brief}\n受众：{job.audience or '团队决策者'}\n"
        f"图片资料：{image_manifest or '无'}\n参考资料：{references or '无'}"
    )


def generate_presentation(job_id: int) -> None:
    db = SessionLocal()
    job = db.get(PresentationJob, job_id)
    if job is None:
        db.close(); return

    def stage(name: str, progress: int) -> None:
        job.stage = name; job.progress = progress; db.commit()

    try:
        job.status = "running"; job.started_at = now_utc(); job.error = ""; db.commit()
        assets = db.query(PresentationJobAsset).filter(PresentationJobAsset.job_id == job.id).order_by(PresentationJobAsset.id.asc()).all()
        image_assets = [asset for asset in assets if asset.role == "reference" and (asset.content_type or "").startswith("image/")]
        text_parts = [f"【{asset.filename}】\n{asset.text_content or ''}" for asset in assets if asset.role == "reference" and asset.text_content]
        reference_text = _clip_multiline("\n\n".join(text_parts), MAX_REFERENCE_TEXT)
        vision_attachments = _presentation_vision_attachments(image_assets)
        image_manifest = "\n".join(
            f"【图片资料：{asset.filename}】"
            for asset in image_assets
            if str(asset.filename or "").strip()
        )
        template_asset = next((asset for asset in assets if asset.role == "template"), None)
        template_path = Path(template_asset.file_path) if template_asset else None
        user_direction = "\n".join((job.title or "", job.purpose or "", job.brief or ""))
        template_disabled = _prompt_disables_native_template(user_direction)
        template_library = _analyse_native_template(template_path) if template_path and not template_disabled else None
        if template_path and not template_disabled and not template_library:
            raise RuntimeError(
                "上传的 PPT 模板无法解析，已停止生成以避免静默改用其他样式。"
                "请确认模板可在 PowerPoint 中正常打开后重新上传。"
            )
        template_context = _template_context(template_library)
        resolved_style, prompt_style, style_source = _resolve_style(job, reference_text, has_images=bool(image_assets), has_template=bool(template_asset))
        # A raw PPTX template is a visual shell, not a source of new evidence.
        # Embedded photos are therefore isolated by default; users can opt in
        # explicitly when they really want to carry those media objects over.
        allow_template_media = _prompt_allows_template_media(user_direction)
        # Uploading a PPTX is an explicit visual choice. Keep its native pages
        # even when the brief contains ordinary style adjectives; abandon it
        # only when the user clearly says the template is not the visual base.
        use_native_template = bool(
            template_path
            and template_library
            and not template_disabled
        )
        mode, features = _job_options(job)
        stage("reading_references", 18)
        service: OpenAIService | None = None
        blueprint: dict[str, Any] | None = None
        # Pass 1 is intentionally design-neutral. It gives the composition
        # pass a compact, source-bound map of facts, sections, metrics and
        # image evidence instead of asking one completion to do everything at
        # once.
        try:
            service = OpenAIService(provider=job.provider or "openai", text_model=job.model or None)
            blueprint_result = service.chat(
                _material_blueprint_prompt(job, reference_text, image_manifest),
                attachments=vision_attachments or None,
            )
            record_token_usage(
                db,
                user_id=job.user_id,
                source="presentation-blueprint",
                provider=job.provider or "openai",
                model=str(blueprint_result.get("model") or service.text_model),
                prompt_tokens=int(blueprint_result.get("prompt_tokens") or 0),
                completion_tokens=int(blueprint_result.get("completion_tokens") or 0),
                total_tokens=int(blueprint_result.get("total_tokens") or 0),
            )
            blueprint = _safe_json(str(blueprint_result.get("text") or ""))
        except Exception as exc:
            logger.warning("Presentation evidence pass fallback for job %s: %s", job.id, exc)
        blueprint_context = _clip_multiline(
            json.dumps(blueprint, ensure_ascii=False, indent=2) if blueprint else "",
            MAX_BLUEPRINT_CONTEXT,
        )
        prompt = (
            "你是资深国企汇报设计师、信息架构师和数据编辑。只返回 JSON，不要 Markdown。先在内部完成资料事实提取、去重、总结、风险核验和页面分配，再输出结构。结构必须是 "
            '{"title":"...","slides":[{"kind":"cover|agenda|section|statement|content|cards|metric|timeline|comparison|quote|chart|table|image|closing",'
            '"layout_id":"cover|agenda|section-divider|statement|content|content-stack|content-rail|content-emphasis|cards|metric|timeline|comparison|chart|bar-chart|line-chart|donut-chart|table|photo-split|closing",'
            '"visual_intent":"hero|evidence|comparison|process|metrics|photo|table|chapter|closing","density":"low|medium|high",'
            '"template_slide":null,"kicker":"...","title":"...","subtitle":"...","summary":"一句页面结论",'
            '"body":"...","bullets":["..."],"metric":"...","metric_label":"...",'
            '"left_title":"...","left_bullets":["..."],"right_title":"...",'
            '"right_bullets":["..."],"steps":[{"title":"...","body":"..."}],'
            '"image_slot":"...","chart":{"type":"bar|line|donut","labels":["..."],"values":[1,2],"unit":"..."},'
            '"table":{"columns":["..."],"rows":[["..."]]} }]}. '
            f"请为{job.audience or '团队决策者'}制作一份正好{job.slide_count}页、语言为{job.language}的演示文稿。"
            f"主题：{job.title}\n目的：{job.purpose}\n简报：{job.brief}\n"
            f"叙事模式：{mode}（{MODE_GUIDANCE.get(mode, MODE_GUIDANCE['pyramid'])}）\n视觉风格候选：{job.style}；当前解析风格：{resolved_style}（来源：{style_source}），画幅：{job.aspect_ratio}\n"
            f"启用能力：{', '.join(FEATURE_GUIDANCE[item] for item in sorted(features) if item in FEATURE_GUIDANCE)}。"
            f"视觉素材位：{'启用' if job.include_images else '关闭；不得输出 image/photo 页面，也不得嵌入参考图片'}。"
            "用户简报中的内容、结构和口径要求优先级最高。上传模板本身代表明确的视觉选择，除非用户明确要求忽略或不用模板，否则模板的背景、字体、主题色、母版结构和页面装饰必须沿用，不得另画一套统一背景。"
            "只使用参考资料中出现的事实和数字，不得补造数据；合并重复观点，避免整段复制原文。每页只表达一个主结论，标题写成判断句，并给出一句可直接汇报的 summary。先用事实蓝图校验口径，再写页面内容。"
            "相邻页面必须轮换版式：封面、目录、章节分隔、结论页、双栏/卡片、指标、时间轴、表格/图表、图片证据和收束页交替出现；普通正文也要在 content、content-stack、content-rail、content-emphasis 之间轮换；不要把所有页面都做成同样的背景、圆角卡片或右侧指标框。"
            "若采用国企汇报方向，变化应来自深蓝章节页、白底正文页、浅色数据页、图片证据页和不同信息骨架，而不是霓虹、毛玻璃、大面积渐变或密集的拟态控件。编号章、章节条、证据标签、图片说明和KPI标记只在对应页面角色中使用，不得每页复制同一组控件。10页以上应在封面后安排 agenda 目录页；超过12页时至少安排2个章节分隔页，且相邻普通正文页不得使用相同 layout_id。"
            "机场岗位外包类材料优先采用点位/图纸证据、人员与成本表、市场调研和可行性判断的页面节奏；年度规划类材料优先采用工作主线、重点举措、季度里程碑和成果目标；安护年度总结类材料优先采用业务数据、举措证据、能力建设、服务成效、问题与计划。"
            "只有在资料存在至少两组可比较数字时才输出 chart 或 table；chart.values 必须逐一来自资料，无法核验时宁可不用图表。图表类型按数据语义选择：时间序列用 line，构成占比用 donut，类别比较用 bar。"
            "visual_intent 描述页面真正的沟通任务，density 描述信息密度；不要为了凑版式把普通正文伪装成指标页。若上传了模板，只有在某页与内容的视觉骨架明显唯一匹配时才填写 template_slide，填写模板摘要中的纯整数页码（例如 24，不要写 P24）；其余页面填 null，由系统按版式能力和相邻去重选择。相邻页面尽量不要重复 template_slide 或 family。"
            "模板页面结构只作为可填充的原生页面库，输出时由系统克隆合适页面并替换槽位，不能改变模板的主色、字体、图片、Logo、表格和页眉页脚。"
            "模板页面库摘要中的示例文字和数字只能用于判断页面类型与槽位容量，不能作为新汇报的数据来源；需要复用模板里的业务内容时，用户应同时把该 PPT 作为参考资料上传。"
            f"事实蓝图：{blueprint_context or '上一阶段未返回蓝图，请直接从参考资料提取。'}\n"
            f"参考资料核验摘录：{_clip_multiline(reference_text, 14000) or '无'}\n"
            f"图片资料清单：{image_manifest or '无'}\n"
            "如果收到图片附件，请先判断其中的场景、流程、表格或可读文字，再决定它适合做证据页、封面还是不使用；不要臆造图片中无法确认的数字。\n"
            f"模板页面库摘要：{template_context or '未上传原生模板'}"
        )
        stage("structuring_story", 40)
        spec: dict[str, Any] | None = None
        try:
            if service is None:
                service = OpenAIService(provider=job.provider or "openai", text_model=job.model or None)
            result = service.chat(prompt, attachments=None if blueprint else (vision_attachments or None))
            record_token_usage(db, user_id=job.user_id, source="presentation", provider=job.provider or "openai", model=str(result.get("model") or service.text_model), prompt_tokens=int(result.get("prompt_tokens") or 0), completion_tokens=int(result.get("completion_tokens") or 0), total_tokens=int(result.get("total_tokens") or 0))
            spec = _safe_json(str(result.get("text") or ""))
        except Exception as exc:
            logger.warning("Presentation model planning fallback for job %s: %s", job.id, exc)
        normalised = _normalise_spec(spec, job, reference_text)
        normalised = _guard_model_data_slides(
            normalised,
            "\n".join((job.title or "", job.brief or "", reference_text)),
        )
        normalised, structure_repairs = _ensure_report_structure(
            normalised,
            max(5, min(100, int(job.slide_count or 10))),
            native_agenda=bool(
                not use_native_template
                or (
                    template_library
                    and any(
                        _native_layout_family(item) == "agenda"
                        for item in (template_library.get("slides") or [])
                        if isinstance(item, dict)
                    )
                )
            ),
        )
        normalised = _enforce_layout_variety(normalised, features)
        if not job.include_images:
            for slide in normalised.get("slides", []):
                if not isinstance(slide, dict):
                    continue
                kind = str(slide.get("kind") or "").casefold()
                layout = str(slide.get("layout_id") or "").casefold()
                if kind in {"image", "photo"} or layout in {"image", "photo-split", "image-split"}:
                    slide["kind"] = "content"
                    slide["layout_id"] = "content"
                    slide["image_slot"] = ""
        inferred_chart = _infer_chart_from_reference(reference_text)
        if inferred_chart and not any(str(item.get("kind") or "") == "chart" for item in normalised.get("slides", []) if isinstance(item, dict)):
            chart_slide = {
                "kind": "chart",
                "layout_id": "bar-chart",
                "kicker": "DATA / EXTRACTED",
                "title": "资料中的量化变化值得单独看",
                "summary": "图表直接取自参考资料中的标签和数字，未补造缺失数据。",
                "chart": inferred_chart,
            }
            normalised.setdefault("slides", []).insert(max(1, len(normalised.get("slides", [])) - 1), chart_slide)
            target_count = max(5, min(100, int(job.slide_count or 10)))
            if len(normalised["slides"]) > target_count:
                normalised["slides"] = normalised["slides"][: max(1, target_count - 1)] + [normalised["slides"][-1]]
        stage("composing_slides", 58)
        template_palette: dict[str, str] = {}
        if use_native_template and template_path and template_path.exists():
            # Palette extraction is metadata only. The actual output is built
            # by cloning native source slides below, so masters and media stay
            # intact instead of being replaced by a blank layout.
            template_prs = Presentation(str(template_path))
            template_palette = _extract_template_palette(template_prs)
            prs = None
        else:
            prs = Presentation()
            if job.aspect_ratio == "4:3": prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
            else: prs.slide_width, prs.slide_height = Inches(CANVAS_W), Inches(CANVAS_H)
        palette = dict(STYLE_PRESETS.get(resolved_style, STYLE_PRESETS["dark-tech"]))
        if template_palette:
            # A declared prompt style controls the palette; otherwise the uploaded template is the visual source of truth.
            if use_native_template:
                palette.update(template_palette)
            elif template_palette.get("font"):
                palette["font"] = template_palette["font"]
            if template_palette.get("bg"):
                palette["_template_bg"] = template_palette["bg"]
        visual_image_assets = image_assets if job.include_images else []
        cover_asset = visual_image_assets[0] if visual_image_assets else None
        cover_image = Path(cover_asset.file_path) if cover_asset else None
        slides = normalised.get("slides") or []
        if cover_image and not use_native_template and len(slides) > 1 and not any(str(item.get("kind") or "") in {"image", "photo"} for item in slides if isinstance(item, dict)):
            slides[1] = {**slides[1], "kind": "image", "layout_id": "photo-split", "image_slot": "参考图片"}
        if not slides: raise RuntimeError("无法构建任何幻灯片")
        output_dir = OUTPUT_ROOT / str(job.user_id); output_dir.mkdir(parents=True, exist_ok=True)
        output_filename = re.sub(r"[^A-Za-z0-9\u4e00-\u9fff._-]+", "-", job.title.strip() or "aiweb-presentation")[:90].strip("-") + ".pptx"
        output_path = output_dir / f"{job.id}-{output_filename}"
        native_plan: dict[str, Any] | None = None
        native_chart_overlay_count = 0
        if use_native_template and template_path:
            native_plan = _build_native_fill_plan(
                template_library,
                [item for item in slides if isinstance(item, dict)],
                allow_template_media=allow_template_media,
            )
            if not native_plan.get("slides"):
                raise RuntimeError("模板没有可填充的原生页面")
            output_path.unlink(missing_ok=True)
            try:
                _apply_native_fill(template_path, native_plan, output_path)
            except Exception as exc:
                raise RuntimeError(
                    "模板原生填充失败，已停止生成以避免保留模板中的旧表格或旧图表数据。"
                ) from exc
            native_chart_overlay_count = _apply_native_chart_overlays(
                output_path,
                native_plan,
                palette,
                canvas_px=(template_library or {}).get("canvas_px") if template_library else None,
            )
            _remove_native_empty_shapes(output_path, native_plan)
            output_slide_count = len(native_plan.get("slides") or [])
        else:
            if prs is None:
                raise RuntimeError("无法初始化 PPT 画布")
            image_use_counts: Counter[int] = Counter()
            for index, slide_spec in enumerate(slides, start=1):
                image_path = _pick_reference_image(slide_spec, visual_image_assets, image_use_counts) if index > 1 else cover_image
                _add_slide(prs, slide_spec, palette, job, index, image_path, features)
            prs.save(str(output_path))
            output_slide_count = len(slides)
        stage("quality_check", 86)
        check = Presentation(str(output_path))
        if len(check.slides) != output_slide_count: raise RuntimeError("生成文件的幻灯片数量校验失败")
        try: metadata = json.loads(job.metadata_json or "{}")
        except (TypeError, ValueError): metadata = {}
        metadata.update({
            "slide_count": output_slide_count,
            "reference_count": len(assets),
            "template_uploaded": bool(template_asset),
            "template_used": bool(native_plan),
            "native_template_fill": bool(native_plan),
            "native_template_requested": bool(template_path and not template_disabled),
            "template_media_reused": allow_template_media,
            "native_chart_overlays": native_chart_overlay_count,
            "template_source_slides": [item.get("source_slide") for item in (native_plan or {}).get("slides", []) if isinstance(item, dict)],
            "mode": mode,
            "features": sorted(features),
            "requested_style": job.style,
            "resolved_style": resolved_style,
            "prompt_style": prompt_style,
            "style_source": style_source,
            "template_palette": template_palette,
            "blueprint_used": bool(blueprint),
            "structure_repairs": structure_repairs,
        })
        job.metadata_json = json.dumps(metadata, ensure_ascii=False); job.output_path = str(output_path); job.output_filename = output_filename
        job.status = "completed"; job.stage = "ready_to_download"; job.progress = 100; job.completed_at = now_utc(); db.commit()
    except Exception as exc:
        logger.exception("Presentation job %s failed", job_id); db.rollback(); job = db.get(PresentationJob, job_id)
        if job is not None:
            job.status = "failed"; job.stage = "failed"; job.progress = 0; job.error = str(exc)[:500] or "PPT 生成失败，请稍后重试。"; job.completed_at = now_utc(); db.commit()
    finally:
        db.close()
