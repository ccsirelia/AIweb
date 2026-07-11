"""Extract readable text from common office / document uploads."""

from __future__ import annotations

import csv
import io
import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree

logger = logging.getLogger(__name__)

# Keep extracted text bounded so chat prompts stay manageable.
DEFAULT_MAX_CHARS = 120_000

OFFICE_EXTENSIONS = {
    ".docx",
    ".doc",
    ".pdf",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".xls",
    ".csv",
    ".txt",
    ".md",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".yaml",
    ".yml",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".css",
}

LEGACY_BINARY_EXTENSIONS = {".doc", ".ppt", ".xls"}


def _clip(text: str, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    cleaned = "\n".join(line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"))
    cleaned = "\n".join(line for line in cleaned.split("\n") if line.strip() or line == "")
    cleaned = cleaned.strip()
    if len(cleaned) <= max_chars:
        return cleaned
    omitted = len(cleaned) - max_chars
    return (
        cleaned[:max_chars].rstrip()
        + f"\n\n[... truncated {omitted} characters to fit model context ...]"
    )


def _decode_text_bytes(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str | None:
    sample = data[: max(max_chars * 4, 200_000)]
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "utf-16", "latin-1"):
        try:
            text = sample.decode(encoding)
            text = _clip(text, max_chars)
            return text or None
        except UnicodeDecodeError:
            continue
    text = _clip(sample.decode("utf-8", errors="ignore"), max_chars)
    return text or None


def _extract_docx(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    from docx import Document

    document = Document(io.BytesIO(data))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            parts.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            cells = [" ".join((cell.text or "").split()) for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Table {table_index}]")
            parts.extend(rows)

    # Headers and footers are not included in document.paragraphs, but often
    # contain document titles, dates, owners, or other important context.
    seen_story_parts: set[int] = set()
    for section_index, section in enumerate(document.sections, start=1):
        for story_name, story in (("Header", section.header), ("Footer", section.footer)):
            story_key = id(story.part)
            if story_key in seen_story_parts:
                continue
            seen_story_parts.add(story_key)

            story_lines = [
                (paragraph.text or "").strip()
                for paragraph in story.paragraphs
                if (paragraph.text or "").strip()
            ]
            for table in story.tables:
                for row in table.rows:
                    cells = [" ".join((cell.text or "").split()) for cell in row.cells]
                    if any(cells):
                        story_lines.append(" | ".join(cells))
            if story_lines:
                parts.append(f"[{story_name} {section_index}]")
                parts.extend(story_lines)

    # Fallback: some docs store most text in headers/footers or custom XML.
    if not parts:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                for name in archive.namelist():
                    if not name.startswith("word/") or not name.endswith(".xml"):
                        continue
                    xml_bytes = archive.read(name)
                    root = ElementTree.fromstring(xml_bytes)
                    texts = [
                        node.text.strip()
                        for node in root.iter()
                        if node.text and node.text.strip()
                    ]
                    if texts:
                        parts.append(f"[{name}]")
                        parts.extend(texts)
        except Exception as exc:  # noqa: BLE001
            logger.warning("docx zip fallback failed: %s", exc)

    return _clip("\n".join(parts), max_chars)


def _extract_pdf(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            logger.warning("pdf page %s extract failed: %s", index, exc)
            text = ""
        text = text.strip()
        if text:
            parts.append(f"[Page {index}]\n{text}")
        if sum(len(item) for item in parts) >= max_chars:
            break
    return _clip("\n\n".join(parts), max_chars)


def _extract_xlsx(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell).strip() for cell in row]
                if not any(values):
                    continue
                parts.append(" | ".join(values))
                row_count += 1
                if row_count >= 500:
                    parts.append("[... sheet rows truncated ...]")
                    break
                if sum(len(item) for item in parts) >= max_chars:
                    break
            if sum(len(item) for item in parts) >= max_chars:
                break
    finally:
        workbook.close()
    return _clip("\n".join(parts), max_chars)


def _extract_pptx(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if text:
                    slide_lines.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    cells = [" ".join((cell.text or "").split()) for cell in row.cells]
                    if any(cells):
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            parts.append(f"[Slide {index}]")
            parts.extend(slide_lines)
        if sum(len(item) for item in parts) >= max_chars:
            break
    return _clip("\n".join(parts), max_chars)


def _extract_csv(data: bytes, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    text = _decode_text_bytes(data, max_chars=max_chars * 2)
    if not text:
        return ""
    try:
        reader = csv.reader(io.StringIO(text))
        rows = []
        for index, row in enumerate(reader, start=1):
            rows.append(" | ".join(cell.strip() for cell in row))
            if index >= 1000 or sum(len(item) for item in rows) >= max_chars:
                if index >= 1000:
                    rows.append("[... csv rows truncated ...]")
                break
        return _clip("\n".join(rows), max_chars)
    except Exception:
        return _clip(text, max_chars)


def extract_document_text(
    filename: str,
    content_type: str = "",
    data: bytes | None = None,
    *,
    max_chars: int = DEFAULT_MAX_CHARS,
) -> str | None:
    """Return extracted text for supported documents, or None if unavailable."""
    if not data:
        return None

    ext = Path(filename or "").suffix.lower()
    content_type = (content_type or "").lower().strip()

    try:
        if ext == ".docx" or "wordprocessingml" in content_type:
            text = _extract_docx(data, max_chars=max_chars)
        elif ext == ".pdf" or content_type == "application/pdf":
            text = _extract_pdf(data, max_chars=max_chars)
        elif ext in {".xlsx", ".xlsm"} or "spreadsheetml" in content_type:
            text = _extract_xlsx(data, max_chars=max_chars)
        elif ext == ".pptx" or "presentationml" in content_type:
            text = _extract_pptx(data, max_chars=max_chars)
        elif ext == ".csv" or content_type in {"text/csv", "application/csv"}:
            text = _extract_csv(data, max_chars=max_chars)
        elif ext in LEGACY_BINARY_EXTENSIONS:
            return (
                f"[Unsupported legacy binary format: {filename}. "
                "Please convert to .docx / .pptx / .xlsx and upload again.]"
            )
        elif ext in OFFICE_EXTENSIONS or content_type.startswith("text/") or "json" in content_type or "xml" in content_type:
            # Plain text family and source code.
            text = _decode_text_bytes(data, max_chars=max_chars) or ""
        else:
            # Last-resort: attempt UTF decode for unknown but text-like uploads.
            text = _decode_text_bytes(data, max_chars=max_chars) or ""
    except Exception as exc:  # noqa: BLE001
        logger.warning("Failed to extract text from %s: %s", filename, exc)
        return None

    text = (text or "").strip()
    if not text:
        return None
    return text
