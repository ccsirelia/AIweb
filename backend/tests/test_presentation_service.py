from __future__ import annotations

import json
import unittest
from collections import Counter
from pathlib import Path
from tempfile import TemporaryDirectory
from types import SimpleNamespace

from pptx import Presentation

from services.presentation_service import (
    _add_slide,
    _build_native_fill_plan,
    _chart_is_supported_by_source,
    _choose_intelligent_style,
    _clip_multiline,
    _color_contrast,
    _contrasting_color,
    _ensure_report_structure,
    _enforce_layout_variety,
    _fallback_spec,
    _guard_model_data_slides,
    _infer_chart_from_reference,
    _job_options,
    _native_layout_family,
    _native_layout_affordance,
    _native_layout_unit_capacity,
    _native_large_picture_count,
    _native_is_chrome,
    _native_chapter_title,
    _native_closing_text,
    _native_compact_chapter_title,
    _native_hidden_empty_shape_ids,
    _native_title_candidate,
    _native_page_role,
    _native_pick_slide,
    _native_row_pairs,
    _native_chart_bounds,
    _native_chart_edits,
    _native_replacements,
    _native_sequence_sources,
    _native_slot_is_body_like,
    _native_slot_is_empty_decor,
    _native_table_edits,
    _template_context,
    _presentation_vision_attachments,
    _prompt_disables_native_template,
    _prompt_allows_template_media,
    _slide_palette,
    _style_from_text,
    _table_is_supported_by_source,
    STYLE_PRESETS,
)


def slot(role: str, text: str, x: float, width: float = 320, height: float = 120) -> dict[str, object]:
    return {
        "slot_id": f"{role}-{x}-{text}",
        "role": role,
        "text": text,
        "geometry": {"x": x, "y": 180, "width": width, "height": height},
        "text_metrics": {"font_size_px": 24},
    }


class PresentationStylePrecedenceTests(unittest.TestCase):
    def test_explicit_empty_feature_set_stays_empty(self) -> None:
        mode, features = _job_options(SimpleNamespace(metadata_json='{"mode":"briefing","features":[]}'))
        self.assertEqual(mode, "briefing")
        self.assertEqual(features, set())

    def test_business_words_do_not_implicitly_override_uploaded_template(self) -> None:
        self.assertIsNone(_style_from_text("航站楼国企项目汇报，内容需要沉稳严谨"))

    def test_explicit_visual_direction_is_detected(self) -> None:
        self.assertEqual(_style_from_text("请采用国企蓝白风格"), "state-briefing")
        self.assertEqual(_style_from_text("做成暗夜科技"), "dark-tech")
        self.assertEqual(_style_from_text("请采用机场蓝白样式"), "aviation-blue")
        self.assertEqual(_style_from_text("视觉使用浅青规划风格"), "aqua-planning")
        self.assertEqual(_style_from_text("设计成安护深蓝"), "security-report")

    def test_intelligent_style_uses_report_semantics(self) -> None:
        self.assertEqual(_choose_intelligent_style("航站楼岗位外包项目", has_images=False, has_template=False), "aviation-blue")
        self.assertEqual(_choose_intelligent_style("2027年度工作计划与部署", has_images=False, has_template=False), "aqua-planning")
        self.assertEqual(_choose_intelligent_style("安检护卫年度总结", has_images=False, has_template=False), "security-report")

    def test_state_palette_changes_by_page_role(self) -> None:
        base = STYLE_PRESETS["state-briefing"]
        self.assertEqual(_slide_palette(base, kind="section", index=2)["bg"], "0B3D70")
        self.assertEqual(_slide_palette(base, kind="chart", index=3)["bg"], "EEF6FC")
        self.assertEqual(_slide_palette(base, kind="content", index=4)["bg"], "F3F7FB")

    def test_reference_derived_palettes_change_by_page_role(self) -> None:
        self.assertEqual(_slide_palette(STYLE_PRESETS["aviation-blue"], kind="section", index=2)["bg"], "0A477F")
        self.assertEqual(_slide_palette(STYLE_PRESETS["aqua-planning"], kind="chart", index=3)["bg"], "E8F7F8")
        self.assertEqual(_slide_palette(STYLE_PRESETS["security-report"], kind="chart", index=3)["bg"], "EEF3FB")

    def test_overlay_text_color_is_corrected_for_local_background(self) -> None:
        corrected = _contrasting_color("F8FAFC", "FFFFFF", minimum=4.5)
        self.assertEqual(corrected, "111827")
        self.assertGreaterEqual(_color_contrast(corrected, "FFFFFF"), 4.5)

    def test_state_renderer_outputs_distinct_page_roles(self) -> None:
        prs = Presentation()
        while prs.slides:
            slide_id = prs.slides._sldIdLst[0]
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)
        job = SimpleNamespace(
            title="专项汇报",
            audience="管理层",
            purpose="推动决策",
            include_images=False,
            metadata_json="{}",
        )
        features = {"kicker_summary", "layout_variety", "visual_decor", "data_story"}
        specs = [
            {"kind": "cover", "layout_id": "cover", "title": "专项汇报"},
            {"kind": "section", "layout_id": "section-divider", "title": "项目背景"},
            {
                "kind": "chart",
                "layout_id": "bar-chart",
                "title": "关键指标形成对比",
                "chart": {"type": "bar", "labels": ["甲", "乙"], "values": [12, 20]},
            },
            {
                "kind": "content",
                "layout_id": "content-stack",
                "title": "三项工作依次推进",
                "bullets": ["明确边界", "完成交接", "持续复盘"],
            },
        ]
        for index, spec in enumerate(specs, start=1):
            _add_slide(prs, spec, STYLE_PRESETS["state-briefing"], job, index, None, features)

        backgrounds = [str(slide.background.fill.fore_color.rgb) for slide in prs.slides]
        self.assertEqual(backgrounds, ["FFFFFF", "0B3D70", "EEF6FC", "FAFCFE"])
        self.assertGreater(len(prs.slides[2].shapes), len(prs.slides[1].shapes))
        self.assertNotEqual(len(prs.slides[0].shapes), len(prs.slides[3].shapes))

    def test_reference_derived_renderers_support_agenda_and_role_variants(self) -> None:
        job = SimpleNamespace(
            title="年度工作汇报",
            audience="管理层",
            purpose="复盘并部署",
            include_images=False,
            metadata_json="{}",
        )
        features = {"kicker_summary", "layout_variety", "visual_decor", "data_story"}
        specs = [
            {"kind": "cover", "layout_id": "cover", "title": "年度工作汇报"},
            {"kind": "agenda", "layout_id": "agenda", "title": "汇报目录", "bullets": ["总体情况", "重点工作", "问题风险", "下一步"]},
            {"kind": "section", "layout_id": "section-divider", "kicker": "01", "title": "重点工作"},
            {"kind": "chart", "layout_id": "bar-chart", "title": "关键指标", "chart": {"type": "bar", "labels": ["甲", "乙"], "values": [12, 20]}},
            {"kind": "closing", "layout_id": "closing", "title": "谢谢"},
        ]
        for style_id in ("aviation-blue", "aqua-planning", "security-report"):
            with self.subTest(style=style_id):
                prs = Presentation()
                while prs.slides:
                    slide_id = prs.slides._sldIdLst[0]
                    prs.part.drop_rel(slide_id.rId)
                    prs.slides._sldIdLst.remove(slide_id)
                for index, spec in enumerate(specs, start=1):
                    _add_slide(prs, spec, STYLE_PRESETS[style_id], job, index, None, features)
                self.assertEqual(len(prs.slides), 5)
                self.assertGreater(len(prs.slides[1].shapes), 8)
                self.assertNotEqual(
                    str(prs.slides[1].background.fill.fore_color.rgb),
                    str(prs.slides[2].background.fill.fore_color.rgb),
                )
                chart_text = "\n".join(
                    shape.text
                    for shape in prs.slides[3].shapes
                    if getattr(shape, "has_text_frame", False)
                )
                expected_control = {
                    "aviation-blue": "专项汇报 / AVIATION BRIEF",
                    "aqua-planning": "年度规划 / ANNUAL PLAN",
                    "security-report": "数据要点",
                }[style_id]
                self.assertIn(expected_control, chart_text)

    def test_template_is_abandoned_only_by_explicit_instruction(self) -> None:
        self.assertFalse(_prompt_disables_native_template("请采用沉稳国企风格并参考上传模板"))
        self.assertTrue(_prompt_disables_native_template("模板仅作内容参考，请重新设计"))


class PresentationDataTests(unittest.TestCase):
    def test_vision_context_is_bounded_and_keeps_data_urls(self) -> None:
        with TemporaryDirectory() as directory:
            path = Path(directory) / "reference.png"
            path.write_bytes(b"png-bytes")
            asset = SimpleNamespace(
                file_path=str(path),
                filename="reference.png",
                content_type="image/png",
            )
            attachments = _presentation_vision_attachments([asset])
        self.assertEqual(len(attachments), 1)
        self.assertTrue(attachments[0]["data_url"].startswith("data:image/png;base64,"))

    def test_multiline_clipping_preserves_table_rows(self) -> None:
        self.assertEqual(_clip_multiline("A | 1\nB | 2", 100), "A | 1\nB | 2")

    def test_chart_inference_keeps_one_coherent_unit(self) -> None:
        source = "\n".join((
            "[Sheet: 经营]",
            "指标 | 2025",
            "营业收入 | 120万元",
            "利润 | 18万元",
            "完成率 | 80%",
            "满意度 | 92%",
        ))
        chart = _infer_chart_from_reference(source)
        self.assertIsNotNone(chart)
        self.assertEqual(chart["labels"], ["营业收入", "利润"])
        self.assertEqual(chart["values"], [120.0, 18.0])
        self.assertEqual(chart["unit"], "万元")

    def test_model_chart_values_must_exist_in_source(self) -> None:
        chart = {"labels": ["收入", "利润"], "values": [120, 18]}
        self.assertTrue(_chart_is_supported_by_source(chart, "营业收入120万元，利润18万元"))
        self.assertFalse(_chart_is_supported_by_source(chart, "营业收入120万元，利润20万元"))

    def test_model_chart_keeps_label_value_pairing(self) -> None:
        chart = {"labels": ["收入", "利润"], "values": [18, 120]}
        self.assertFalse(_chart_is_supported_by_source(chart, "收入120万元\n利润18万元"))

    def test_model_table_cells_must_exist_in_source(self) -> None:
        table = {"columns": ["岗位", "人数"], "rows": [["护卫岗", "100"], ["管理岗", "4"]]}
        source = "岗位 | 人数\n护卫岗 | 100\n管理岗 | 4"
        self.assertTrue(_table_is_supported_by_source(table, source))
        table["rows"][1][1] = "8"
        self.assertFalse(_table_is_supported_by_source(table, source))

    def test_model_table_headers_must_be_traceable_or_structural(self) -> None:
        table = {"columns": ["岗位", "人数", "来源系统"], "rows": [["护卫岗", "100", "台账"]]}
        self.assertFalse(_table_is_supported_by_source(table, "岗位 | 人数\n护卫岗 | 100\n台账"))
        table["columns"][-1] = "说明"
        self.assertTrue(_table_is_supported_by_source(table, "岗位 | 人数\n护卫岗 | 100\n台账"))

    def test_unverified_model_table_is_downgraded(self) -> None:
        spec = {"slides": [{"kind": "table", "layout_id": "table", "title": "测算", "table": {"columns": ["岗位", "人数"], "rows": [["护卫岗", "999"]]}}]}
        guarded = _guard_model_data_slides(spec, "护卫岗 | 100")
        self.assertEqual(guarded["slides"][0]["kind"], "content")
        self.assertIsNone(guarded["slides"][0]["table"])

    def test_fallback_uses_uploaded_facts_instead_of_generic_copy(self) -> None:
        job = SimpleNamespace(
            title="经营汇报",
            brief="总结经营情况",
            audience="管理层",
            purpose="决策",
            slide_count=6,
        )
        result = _fallback_spec(job, "营业收入 | 120万元\n利润 | 18万元\n主要风险为供应链延迟")
        visible = str(result["slides"])
        self.assertIn("营业收入", visible)
        self.assertIn("供应链延迟", visible)


class NativeTemplateLayoutTests(unittest.TestCase):
    def test_template_context_exposes_structure_without_example_copy(self) -> None:
        context = _template_context({
            "slides": [{
                "slide_index": 7,
                "page_type": "content_candidate",
                "text_summary": "旧项目成本 1094 万元/年，不应成为新汇报事实",
                "slots": [slot("body_candidate", "旧项目成本 1094 万元/年", 100, 900, 300)],
                "tables": [], "charts": [], "diagrams": [], "pictures": [],
            }],
        })
        self.assertIn("family=hero-body", context)
        self.assertIn("affordance=单一结论或图文说明", context)
        self.assertNotIn("1094", context)
        self.assertNotIn("旧项目成本", context)

    def test_table_shell_wins_over_misclassified_chapter_label(self) -> None:
        item = {
            "page_type": "chapter_candidate",
            "slots": [],
            "tables": [{"row_count": 24, "column_count": 11}],
        }
        self.assertTrue(_native_layout_family(item).startswith("table-"))
        self.assertEqual(_native_page_role(item), "data")

    def test_generic_content_does_not_select_template_photo(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": [{"geometry": {"x": 0, "y": 0, "width": 900, "height": 500}}]},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": []},
            ],
        }
        selected = _native_pick_slide(library, {"kind": "content", "layout_id": "content"}, Counter(), None)
        self.assertEqual(selected["slide_index"], 2)
        selected_image = _native_pick_slide(library, {"kind": "image", "layout_id": "photo-split", "visual_intent": "photo", "_allow_template_media": True}, Counter(), None)
        self.assertEqual(selected_image["slide_index"], 1)

    def test_template_photo_requires_explicit_opt_in(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": [{"geometry": {"x": 0, "y": 0, "width": 900, "height": 500}}]},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": []},
            ],
        }
        selected = _native_pick_slide(library, {"kind": "image", "layout_id": "photo-split", "visual_intent": "photo"}, Counter(), None)
        self.assertEqual(selected["slide_index"], 2)
        self.assertFalse(_prompt_allows_template_media("请生成图片证据页"))
        self.assertTrue(_prompt_allows_template_media("请沿用模板图片"))

    def test_incompatible_explicit_template_page_is_ignored(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": [{"geometry": {"x": 0, "y": 0, "width": 900, "height": 500}}]},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("label_candidate", "卡片", 100, 180, 80)] * 6, "pictures": []},
            ],
        }
        selected = _native_pick_slide(library, {"kind": "image", "layout_id": "photo-split", "visual_intent": "photo", "template_slide": "P2", "_allow_template_media": True}, Counter(), None)
        self.assertEqual(selected["slide_index"], 1)

    def test_chart_bounds_stay_inside_source_body_slot(self) -> None:
        bounds = _native_chart_bounds(
            {"slots": [{"role": "body_candidate", "geometry": {"x": 160, "y": 180, "width": 760, "height": 360}}]},
            {"width": 1280, "height": 720},
        )
        left, top, width, height = bounds
        self.assertGreaterEqual(left, 160 / 1280 * 13.333)
        self.assertGreaterEqual(top, 180 / 720 * 7.5)
        self.assertLessEqual(left + width, (160 + 760) / 1280 * 13.333)
        self.assertLessEqual(top + height, (180 + 360) / 720 * 7.5)

    def test_repeated_template_hint_rotates_before_reuse(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": []},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)], "pictures": []},
            ],
        }
        used = Counter()
        first = _native_pick_slide(library, {"kind": "content", "layout_id": "content", "template_slide": 1}, used, None)
        used[first["slide_index"]] += 1
        second = _native_pick_slide(library, {"kind": "content", "layout_id": "content", "template_slide": 1}, used, first["slide_index"])
        self.assertEqual(first["slide_index"], 1)
        self.assertEqual(second["slide_index"], 2)

    def test_structure_repair_does_not_drop_authored_evidence(self) -> None:
        slides = [{"kind": "cover", "title": "封面"}]
        slides.extend({"kind": "content", "title": f"事实{index}", "body": f"唯一证据 TOKEN-{index}"} for index in range(1, 8))
        slides.append({"kind": "chart", "title": "数据页", "chart": {"labels": ["甲", "乙"], "values": [1, 2]}})
        slides.append({"kind": "closing", "title": "结束"})
        repaired, repairs = _ensure_report_structure({"slides": slides}, 10, native_agenda=True)
        visible = json.dumps(repaired, ensure_ascii=False)
        for index in range(1, 8):
            self.assertIn(f"TOKEN-{index}", visible)
        self.assertIn("数据页", visible)
        self.assertEqual(repairs, [])

    def test_repeated_chapter_heading_is_preserved_as_template_chrome(self) -> None:
        slot_data = slot("title_candidate", "福州空港概况", 357, 574, 74)
        slot_data["geometry"]["y"] = 393
        slot_data["text_metrics"] = {"font_size_px": 53.33}
        self.assertTrue(_native_is_chrome(slot_data, Counter({"福州空港概况": 4}), 31))

    def test_chapter_title_maps_to_large_title_slot(self) -> None:
        chapter = {
            "slide_index": 3,
            "page_type": "chapter_candidate",
            "slots": [
                {"slot_id": "chapter-title", "role": "title_candidate", "text": "福州空港概况", "geometry": {"x": 357, "y": 393, "width": 574, "height": 74}, "text_metrics": {"font_size_px": 53.33}},
                {"slot_id": "chapter-label", "role": "label_candidate", "text": "项目说明", "geometry": {"x": 502, "y": 324, "width": 275, "height": 110}, "text_metrics": {"font_size_px": 64}},
            ],
        }
        replacements = _native_replacements(
            chapter,
            {"title": "项目范围覆盖重点护卫场景", "summary": "岗位边界已经明确。"},
            Counter({"福州空港概况": 4}),
            31,
        )
        self.assertEqual(replacements[0]["slot_id"], "chapter-label")
        self.assertIn({"slot_id": "chapter-title", "text": ""}, replacements)

    def test_chapter_counter_is_rewritten_instead_of_leaking_source_number(self) -> None:
        chapter = {
            "slide_index": 3,
            "page_type": "chapter_candidate",
            "slots": [
                {"slot_id": "chapter-title", "role": "title_candidate", "text": "旧章节", "geometry": {"x": 357, "y": 393, "width": 574, "height": 74}, "text_metrics": {"font_size_px": 53.33}},
                {"slot_id": "part-counter", "role": "label_candidate", "text": "Part Three", "geometry": {"x": 500, "y": 100, "width": 280, "height": 60}, "text_metrics": {"font_size_px": 24}},
            ],
        }
        replacements = _native_replacements(
            chapter,
            {"title": "项目范围", "kicker": "PART 01"},
            Counter({"Part Three": 4}),
            31,
        )
        self.assertIn({"slot_id": "part-counter", "text": "PART 01"}, replacements)

    def test_chapter_title_prefers_complete_phrase_over_ellipsis(self) -> None:
        self.assertEqual(_native_compact_chapter_title("项目范围与岗位边界"), "项目范围")
        self.assertEqual(_native_compact_chapter_title("可行性研究报告"), "可行性研究报告")

    def test_native_closing_copy_stays_on_one_line(self) -> None:
        narrow_display_slot = {
            "geometry": {"x": 873, "y": 240, "width": 390, "height": 165},
            "text_metrics": {"font_size_px": 128},
        }
        wide_display_slot = {
            "geometry": {"x": 700, "y": 240, "width": 720, "height": 165},
            "text_metrics": {"font_size_px": 96},
        }
        self.assertEqual(_native_closing_text("感谢聆听", narrow_display_slot), "谢谢")
        self.assertEqual(_native_closing_text("谢谢", narrow_display_slot), "谢谢")
        self.assertEqual(_native_closing_text("感谢聆听", wide_display_slot), "感谢聆听")

    def test_empty_grouped_visual_slot_is_marked_for_removal(self) -> None:
        source = {
            "slots": [
                {
                    "slot_id": "s28_sh49",
                    "shape_id": "49",
                    "role": "label_candidate",
                    "text": "",
                    "group_depth": 1,
                    "text_node_count": 0,
                    "geometry": {"x": 77, "y": 518, "width": 213, "height": 108},
                },
                {
                    "slot_id": "s28_sh50",
                    "shape_id": "50",
                    "role": "label_candidate",
                    "text": "市场成熟",
                    "group_depth": 1,
                    "text_node_count": 1,
                    "geometry": {"x": 98, "y": 534, "width": 172, "height": 74},
                },
            ]
        }
        self.assertEqual(
            _native_hidden_empty_shape_ids(
                source,
                [{"slot_id": "s28_sh49", "text": ""}, {"slot_id": "s28_sh50", "text": ""}],
            ),
            [49],
        )

    def test_generic_chapter_title_uses_nearby_topic(self) -> None:
        slides = [
            {"title": "项目范围与岗位边界", "summary": "范围已明确"},
            {"title": "第 02 页核心判断", "summary": "岗位配置按风险和客流分层"},
            {"title": "成本测算显示节约空间", "summary": ""},
        ]
        self.assertEqual(_native_title_candidate(slides[1]), "岗位配置按风险和客流分层")
        self.assertEqual(_native_chapter_title(slides, 1), "岗位配置按风险和客流分层")

    def test_generic_body_title_uses_page_summary(self) -> None:
        body = {
            "slide_index": 4,
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "body-title", "role": "title_candidate", "text": "原示例标题", "geometry": {"x": 100, "y": 20, "width": 800, "height": 50}, "text_metrics": {"font_size_px": 32}},
            ],
        }
        replacements = _native_replacements(
            body,
            {"title": "第 04 页核心判断", "summary": "成本测算支持方案决策"},
            Counter(),
            10,
        )
        self.assertEqual(replacements[0]["text"], "成本测算支持方案决策")

    def test_large_misclassified_title_frame_receives_body_copy(self) -> None:
        shell = {
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "page-title", "role": "title_candidate", "text": "旧标题", "geometry": {"x": 240, "y": 13, "width": 820, "height": 55}, "text_metrics": {"font_size_px": 37}},
                {"slot_id": "side-label", "role": "label_candidate", "text": "旧侧标", "geometry": {"x": 40, "y": 260, "width": 330, "height": 220}, "text_metrics": {"font_size_px": 48}},
                {"slot_id": "body", "role": "title_candidate", "text": "原来的多行正文内容，用于说明岗位配置与职责边界。", "paragraph_count": 3, "geometry": {"x": 375, "y": 148, "width": 793, "height": 478}, "text_metrics": {"font_size_px": 33}},
            ],
        }
        self.assertTrue(_native_slot_is_body_like(shell["slots"][2]))
        self.assertEqual(_native_layout_family(shell), "hero-body")
        replacements = _native_replacements(
            shell,
            {"kind": "statement", "title": "岗位职责需要覆盖全天候运行", "body": "巡逻、交接与异常处置形成完整闭环。"},
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["body"], "巡逻、交接与异常处置形成完整闭环。")
        self.assertEqual(mapping["side-label"], "岗位职责")

    def test_empty_slot_detection_keeps_text_placeholders_but_rejects_artwork(self) -> None:
        placeholder = {"role": "label_candidate", "shape_name": "文本框 1", "text": "", "geometry": {"width": 260, "height": 90}}
        large_mask = {"role": "title_candidate", "shape_name": "任意多边形 1", "text": "", "geometry": {"width": 930, "height": 520}}
        grouped_tile = {"role": "label_candidate", "shape_name": "平行四边形 1", "group_depth": 1, "text": "", "geometry": {"width": 213, "height": 108}}
        self.assertFalse(_native_slot_is_empty_decor(placeholder))
        self.assertTrue(_native_slot_is_empty_decor(large_mask))
        self.assertTrue(_native_slot_is_empty_decor(grouped_tile))

    def test_chart_without_native_chart_does_not_choose_table_shell(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "cover_candidate", "slots": []},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
                {"slide_index": 3, "page_type": "content_candidate", "slots": [slot("body_candidate", "双栏左", 100), slot("body_candidate", "双栏右", 700)]},
                {"slide_index": 4, "page_type": "content_candidate", "slots": [], "tables": [{"table_id": "example", "row_count": 5, "column_count": 4}]},
                {"slide_index": 5, "page_type": "ending_candidate", "slots": []},
            ]
        }
        selected = _native_pick_slide(
            library,
            {"kind": "chart", "layout_id": "bar-chart", "chart": {"labels": ["甲", "乙"], "values": [1, 2]}},
            Counter(),
            None,
        )
        self.assertNotEqual(selected["slide_index"], 4)

    def test_chart_overlay_moves_from_compact_split_to_spacious_shell(self) -> None:
        library = {
            "canvas_px": {"width": 1280, "height": 720},
            "slides": [
                {"slide_index": 1, "page_type": "cover_candidate", "slots": []},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "紧凑双栏", 300, 700, 90)]},
                {"slide_index": 3, "page_type": "content_candidate", "slots": [slot("body_candidate", "宽阔正文", 120, 1000, 420)]},
                {"slide_index": 4, "page_type": "ending_candidate", "slots": []},
            ],
        }
        plan = _build_native_fill_plan(
            library,
            [
                {"kind": "cover", "layout_id": "cover", "title": "封面"},
                {"kind": "chart", "layout_id": "donut-chart", "title": "构成", "chart": {"type": "donut", "labels": ["甲", "乙"], "values": [1, 2]}},
                {"kind": "closing", "layout_id": "closing", "title": "结束"},
            ],
        )
        chart_item = next(item for item in plan["slides"] if item.get("shape_chart"))
        self.assertEqual(chart_item["source_slide"], 3)
        self.assertGreaterEqual(chart_item["shape_chart"]["bounds"][2], 6)
        self.assertGreaterEqual(chart_item["shape_chart"]["bounds"][3], 2.2)

    def test_split_shell_maps_mixed_label_and_body_frames_in_reading_order(self) -> None:
        shell = {
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "left-heading", "role": "label_candidate", "text": "旧左标题", "geometry": {"x": 80, "y": 120, "width": 180, "height": 60}, "text_metrics": {"font_size_px": 24}},
                {"slot_id": "left-body", "role": "label_candidate", "text": "旧左正文", "geometry": {"x": 80, "y": 230, "width": 180, "height": 60}, "text_metrics": {"font_size_px": 24}},
                {"slot_id": "right-body", "role": "body_candidate", "text": "旧右正文", "geometry": {"x": 400, "y": 120, "width": 600, "height": 80}, "text_metrics": {"font_size_px": 24}},
                {"slot_id": "right-detail", "role": "label_candidate", "text": "旧右细节", "geometry": {"x": 400, "y": 240, "width": 600, "height": 80}, "text_metrics": {"font_size_px": 24}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {"kind": "comparison", "title": "判断", "left_title": "现状", "left_bullets": ["左证据"], "right_title": "方案", "right_bullets": ["右证据"]},
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["right-body"], "方案")
        self.assertEqual(mapping["right-detail"], "右证据")

    def test_comparison_uses_all_rows_of_a_vertical_stack(self) -> None:
        slots: list[dict[str, object]] = [
            {"slot_id": "page-title", "role": "title_candidate", "text": "旧标题", "geometry": {"x": 240, "y": 13, "width": 820, "height": 55}, "text_metrics": {"font_size_px": 37}},
        ]
        for index in range(4):
            y = 130 + index * 110
            slots.extend((
                {"slot_id": f"row-label-{index}", "role": "label_candidate", "text": f"旧标签{index}", "geometry": {"x": 72, "y": y, "width": 178, "height": 82}, "text_metrics": {"font_size_px": 26}},
                {"slot_id": f"row-body-{index}", "role": "label_candidate", "text": f"原来的第{index + 1}行详细说明内容。", "geometry": {"x": 320, "y": y + 8, "width": 820, "height": 62}, "text_metrics": {"font_size_px": 24}},
            ))
        shell = {"slide_index": 2, "page_type": "content_candidate", "slots": slots}
        comparison_spec = {
            "kind": "comparison",
            "title": "现状与方案",
            "left_title": "现状",
            "left_bullets": ["岗位职责边界不清", "高峰调度依赖经验"],
            "right_title": "方案",
            "right_bullets": ["按区域分层配置", "建立交接检查"],
        }
        replacements = _native_replacements(
            shell,
            comparison_spec,
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["row-label-0"], "现状")
        self.assertEqual(mapping["row-body-0"], "岗位职责边界不清")
        self.assertEqual(mapping["row-label-2"], "方案")
        self.assertEqual(mapping["row-body-3"], "建立交接检查")
        split = {
            "slide_index": 1,
            "page_type": "content_candidate",
            "slots": [
                slot("body_candidate", "左侧正文", 80, 480, 300),
                slot("body_candidate", "右侧正文", 680, 480, 300),
            ],
        }
        selected = _native_pick_slide(
            {"slides": [split, shell]},
            comparison_spec,
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_layout_family_detects_grid_and_split(self) -> None:
        grid = {
            "page_type": "content_candidate",
            "slots": [
                *(slot("label_candidate", f"L{index}", (index % 3) * 400) for index in range(6)),
                *(slot("body_candidate", f"B{index}", index * 400) for index in range(3)),
            ],
            "tables": [],
        }
        split = {
            "page_type": "content_candidate",
            "slots": [slot("body_candidate", "左", 40), slot("body_candidate", "右", 700)],
            "tables": [],
        }
        self.assertTrue(_native_layout_family(grid).startswith("grid-"))
        self.assertEqual(_native_layout_family(split), "split")

    def test_toc_is_not_selected_for_ordinary_content(self) -> None:
        toc = {
            "slide_index": 2,
            "page_type": "content_candidate",
            "text_summary": "目    录 | 第一部分 | 第二部分",
            "slots": [slot("label_candidate", f"目录{index}", 420, 640, 80) for index in range(10)],
            "tables": [],
            "charts": [],
            "diagrams": [],
        }
        content = {
            "slide_index": 6,
            "page_type": "content_candidate",
            "text_summary": "项目范围",
            "slots": [slot("body_candidate", "正文", 240, 850, 380)],
            "tables": [],
            "charts": [],
            "diagrams": [],
        }
        selected = _native_pick_slide(
            {"slides": [toc, content]},
            {"kind": "content", "layout_id": "content"},
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 6)

    def test_native_chart_is_not_selected_without_verified_payload(self) -> None:
        chart_page = {
            "slide_index": 1,
            "page_type": "content_candidate",
            "text_summary": "示例图表",
            "slots": [slot("title_candidate", "标题", 120, 700, 90)],
            "tables": [],
            "charts": [{"chart_id": "s01_ch2"}],
            "diagrams": [],
        }
        text_page = {
            "slide_index": 2,
            "page_type": "content_candidate",
            "text_summary": "正文",
            "slots": [slot("body_candidate", "正文", 120, 850, 360)],
            "tables": [],
            "charts": [],
            "diagrams": [],
        }
        selected = _native_pick_slide(
            {"slides": [chart_page, text_page]},
            {"kind": "content", "layout_id": "content"},
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_smartart_page_is_excluded_from_generated_content(self) -> None:
        diagram_page = {
            "slide_index": 1,
            "page_type": "content_candidate",
            "slots": [slot("body_candidate", "旧流程", 120, 850, 360)],
            "diagrams": [{"diagram_id": "s01_d1"}],
        }
        text_page = {
            "slide_index": 2,
            "page_type": "content_candidate",
            "slots": [slot("body_candidate", "正文", 120, 850, 360)],
            "diagrams": [],
        }
        selected = _native_pick_slide(
            {"slides": [diagram_page, text_page]},
            {"kind": "timeline", "layout_id": "timeline", "steps": [{"title": "启动", "body": "执行"}]},
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_native_chart_edits_require_every_chart_to_be_supported(self) -> None:
        payload = {"chart": {"labels": ["甲", "乙"], "values": [1, 2]}}
        supported = {"chart_id": "s01_ch1", "edit_capability": {"supported": True}}
        unsupported = {"chart_id": "s01_ch2", "edit_capability": {"supported": False}}
        self.assertEqual(_native_chart_edits({"charts": [supported, unsupported]}, payload), [])
        edits = _native_chart_edits({"charts": [supported, {**supported, "chart_id": "s01_ch3"}]}, payload)
        self.assertEqual([item["chart_id"] for item in edits], ["s01_ch1", "s01_ch3"])

    def test_metric_does_not_use_dense_research_grid_shell(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": [*(slot("label_candidate", f"卡片{index}", (index % 3) * 400) for index in range(6)), *(slot("body_candidate", f"内容{index}", (index % 3) * 400) for index in range(3))]},
                {"slide_index": 2, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
            ],
        }
        selected = _native_pick_slide(
            library,
            {"kind": "metric", "layout_id": "metric", "metric": "92%", "metric_label": "满意度"},
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_statement_prefers_single_shell_over_three_columns(self) -> None:
        library = {
            "slides": [
                {
                    "slide_index": 1,
                    "page_type": "content_candidate",
                    "slots": [
                        slot("label_candidate", "一", 80),
                        slot("label_candidate", "二", 480),
                        slot("label_candidate", "三", 880),
                    ],
                },
                {
                    "slide_index": 2,
                    "page_type": "content_candidate",
                    "slots": [slot("body_candidate", "单一结论", 120, 900, 320)],
                },
                {
                    "slide_index": 3,
                    "page_type": "content_candidate",
                    "slots": [
                        slot("body_candidate", "左侧正文", 80, 480, 300),
                        slot("body_candidate", "右侧正文", 680, 480, 300),
                    ],
                },
            ]
        }
        selected = _native_pick_slide(
            library,
            {"kind": "statement", "layout_id": "statement", "title": "核心判断", "body": "结论"},
            Counter({2: 2}),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_card_selection_matches_item_count_instead_of_grid_columns(self) -> None:
        def content_slot(role: str, text: str, x: int, y: int) -> dict[str, object]:
            value = slot(role, text, x, 280, 100)
            value["geometry"]["y"] = y
            return value

        three_cards = [
            *(content_slot("label_candidate", f"三项标题{index}", x, 180) for index, x in enumerate((80, 480, 880))),
            *(content_slot("body_candidate", f"三项正文{index}", x, 300) for index, x in enumerate((80, 480, 880))),
        ]
        six_positions = ((80, 150), (480, 150), (880, 150), (80, 390), (480, 390), (880, 390))
        six_cards = [
            *(content_slot("label_candidate", f"六项标题{index}", x, y) for index, (x, y) in enumerate(six_positions)),
            *(content_slot("body_candidate", f"六项正文{index}", x, y + 90) for index, (x, y) in enumerate(six_positions)),
        ]
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "content_candidate", "slots": six_cards},
                {"slide_index": 2, "page_type": "content_candidate", "slots": three_cards},
            ]
        }
        self.assertEqual(_native_layout_unit_capacity(library["slides"][0]), 6)
        self.assertEqual(_native_layout_unit_capacity(library["slides"][1]), 3)
        selected = _native_pick_slide(
            library,
            {"kind": "cards", "layout_id": "cards", "bullets": ["甲", "乙", "丙"]},
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 2)

    def test_timeline_steps_map_to_native_title_and_body_slots(self) -> None:
        shell = {
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "page-title", "role": "title_candidate", "text": "旧标题", "geometry": {"x": 80, "y": 20, "width": 1000, "height": 60}, "text_metrics": {"font_size_px": 32}},
                {"slot_id": "step-1-title", "role": "label_candidate", "text": "旧步骤一", "geometry": {"x": 90, "y": 170, "width": 230, "height": 50}, "text_metrics": {"font_size_px": 22}},
                {"slot_id": "step-1-body", "role": "body_candidate", "text": "旧说明一", "geometry": {"x": 350, "y": 170, "width": 760, "height": 70}, "text_metrics": {"font_size_px": 20}},
                {"slot_id": "step-2-title", "role": "label_candidate", "text": "旧步骤二", "geometry": {"x": 90, "y": 310, "width": 230, "height": 50}, "text_metrics": {"font_size_px": 22}},
                {"slot_id": "step-2-body", "role": "body_candidate", "text": "旧说明二", "geometry": {"x": 350, "y": 310, "width": 760, "height": 70}, "text_metrics": {"font_size_px": 20}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {
                "kind": "timeline",
                "layout_id": "timeline",
                "title": "实施计划",
                "steps": [
                    {"title": "启动", "body": "完成动员"},
                    {"title": "验收", "body": "核对成果"},
                ],
            },
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["step-1-title"], "启动")
        self.assertEqual(mapping["step-1-body"], "完成动员")
        self.assertEqual(mapping["step-2-title"], "验收")
        self.assertEqual(mapping["step-2-body"], "核对成果")

    def test_row_pair_geometry_recognizes_three_four_and_five_step_shells(self) -> None:
        def row_shell(slide_index: int, count: int) -> dict[str, object]:
            slots: list[dict[str, object]] = [
                {"slot_id": f"title-{slide_index}", "role": "title_candidate", "text": "实施计划", "geometry": {"x": 240, "y": 13, "width": 820, "height": 55}, "text_metrics": {"font_size_px": 37}},
            ]
            for index in range(count):
                y = 130 + index * 100
                slots.extend((
                    {"slot_id": f"label-{slide_index}-{index}", "role": "label_candidate", "text": f"阶段{index + 1}", "geometry": {"x": 72, "y": y, "width": 178, "height": 82}, "text_metrics": {"font_size_px": 26}},
                    {"slot_id": f"body-{slide_index}-{index}", "role": "label_candidate", "text": f"第{index + 1}阶段需要完成一项完整的实施任务。", "geometry": {"x": 300, "y": y + 8, "width": 820, "height": 62}, "text_metrics": {"font_size_px": 26}},
                ))
            return {"slide_index": slide_index, "page_type": "content_candidate", "slots": slots}

        shells = [row_shell(3, 3), row_shell(4, 4), row_shell(5, 5)]
        next(slot for slot in shells[0]["slots"] if slot["slot_id"] == "body-3-0")["text"] = "已完成"
        shells[0]["slots"].append(
            {"slot_id": "unpaired-badge", "role": "label_candidate", "text": "管理提示", "geometry": {"x": 1140, "y": 640, "width": 110, "height": 42}, "text_metrics": {"font_size_px": 18}}
        )
        for expected, shell in zip((3, 4, 5), shells):
            self.assertEqual(len(_native_row_pairs(shell["slots"])), expected)
            self.assertEqual(_native_layout_family(shell), "stack")
            self.assertEqual(_native_layout_unit_capacity(shell), expected)

        selected = _native_pick_slide(
            {"slides": [shells[2], shells[0], shells[1]]},
            {
                "kind": "timeline",
                "layout_id": "timeline",
                "title": "三阶段实施",
                "steps": [
                    {"title": "启动", "body": "完成动员与人员到位。"},
                    {"title": "交接", "body": "统一巡查和异常处置记录。"},
                    {"title": "复盘", "body": "按月评估并持续优化。"},
                ],
            },
            Counter(),
            None,
        )
        self.assertEqual(selected["slide_index"], 3)
        replacements = _native_replacements(
            selected,
            {
                "kind": "timeline",
                "layout_id": "timeline",
                "title": "三阶段实施",
                "steps": [
                    {"title": "启动", "body": "完成动员与人员到位。"},
                    {"title": "交接", "body": "统一巡查和异常处置记录。"},
                    {"title": "复盘", "body": "按月评估并持续优化。"},
                ],
            },
            Counter(),
            3,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["label-3-0"], "启动")
        self.assertEqual(mapping["body-3-0"], "完成动员与人员到位。")
        self.assertEqual(mapping["label-3-2"], "复盘")
        self.assertEqual(mapping["body-3-2"], "按月评估并持续优化。")

    def test_chart_hero_shell_keeps_a_series_side_label(self) -> None:
        shell = {
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "page-title", "role": "title_candidate", "text": "旧标题", "geometry": {"x": 240, "y": 13, "width": 820, "height": 55}, "text_metrics": {"font_size_px": 37}},
                {"slot_id": "side-label", "role": "label_candidate", "text": "旧侧标", "geometry": {"x": 40, "y": 260, "width": 330, "height": 220}, "text_metrics": {"font_size_px": 48}},
                {"slot_id": "body", "role": "title_candidate", "text": "原来的图表说明和数据区域", "geometry": {"x": 375, "y": 148, "width": 793, "height": 478}, "text_metrics": {"font_size_px": 28}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {"kind": "chart", "layout_id": "bar-chart", "title": "岗位构成", "metric": "构成", "chart": {"series_name": "岗位数量", "labels": ["甲", "乙"], "values": [1, 2]}},
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["side-label"], "岗位数量")

    def test_long_closing_message_uses_a_clean_native_thank_you_label(self) -> None:
        shell = {
            "page_type": "ending_candidate",
            "slots": [
                {"slot_id": "closing", "role": "label_candidate", "text": "谢 谢！", "geometry": {"x": 873, "y": 240, "width": 390, "height": 165}, "text_metrics": {"font_size_px": 128}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {"kind": "closing", "title": "形成可执行、可追溯、可复盘的护卫服务体系"},
            Counter(),
            1,
        )
        self.assertEqual(replacements, [{"slot_id": "closing", "text": "谢谢"}])

    def test_native_page_badge_is_rewritten_after_slide_reordering(self) -> None:
        shell = {
            "page_type": "content_candidate",
            "slots": [
                {"slot_id": "badge", "role": "title_candidate", "text": "08", "geometry": {"x": 159, "y": 11, "width": 87, "height": 61}, "text_metrics": {"font_size_px": 42}},
                {"slot_id": "page-title", "role": "title_candidate", "text": "旧标题", "geometry": {"x": 240, "y": 13, "width": 820, "height": 55}, "text_metrics": {"font_size_px": 37}},
                {"slot_id": "body", "role": "body_candidate", "text": "旧正文", "geometry": {"x": 120, "y": 180, "width": 960, "height": 320}, "text_metrics": {"font_size_px": 24}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {"kind": "content", "title": "新标题", "body": "新正文", "_native_page_marker": "03"},
            Counter({"08": 4}),
            10,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["badge"], "03")

    def test_agenda_heading_stays_compact_in_the_native_title_box(self) -> None:
        shell = {
            "page_type": "toc_candidate",
            "slots": [
                {"slot_id": "agenda-title", "role": "title_candidate", "text": "目录", "geometry": {"x": 500, "y": 60, "width": 120, "height": 90}, "text_metrics": {"font_size_px": 36}},
                {"slot_id": "agenda-item", "role": "label_candidate", "text": "第一部分", "geometry": {"x": 440, "y": 190, "width": 500, "height": 70}, "text_metrics": {"font_size_px": 24}},
            ],
        }
        replacements = _native_replacements(
            shell,
            {"kind": "agenda", "title": "汇报目录", "bullets": ["背景与范围"]},
            Counter(),
            1,
        )
        mapping = {item["slot_id"]: item["text"] for item in replacements}
        self.assertEqual(mapping["agenda-title"], "目录")

    def test_native_table_edits_write_first_column_and_clear_secondary_table(self) -> None:
        def table(table_id: str, row_count: int, column_count: int) -> dict[str, object]:
            return {
                "table_id": table_id,
                "row_count": row_count,
                "column_count": column_count,
                "rows": [
                    {
                        "row": row_index,
                        "cells": [
                            {
                                "row": row_index,
                                "col": col_index,
                                "text": f"旧{row_index}-{col_index}",
                                "is_merge_slave": False,
                                "col_span": 1,
                            }
                            for col_index in range(column_count)
                        ],
                    }
                    for row_index in range(row_count)
                ],
            }

        edits = _native_table_edits(
            {"tables": [table("large", 6, 4), table("fit", 3, 2)]},
            {
                "title": "岗位测算",
                "table": {
                    "columns": ["岗位", "人数"],
                    "rows": [["护卫", "100"], ["管理", "4"]],
                },
            },
        )
        self.assertEqual(len(edits), 2)
        by_id = {
            item["table_id"]: {
                (cell["row"], cell["col"]): cell["text"] for cell in item["cells"]
            }
            for item in edits
        }
        self.assertEqual(by_id["fit"][(0, 0)], "岗位")
        self.assertEqual(by_id["fit"][(1, 0)], "护卫")
        self.assertEqual(by_id["fit"][(2, 0)], "管理")
        self.assertTrue(all(value == "" for value in by_id["large"].values()))

    def test_native_sequence_preserves_source_page_families(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "cover_candidate", "slots": []},
                {"slide_index": 2, "page_type": "chapter_candidate", "slots": []},
                {"slide_index": 3, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
                {"slide_index": 4, "page_type": "content_candidate", "slots": [slot("body_candidate", "左", 100), slot("body_candidate", "右", 700)]},
                {"slide_index": 5, "page_type": "chapter_candidate", "slots": []},
                {"slide_index": 6, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
                {"slide_index": 7, "page_type": "ending_candidate", "slots": []},
            ]
        }
        sequence = _native_sequence_sources(library, 10)
        self.assertEqual(sequence[0]["slide_index"], 1)
        self.assertEqual(sequence[-1]["slide_index"], 7)
        self.assertIn("chapter", [_native_layout_family(item) for item in sequence])
        self.assertGreaterEqual(len({_native_layout_family(item) for item in sequence}), 4)

    def test_native_fill_uses_semantic_shells_for_matching_count(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "cover_candidate", "slots": [slot("title_candidate", "封面", 100)]},
                {"slide_index": 2, "page_type": "chapter_candidate", "slots": [slot("title_candidate", "章节", 100)]},
                {"slide_index": 3, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
                {"slide_index": 4, "page_type": "ending_candidate", "slots": [slot("title_candidate", "结束", 100)]},
            ]
        }
        slides = [
            {"kind": "cover", "layout_id": "cover", "title": "封面"},
            {"kind": "content", "layout_id": "content", "title": "第一部分"},
            {"kind": "content", "layout_id": "content", "title": "正文", "body": "内容"},
            {"kind": "closing", "layout_id": "closing", "title": "结束"},
        ]
        plan = _build_native_fill_plan(library, slides)
        source_slides = [item["source_slide"] for item in plan["slides"]]
        self.assertEqual(source_slides[0], 1)
        self.assertEqual(source_slides[-1], 4)
        self.assertEqual(len(source_slides), 4)

    def test_native_sequence_keeps_cover_and_ending_at_edges_for_long_deck(self) -> None:
        library = {
            "slides": [
                {"slide_index": 1, "page_type": "cover_candidate", "slots": []},
                {"slide_index": 2, "page_type": "chapter_candidate", "slots": []},
                {"slide_index": 3, "page_type": "content_candidate", "slots": [slot("body_candidate", "正文", 100, 900, 300)]},
                {"slide_index": 4, "page_type": "content_candidate", "slots": [slot("body_candidate", "左", 100), slot("body_candidate", "右", 700)]},
                {"slide_index": 5, "page_type": "chapter_candidate", "slots": []},
                {"slide_index": 6, "page_type": "content_candidate", "slots": [slot("body_candidate", "数据", 100, 900, 300)], "charts": [{"type": "bar"}]},
                {"slide_index": 7, "page_type": "ending_candidate", "slots": []},
            ]
        }
        sequence = _native_sequence_sources(library, 100)
        self.assertEqual(len(sequence), 100)
        roles = [_native_page_role(item) for item in sequence]
        self.assertEqual(roles[0], "cover")
        self.assertEqual(roles[-1], "ending")
        self.assertNotIn("cover", roles[1:-1])
        self.assertNotIn("ending", roles[1:-1])
        self.assertGreaterEqual(roles[1:-1].count("chapter"), 2)


class PresentationLayoutVarietyTests(unittest.TestCase):
    def test_generic_content_pages_rotate_real_layouts(self) -> None:
        spec = {
            "slides": [
                {"kind": "cover", "layout_id": "cover", "title": "封面"},
                *(
                    {"kind": "content", "layout_id": "content", "title": f"正文{index}", "body": "结论", "bullets": ["证据一", "证据二"]}
                    for index in range(1, 6)
                ),
                {"kind": "closing", "layout_id": "closing", "title": "结束"},
            ]
        }
        enriched = _enforce_layout_variety(spec, {"layout_variety"})
        layouts = [item["layout_id"] for item in enriched["slides"][1:-1]]
        self.assertGreaterEqual(len(set(layouts)), 4)
        self.assertIn("content-stack", layouts)
        self.assertIn("content-emphasis", layouts)


if __name__ == "__main__":
    unittest.main()
